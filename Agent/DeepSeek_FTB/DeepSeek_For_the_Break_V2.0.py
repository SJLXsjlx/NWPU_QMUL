#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import uuid
import json
import sqlite3
import subprocess
import shutil
import traceback
import textwrap
import warnings
from pathlib import Path
from typing import Optional, Dict

# ── 第三方库 ──────────────────────────────────────────────
import pandas as pd
import docx
import requests
from docx import Document
from docx.shared import Inches
from duckduckgo_search import DDGS
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_community.tools import WriteFileTool
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.sqlite import SqliteSaver
from pptx import Presentation
from PIL import Image
import pytesseract
import fitz  # PyMuPDF

# ── 可选依赖：matplotlib ─────────────────────────────────
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt      # noqa: F401
    _MPL_AVAILABLE = True
except ImportError:
    _MPL_AVAILABLE = False
    warnings.warn("⚠ matplotlib 未安装，图表生成功能不可用。请执行: pip install matplotlib")

# ── 可选依赖：Mermaid CLI ────────────────────────────────
_MMDC_PATH = shutil.which("mmdc")

# ================== 0. 智能路径配置 ==================

def _detect_tesseract() -> Optional[str]:
    from_path = shutil.which("tesseract")
    if from_path:
        return from_path
    if sys.platform == "win32":
        candidates = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Tesseract-OCR\tesseract.exe",
            os.path.expandvars(r"%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe"),
            os.path.expandvars(r"%PROGRAMFILES%\Tesseract-OCR\tesseract.exe"),
        ]
    else:
        candidates = [
            "/usr/bin/tesseract",
            "/usr/local/bin/tesseract",
            "/opt/homebrew/bin/tesseract",
            "/opt/local/bin/tesseract",
        ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


def _detect_latex_compilers() -> Dict[str, Optional[str]]:
    compilers: Dict[str, Optional[str]] = {}
    for name in ["pdflatex", "xelatex", "lualatex"]:
        exe_name = name + (".exe" if sys.platform == "win32" else "")
        compilers[name] = shutil.which(exe_name) or shutil.which(name)
    return compilers


# ── 初始化检测 ────────────────────────────────────────────
_TESSERACT_PATH = _detect_tesseract()
_LATEX_COMPILERS = _detect_latex_compilers()

if _TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH
    print(f"🔍 自动检测到 Tesseract: {_TESSERACT_PATH}")
else:
    print("⚠ 警告: 未检测到 Tesseract-OCR 引擎！")
    print("  请安装 Tesseract: https://github.com/tesseract-ocr/tesseract")
    print("  Windows 用户安装时请勾选中文语言包。")

_available_latex = [k for k, v in _LATEX_COMPILERS.items() if v]
if _available_latex:
    print(f"🔍 检测到 LaTeX 编译器: {', '.join(_available_latex)}")
else:
    print("⚠ 警告: 未检测到 LaTeX 编译器！CompileLatexTool 将不可用。")
    print("  请安装 MiKTeX (Windows) 或 TeX Live (macOS/Linux)")

# ================== 1. 核心目录与环境配置 ==================
load_dotenv(override=True)

api_key = os.getenv("DEEPSEEK_FOR_THE_BREAK")
if not api_key or not api_key.startswith("sk-"):
    print("❌ API Key 未加载或格式不正确！请检查 .env 文件。")
    sys.exit(1)

BASE_WORKSPACE = os.getenv(
    "DFTB_WORKSPACE",
    r"D:\.API Keys\DeepSeek_For_the_Break\outputs"
)
CUSTOM_TEMP_DIR = os.path.join(BASE_WORKSPACE, "temp")
os.makedirs(CUSTOM_TEMP_DIR, exist_ok=True)

if CUSTOM_TEMP_DIR not in sys.path:
    sys.path.insert(0, CUSTOM_TEMP_DIR)

print("✅ ds学术已启动（增强版 v2.0 — 沙箱执行 / 搜索 / LaTeX编译 / 图表生成）")
print(f"📁 工作区: {BASE_WORKSPACE}")
print(f"📁 临时目录: {CUSTOM_TEMP_DIR}")

# ================== 2. 系统 Prompt 与模型 ==================
SYSTEM_PROMPT = """你是一位高度专业、严谨且全面的**学术AI助手**（Academic AI Assistant），专注于为科研人员、学生和学者提供全方位的学术支持。你的核心定位是：

- **研究辅助专家**：帮助进行文献分析、论文构思与写作、数据处理与可视化、实验设计建议、学术幻灯片/报告制作等。
- **多格式文档处理专家**：精通PDF（使用PyMuPDF精确提取布局与文本）、Word、PowerPoint、Excel、LaTeX源码、图片OCR（中英双语）等文件的深度读取、解析与内容综合。
- **学术排版与生成专家**：擅长使用LaTeX生成高质量、可直接编译的学术论文、报告、公式排版；同时支持Word文档的结构化创建与编辑。
- **视觉辅助支持**：可通过工具搜索并下载与学术主题相关的专业图片、图表或示意图。
- **代码与工具集成专家**：能编写科研相关Python脚本、数据分析流程，并通过安全工具持久化保存所有生成内容到指定工作区。

**总体行为准则**（必须严格遵守）：

1. **客观严谨，结构清晰**：所有回答必须客观、中立、逻辑严密、层次分明。主动承认证据局限性、知识时效性或工具输出的不确定性。避免任何形式的猜测或过度泛化。

2. **【核心铁律 - 工具调用严格限制】**：
   - **SmartReadPathTool（文件/文件夹读取工具）** 是唯一允许读取本地文件的工具。
   - **只有在以下两种情况才允许调用**：
     a) 用户**明确提供了具体文件路径或文件夹路径**；
     b) 用户**显式要求**"读取"、"查看"、"打开"、"分析这个文件"、"检查路径下的内容"等操作。
   - **绝对禁止**：用户仅要求"写一个程序"、"解释某个概念"、"如何实现XX"、"总结XX领域"或进行任何常规对话时，**绝不能主动猜测路径、调用读取工具或编造文件位置**。此时应直接基于你的知识输出代码、解释或文本。
   - 当用户提供真实目录/文件夹时，必须先调用SmartReadPathTool进行递归遍历，然后**结合各子文件夹的相对路径结构 + 文件内容**进行整体归纳、交叉分析或针对性解答，并在回答中引用具体相对路径以便用户追溯。
   - 如果工具返回错误（如"找不到路径"），**立即停止一切路径猜测尝试**，直接基于已有知识回答，或礼貌询问用户提供正确路径。绝不反复尝试修改路径重试。

3. **学术写作与持久化输出优先**：
   - 涉及撰写学术论文、学位论文、研究报告、技术文档、或包含大量数学公式/方程的内容时，**必须主动、优先使用 EditTexFileTool** 生成结构完整、排版专业、可直接用 pdflatex/xelatex 编译的 .tex 源码。
   - 对于需要富文本、图片插入、表格的文档，优先使用 EditWordDocTool 创建或追加内容。
   - 所有生成的内容（代码、分析、文档）都应通过相应工具保存到 CUSTOM_TEMP_DIR 或工作区，实现持久化，方便用户后续使用。
   - 纯文本/笔记类内容优先使用 SaveMarkdownTool 保存为 .md 文件。
   - 如需生成数据图表，使用 ChartGenerationTool；如需流程图，使用 MermaidTool。

4. **错误处理与用户引导**：
   - 任何工具调用失败或参数错误时，清晰告知用户具体问题，提供替代方案或请求澄清。
   - 永远不要为了"完成任务"而编造信息或强行调用工具。

5. **响应风格与语言**：
   - 默认使用与用户查询相同的语言（本系统以中文为主）。
   - 采用正式、专业的学术语气。
   - 结构化输出：善用标题、 bullet points、编号列表、代码块、表格和 Markdown 格式提升可读性。
   - 复杂任务时，先在思考中拆解步骤，再逐步执行或输出。
   - 主动提供后续建议，如"需要我将以上分析保存为Word文档吗？"或直接调用工具持久化。

6. **多轮对话与记忆管理**：
   - 利用系统提供的 SQLite 持久化记忆（SqliteSaver）保持上下文连贯性。
   - 用户输入 'new'、'clear' 等时，系统会重置 thread_id，你应配合给出全新对话的提示。

7. **学术诚信与透明度**：
   - 明确区分"基于用户提供文件的内容"和"基于我的训练知识"的信息来源。
   - 对于OCR结果、PDF提取等，提醒可能存在的识别误差。
   - 不协助任何非学术或违反学术伦理的任务。
   - 主动说明工具的沙箱限制（所有文件操作仅限于指定工作区）。

8. **任务执行优先级**：
   - 文件读取 → 内容综合分析 → 结构化输出/保存
   - 学术写作 → 优先 LaTeX 工具
   - 视觉需求 → FetchWebImageTool（需提供英文关键词）
   - 通用代码/文本 → safe_write_tool 保存

9. **新增工具使用指南**：
   - 搜索学术资料 → DuckDuckGoSearchTool（支持中英文）
   - 执行数据分析/科学计算代码 → PythonSandboxTool
   - .tex 编译为 PDF → CompileLatexTool
   - 生成图表 → ChartGenerationTool (matplotlib) 或 MermaidTool (流程图)
   - 快速笔记 → SaveMarkdownTool

请始终进行**逐步推理（Chain-of-Thought）**后再决定是直接回复还是调用工具。你的目标是成为用户在学术道路上最可靠、最高效、最严谨的智能伙伴，帮助他们在科研中实现突破（For the Break）。

严格遵守以上所有规则，这是构建用户信任和系统稳定性的基础。"""

llm = ChatOpenAI(
    model="deepseek-v4-pro",
    base_url="https://api.deepseek.com",
    api_key=api_key,
    temperature=0.5,
    max_tokens=4096,
    model_kwargs={
        "reasoning_effort": "max",
        "extra_body": {
            "thinking": {"type": "enabled"}
        }
    }
)

# ================== 3. 自定义工具箱 ==================

def _safe_read_text(file_path: str) -> str:
    for encoding in ["utf-8", "gbk", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _resolve_path(path: str) -> str:
    p = Path(path)
    if p.is_absolute():
        return str(p)
    cwd_candidate = Path.cwd() / p
    if cwd_candidate.exists():
        return str(cwd_candidate.resolve())
    temp_candidate = Path(CUSTOM_TEMP_DIR) / p
    if temp_candidate.exists():
        return str(temp_candidate.resolve())
    return os.path.abspath(path)


def _ocr_image(file_path: str, lang: str = "chi_sim+eng") -> str:
    if not _TESSERACT_PATH:
        return (
            "❌ OCR 不可用：未检测到 Tesseract-OCR 引擎。\n"
            "   请访问 https://github.com/tesseract-ocr/tesseract 下载安装。\n"
            "   Windows 用户安装时请勾选中文语言包。"
        )
    try:
        img = Image.open(file_path)
        meta = f"📷 [图片信息] 格式: {img.format}, 尺寸: {img.size}, 模式: {img.mode}\n"
        try:
            ocr_text = pytesseract.image_to_string(img, lang=lang)
        except pytesseract.TesseractError as te:
            if "Failed loading language" in str(te):
                return (
                    meta
                    + f"⚠ OCR 语言包缺失: {lang}。\n"
                    + "   请下载中文简体语言包 (chi_sim): "
                    + "https://github.com/tesseract-ocr/tessdata\n"
                    + f"   原始错误: {te}"
                )
            return meta + f"❌ Tesseract 错误: {te}"
        if ocr_text.strip():
            return meta + f"【OCR 提取文本】:\n{ocr_text.strip()}"
        else:
            return meta + "【OCR 提示】: 未能在图片中识别出清晰的文本（可能为纯图/低分辨率/非文字内容）。"
    except FileNotFoundError:
        return f"❌ 图片文件不存在: {file_path}"
    except Exception as e:
        return f"❌ 图片 OCR 解析错误: {e}\n   堆栈: {traceback.format_exc()}"


# ── 工具 1: 智能文件读取 ─────────────────────────────────

@tool
def SmartReadPathTool(path: str) -> str:
    """
    【仅在用户明确给出路径时触发】读取指定路径的内容。支持单文件读取，也支持整个文件夹的深度递归遍历。
    支持格式: .txt, .md, .py, .tex, .json, .csv, .pdf, .docx, .xlsx, .pptx, .png, .jpg, .jpeg, .bmp。
    警告：绝对不要编造或猜测路径！如果没有明确的路径输入，请勿调用此工具！
    """
    resolved = _resolve_path(path)
    if not os.path.exists(resolved):
        return (
            f"❌ 找不到指定的路径: '{path}'\n"
            f"   解析后的路径: '{resolved}'\n"
            f"   请检查路径是否正确，或提供新的路径。\n"
            f"   💡 提示: 请停止尝试猜测路径，直接基于已有知识回答用户问题。"
        )

    def _read_single_file(file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        if ext in {"txt", "md", "py", "json", "tex", "yaml", "yml", "cfg", "ini", "log"}:
            return _safe_read_text(file_path)
        elif ext == "csv":
            for enc in ["utf-8", "gbk", "latin-1"]:
                try:
                    df = pd.read_csv(file_path, encoding=enc)
                    return f"[CSV: {len(df)} 行 × {len(df.columns)} 列]\n{df.to_string(max_rows=200)}"
                except (UnicodeDecodeError, UnicodeError):
                    continue
            return f"❌ CSV 编码读取失败: {file_path}"
        elif ext == "pdf":
            try:
                doc = fitz.open(file_path)
                text_parts = []
                total_pages = len(doc)
                for i, page in enumerate(doc):
                    page_text = page.get_text("text")
                    if page_text.strip():
                        text_parts.append(f"── Page {i+1}/{total_pages} ──\n{page_text.strip()}")
                if text_parts:
                    return "\n\n".join(text_parts)
                else:
                    ocr_results = []
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=200)
                        temp_img = os.path.join(CUSTOM_TEMP_DIR, f"_pdf_ocr_page_{i+1}.png")
                        pix.save(temp_img)
                        ocr_results.append(f"── Page {i+1}/{total_pages} (OCR) ──\n{_ocr_image(temp_img)}")
                    return "\n\n".join(ocr_results) if ocr_results else "[空 PDF 或无法提取内容]"
            except Exception as e:
                return f"❌ PDF 解析错误: {e}\n   {traceback.format_exc()}"
        elif ext == "docx":
            try:
                doc = docx.Document(file_path)
                content = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        content.append(p.text.strip())
                if doc.tables:
                    content.append("\n[文档表格]:")
                    for ti, table in enumerate(doc.tables):
                        content.append(f"  表格 {ti+1}:")
                        for row in table.rows:
                            row_data = [
                                cell.text.strip().replace("\n", " ")
                                for cell in row.cells if cell.text.strip()
                            ]
                            if row_data:
                                content.append("    | " + " | ".join(row_data))
                return "\n".join(content) if content else "[空 Word 文档]"
            except Exception as e:
                return f"❌ DOCX 解析错误: {e}"
        elif ext in {"xlsx", "xls"}:
            try:
                excel_file = pd.ExcelFile(file_path)
                sheet_outputs = []
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    max_rows = 500
                    if len(df) > max_rows:
                        sheet_outputs.append(
                            f"── Sheet: {sheet_name} (共 {len(df)} 行，仅显示前 {max_rows} 行) ──\n"
                            + df.head(max_rows).to_string()
                        )
                    else:
                        sheet_outputs.append(f"── Sheet: {sheet_name} ──\n{df.to_string()}")
                return "\n\n".join(sheet_outputs)
            except Exception as e:
                return f"❌ Excel 解析错误: {e}"
        elif ext == "pptx":
            try:
                prs = Presentation(file_path)
                ppt_outputs = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_data = [
                                    cell.text_frame.text.strip().replace("\n", " ")
                                    for cell in row.cells
                                    if cell.text_frame.text.strip()
                                ]
                                if row_data:
                                    slide_text.append(" | ".join(row_data))
                    header = f"── Slide {i+1} ──"
                    body = "\n".join(slide_text) if slide_text else "[无文本/非文本组件]"
                    ppt_outputs.append(f"{header}\n{body}")
                return "\n\n".join(ppt_outputs)
            except Exception as e:
                return f"❌ PPTX 解析错误: {e}"
        elif ext in {"png", "jpg", "jpeg", "bmp", "tiff", "tif"}:
            return _ocr_image(file_path)
        else:
            return f"[暂不支持的格式: .{ext}]"

    if os.path.isfile(resolved):
        result = _read_single_file(resolved)
        return f"📖 成功读取单文件: {resolved}\n{'='*50}\n{result}"
    elif os.path.isdir(resolved):
        supported_exts = {
            "txt", "md", "py", "tex", "json", "csv", "yaml", "yml",
            "pdf", "docx", "xlsx", "xls", "pptx", "png", "jpg", "jpeg", "bmp", "tiff",
        }
        results = [f"📂 递归遍历: {resolved}\n{'='*60}"]
        file_count = 0
        for root, dirs, files in os.walk(resolved):
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in {"__pycache__", "node_modules", ".git"}]
            for file in files:
                ext = os.path.splitext(file)[1].lower().lstrip(".")
                if ext in supported_exts:
                    file_count += 1
                    full_path = os.path.join(root, file)
                    rel_path = os.path.relpath(full_path, resolved)
                    results.append(f"\n📍 相对路径: {rel_path}\n{'-'*40}")
                    try:
                        results.append(_read_single_file(full_path))
                    except Exception as e:
                        results.append(f"❌ 读取错误: {e}")
                    results.append("=" * 60)
        if file_count == 0:
            return f"📁 目标文件夹 '{resolved}' 内没有找到支持的文件。"
        results.insert(1, f"   (共发现 {file_count} 个支持的文件)\n")
        return "\n".join(results)
    return f"❌ 路径 '{resolved}' 不是有效的文件或文件夹。"


# ── 工具 2: LaTeX 编辑 ──────────────────────────────────

@tool
def EditTexFileTool(tex_filename: str, latex_content: str, append: bool = False) -> str:
    """
    【学术专用】创建新 LaTeX 文件或修改已有 .tex 文件。
    参数:
    - tex_filename: .tex 文件名 (例如 'main.tex' 或 'sections/intro.tex')
    - latex_content: 要写入的完整 LaTeX 代码块或追加的文本段落
    - append: True=追加, False=覆盖/新建
    """
    if not tex_filename.lower().endswith(".tex"):
        tex_filename += ".tex"
    save_path = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    file_dir = os.path.dirname(save_path)
    if file_dir:
        os.makedirs(file_dir, exist_ok=True)
    mode = "a" if (append and os.path.exists(save_path)) else "w"
    try:
        with open(save_path, mode, encoding="utf-8") as f:
            if mode == "a":
                f.write("\n" + latex_content)
            else:
                f.write(latex_content)
        return (
            f"✅ LaTeX 文件保存成功: {save_path}\n"
            f"   💡 如需编译为 PDF，请使用 CompileLatexTool，参数 tex_filename='{tex_filename}'"
        )
    except Exception as e:
        return f"❌ LaTeX 文件操作失败: {e}"


# ── 工具 3: 一键编译 LaTeX → PDF ────────────────────────

@tool
def CompileLatexTool(tex_filename: str, compiler: str = "xelatex", clean_aux: bool = True) -> str:
    """
    【新增】一键编译 .tex 文件生成 PDF。
    参数:
    - tex_filename: 主 .tex 文件名 (如 'main.tex')
    - compiler: 编译器选择 ('pdflatex' / 'xelatex' / 'lualatex')，默认 xelatex（支持中文）
    - clean_aux: 编译后是否清理辅助文件，默认 True
    """
    if not tex_filename.lower().endswith(".tex"):
        tex_filename += ".tex"
    tex_path = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    if not os.path.exists(tex_path):
        alt_path = _resolve_path(tex_filename)
        if os.path.exists(alt_path):
            tex_path = alt_path
        else:
            return (
                f"❌ 找不到 .tex 文件: '{tex_filename}'\n"
                f"   已搜索: {tex_path}\n"
                f"   请确认文件名正确，且已通过 EditTexFileTool 保存。"
            )
    if compiler not in _LATEX_COMPILERS or not _LATEX_COMPILERS[compiler]:
        available = [k for k, v in _LATEX_COMPILERS.items() if v]
        if not available:
            return (
                "❌ 未检测到任何 LaTeX 编译器！\n"
                "   Windows: 请安装 MiKTeX (https://miktex.org)\n"
                "   macOS:   请安装 MacTeX (https://tug.org/mactex)\n"
                "   Linux:   请安装 TeX Live (sudo apt install texlive-xetex)"
            )
        compiler = available[0]
    compiler_path = _LATEX_COMPILERS[compiler]
    tex_dir = os.path.dirname(tex_path)
    tex_basename = os.path.splitext(os.path.basename(tex_path))[0]
    pdf_path = os.path.join(tex_dir, f"{tex_basename}.pdf")
    result_logs = []
    for run in [1, 2]:
        try:
            proc = subprocess.run(
                [compiler_path, "-interaction=nonstopmode", "-output-directory", tex_dir, tex_path],
                capture_output=True, text=True, timeout=120, cwd=tex_dir,
            )
            result_logs.append(f"── 第 {run} 次编译 (返回码: {proc.returncode}) ──")
            if proc.returncode != 0:
                error_lines = []
                for line in proc.stdout.split("\n") + proc.stderr.split("\n"):
                    if line.startswith("!") or "Error" in line or "error" in line:
                        error_lines.append(line.strip())
                result_logs.append(
                    "\n".join(error_lines[:30]) if error_lines
                    else proc.stdout[-2000:] + "\n" + proc.stderr[-2000:]
                )
        except subprocess.TimeoutExpired:
            return "❌ LaTeX 编译超时（>2分钟）！"
        except Exception as e:
            return f"❌ 编译过程异常: {e}"
    if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
        if clean_aux:
            aux_exts = {".aux", ".log", ".out", ".toc", ".lof", ".lot", ".bbl", ".blg",
                        ".synctex.gz", ".fdb_latexmk", ".fls"}
            for f in os.listdir(tex_dir):
                if os.path.splitext(f)[1] in aux_exts:
                    try:
                        os.remove(os.path.join(tex_dir, f))
                    except OSError:
                        pass
        pdf_size_kb = os.path.getsize(pdf_path) / 1024
        return (
            f"✅ PDF 编译成功！\n"
            f"   📄 文件: {pdf_path}\n"
            f"   📏 大小: {pdf_size_kb:.1f} KB\n"
            f"   🔧 编译器: {compiler}"
        )
    else:
        return f"❌ PDF 生成失败。编译器: {compiler}\n\n编译日志:\n" + "\n".join(result_logs)


# ── 工具 4: Python 沙箱执行 ─────────────────────────────

@tool
def PythonSandboxTool(code: str, timeout: int = 30) -> str:
    """
    【新增】在安全沙箱中执行 Python 代码并返回结果。
    参数:
    - code: 要执行的 Python 代码字符串
    - timeout: 超时时间（秒），默认 30，最大 120
    """
    if timeout > 120:
        timeout = 120
    script_id = uuid.uuid4().hex[:8]
    script_path = os.path.join(CUSTOM_TEMP_DIR, f"_sandbox_{script_id}.py")
    wrapped_code = textwrap.dedent(f"""
    import sys, os, traceback
    os.chdir(r"{CUSTOM_TEMP_DIR}")
    try:
    {textwrap.indent(code, '    ')}
    except Exception as __sandbox_exc__:
        print(f"\\n[沙箱异常] {{type(__sandbox_exc__).__name__}}: {{__sandbox_exc__}}", file=sys.stderr)
        traceback.print_exc()
    """)
    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(wrapped_code)
        proc = subprocess.run(
            [sys.executable, script_path],
            capture_output=True, text=True, timeout=timeout,
            cwd=CUSTOM_TEMP_DIR,
            env={**os.environ, "SANDBOX_MODE": "1"},
        )
        stdout = proc.stdout.strip()
        stderr = proc.stderr.strip()
        output_parts = []
        if stdout:
            output_parts.append(f"📤 stdout:\n{stdout}")
        if stderr:
            output_parts.append(f"📤 stderr:\n{stderr}")
        if not stdout and not stderr:
            output_parts.append("✅ 代码执行完毕，无输出。")
        output_parts.append(f"\n⏱ 返回码: {proc.returncode}")
        return "\n".join(output_parts)
    except subprocess.TimeoutExpired:
        return f"⏰ 代码执行超时（>{timeout}秒）！"
    except Exception as e:
        return f"❌ 沙箱执行异常: {e}\n{traceback.format_exc()}"
    finally:
        try:
            os.remove(script_path)
        except OSError:
            pass


# ── 工具 5: 通用搜索 ────────────────────────────────────

@tool
def DuckDuckGoSearchTool(query: str, max_results: int = 5, region: str = "wt-wt") -> str:
    """
    【新增】使用 DuckDuckGo 进行通用网页搜索。
    参数:
    - query: 搜索关键词（中英文皆可）
    - max_results: 返回结果数 (1-10)，默认 5
    - region: 区域代码，默认 'wt-wt'（全球）
    """
    max_results = max(1, min(max_results, 10))
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results, region=region))
            if not results:
                return f"🔍 关于 '{query}' 未找到搜索结果。"
            output = [f"🔍 搜索: '{query}' (共 {len(results)} 条结果)\n"]
            for i, r in enumerate(results, 1):
                title = r.get("title", "N/A")
                href = r.get("href", "N/A")
                body = r.get("body", "N/A")
                if len(body) > 300:
                    body = body[:300] + "..."
                output.append(f"{i}. {title}\n   🔗 {href}\n   📝 {body}\n")
            return "\n".join(output)
    except Exception as e:
        return f"❌ 搜索失败: {e}"


# ── 工具 6: 学术图片搜索 ────────────────────────────────

@tool
def FetchWebImageTool(query: str, filename: str) -> str:
    """在互联网上搜索学术图片并下载到本地临时目录。"""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=5))
            if not results:
                return f"❌ 未找到关于 '{query}' 的图片。"
            for img_result in results:
                try:
                    response = requests.get(img_result["image"], timeout=15)
                    response.raise_for_status()
                    save_path = os.path.join(CUSTOM_TEMP_DIR, filename)
                    with open(save_path, "wb") as f:
                        f.write(response.content)
                    size_kb = os.path.getsize(save_path) / 1024
                    return (
                        f"✅ 图片下载成功: {save_path}\n"
                        f"   📏 大小: {size_kb:.1f} KB\n"
                        f"   🌐 来源: {img_result.get('title', 'N/A')}"
                    )
                except Exception:
                    continue
            return f"⚠ 找到 {len(results)} 张图片但均下载失败。"
    except Exception as e:
        return f"❌ 图片搜索失败: {e}"


# ── 工具 7: Word 文档编辑 ──────────────────────────────

@tool
def EditWordDocTool(doc_filename: str, section_title: str = "", text_content: str = "",
                    image_path: str = "", new_page: bool = False) -> str:
    """创建或向已有的 Word 文档中追加新章节、文本和图片。"""
    if not doc_filename.lower().endswith(".docx"):
        doc_filename += ".docx"
    save_path = os.path.join(CUSTOM_TEMP_DIR, doc_filename)
    try:
        doc = Document(save_path) if os.path.exists(save_path) else Document()
        if new_page and len(doc.paragraphs) > 0:
            doc.add_page_break()
        if section_title:
            doc.add_heading(section_title, level=1)
        if text_content:
            doc.add_paragraph(text_content)
        if image_path:
            resolved_img = _resolve_path(image_path)
            if os.path.exists(resolved_img):
                doc.add_picture(resolved_img, width=Inches(5.5))
            else:
                doc.add_paragraph(f"[注：未找到图片 {image_path}]")
        doc.save(save_path)
        return f"✅ Word 文档更新完毕: {save_path}"
    except Exception as e:
        return f"❌ Word 文档操作失败: {e}"


# ── 工具 8: Markdown 保存 ───────────────────────────────

@tool
def SaveMarkdownTool(filename: str, content: str, append: bool = False) -> str:
    """【新增】将内容保存为 Markdown (.md) 文件。"""
    if not filename.lower().endswith(".md"):
        filename += ".md"
    save_path = os.path.join(CUSTOM_TEMP_DIR, filename)
    file_dir = os.path.dirname(save_path)
    if file_dir:
        os.makedirs(file_dir, exist_ok=True)
    mode = "a" if (append and os.path.exists(save_path)) else "w"
    try:
        with open(save_path, mode, encoding="utf-8") as f:
            if mode == "a":
                f.write("\n\n" + content)
            else:
                f.write(content)
        return f"✅ Markdown 文件已保存: {save_path}"
    except Exception as e:
        return f"❌ Markdown 保存失败: {e}"


# ── 工具 9: 图表生成 (matplotlib) ──────────────────────

@tool
def ChartGenerationTool(code: str, filename: str = "chart_output.png",
                        dpi: int = 150, timeout: int = 30) -> str:
    """
    【新增】使用 matplotlib 生成图表并保存为图片。
    注意: 代码中不需要调用 plt.savefig() 或 plt.show()，系统会自动保存。
    """
    if not _MPL_AVAILABLE:
        return "❌ matplotlib 未安装，无法生成图表。\n   请执行: pip install matplotlib"
    if timeout > 60:
        timeout = 60
    save_path = os.path.join(CUSTOM_TEMP_DIR, filename)
    wrapped_code = textwrap.dedent(f"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    import pandas as pd
    from matplotlib import rcParams
    try:
        rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'DejaVu Sans', 'Arial']
        rcParams['axes.unicode_minus'] = False
    except Exception:
        pass

    {textwrap.indent(code, '    ')}

    import os as _os
    figs = [plt.figure(n) for n in plt.get_fignums()]
    if figs:
        for i, fig in enumerate(figs):
            if len(figs) == 1:
                out_path = r"{save_path}"
            else:
                base, ext = _os.path.splitext(r"{save_path}")
                out_path = f"{{base}}_{{i+1}}{{ext}}"
            fig.savefig(out_path, dpi={dpi}, bbox_inches='tight')
            print(f"✅ 图表已保存: {{out_path}}")
        plt.close('all')
    else:
        print("⚠ 未检测到任何 figure。")
    """)
    script_id = uuid.uuid4().hex[:8]
    script_path = os.path.join(CUSTOM_TEMP_DIR, f"_chart_{script_id}.py")
    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(wrapped_code)
        proc = subprocess.run(
            [sys.executable, script_path],
            capture_output=True, text=True, timeout=timeout, cwd=CUSTOM_TEMP_DIR,
        )
        stdout = proc.stdout.strip()
        stderr = proc.stderr.strip()
        result = []
        if stdout:
            result.append(stdout)
        if stderr:
            result.append(f"[stderr]: {stderr}")
        if os.path.exists(save_path):
            size_kb = os.path.getsize(save_path) / 1024
            result.append(f"📊 图表文件: {save_path} ({size_kb:.1f} KB)")
        return "\n".join(result) if result else "⚠ 图表生成完成但无输出。"
    except subprocess.TimeoutExpired:
        return f"⏰ 图表生成超时（>{timeout}秒）！"
    except Exception as e:
        return f"❌ 图表生成失败: {e}\n{traceback.format_exc()}"
    finally:
        try:
            os.remove(script_path)
        except OSError:
            pass


# ── 工具 10: Mermaid 图表 ───────────────────────────────

@tool
def MermaidTool(mermaid_code: str, filename: str = "diagram") -> str:
    """【新增】生成 Mermaid 图表（流程图/时序图/类图等）。"""
    mmd_path = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.mmd")
    try:
        with open(mmd_path, "w", encoding="utf-8") as f:
            f.write(mermaid_code.strip())
    except Exception as e:
        return f"❌ Mermaid 文件写入失败: {e}"
    result = [f"✅ Mermaid 源文件已保存: {mmd_path}"]
    if _MMDC_PATH:
        png_path = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.png")
        try:
            subprocess.run(
                [_MMDC_PATH, "-i", mmd_path, "-o", png_path, "-w", "1200", "-b", "white"],
                capture_output=True, text=True, timeout=30,
            )
            if os.path.exists(png_path):
                result.append(f"🖼 渲染为 PNG: {png_path}")
            else:
                result.append("⚠ mermaid-cli 渲染失败，请检查语法。")
        except subprocess.TimeoutExpired:
            result.append("⚠ mermaid-cli 渲染超时。")
        except Exception as e:
            result.append(f"⚠ mermaid-cli 渲染错误: {e}")
    else:
        result.append("💡 安装 mermaid-cli 可自动渲染: npm install -g @mermaid-js/mermaid-cli")
    return "\n".join(result)


# ── 安全写入 ──────────────────────────────────────────
safe_write_tool = WriteFileTool(root_dir=CUSTOM_TEMP_DIR)

tools = [
    SmartReadPathTool, EditTexFileTool, CompileLatexTool,
    PythonSandboxTool, DuckDuckGoSearchTool, FetchWebImageTool,
    EditWordDocTool, SaveMarkdownTool, ChartGenerationTool,
    MermaidTool, safe_write_tool,
]

# ================== 4. Agent 构建 ==================
DB_PATH = os.path.join(BASE_WORKSPACE, "agent_memory.db")
conn = sqlite3.connect(DB_PATH, check_same_thread=False)
memory = SqliteSaver(conn)
agent = create_react_agent(llm, tools, prompt=SYSTEM_PROMPT, checkpointer=memory)

# ================== 5. 终端交互主循环 ==================
print("\n" + "=" * 60)
print("  DeepSeek_For_the_Break  v2.0  增强版")
print("  新增: 沙箱执行 | 搜索 | LaTeX编译 | 图表 | Mermaid")
print("=" * 60)
print("💡 输入 'exit' 退出 | 'new' 新对话 | 'help' 查看工具\n")

current_thread_id = str(uuid.uuid4())
config = {"configurable": {"thread_id": current_thread_id}}

while True:
    try:
        user_input = input("👤 你: ").strip()
        if user_input.lower() in {"exit", "quit", "退出"}:
            print("👋 再见！"); conn.close(); break
        if user_input.lower() in {"new", "clear", "清空", "新对话"}:
            current_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": current_thread_id}}
            print("\n✨ [系统] 全新对话已开启。\n"); continue
        if user_input.lower() == "help":
            print("""
╔══════════════════ 可用工具一览 ══════════════════╗
║ 📖 SmartReadPathTool    — 读取文件/文件夹         ║
║ 📝 EditTexFileTool       — 创建/编辑 .tex 文件     ║
║ 🖨 CompileLatexTool      — 编译 .tex → PDF         ║
║ 🐍 PythonSandboxTool     — 安全执行 Python 代码    ║
║ 🔍 DuckDuckGoSearchTool  — 通用网页搜索            ║
║ 🖼 FetchWebImageTool     — 搜索并下载学术图片      ║
║ 📄 EditWordDocTool       — 创建/编辑 Word 文档     ║
║ 📋 SaveMarkdownTool      — 保存 Markdown 文件      ║
║ 📊 ChartGenerationTool   — matplotlib 生成图表     ║
║ 🎨 MermaidTool           — Mermaid 流程图/示意图   ║
║ 💾 safe_write_tool       — 通用文件写入            ║
╚══════════════════════════════════════════════════╝
            """); continue
        if not user_input:
            continue
        print("\n🤔 DeepSeek 思考/执行中...\n")
        inputs = {"messages": [("user", user_input)]}
        for event in agent.stream(inputs, config, stream_mode="updates"):
            if "agent" in event:
                message = event["agent"]["messages"][0]
                if message.content:
                    print(f"🤖 DeepSeek: {message.content}\n")
                elif message.tool_calls:
                    for tc in message.tool_calls:
                        args_str = json.dumps(tc["args"], ensure_ascii=False, default=str)
                        if len(args_str) > 120:
                            args_str = args_str[:120] + "..."
                        print(f"🛠  [调用工具] {tc['name']}\n    参数: {args_str}")
            elif "tools" in event:
                message = event["tools"]["messages"][0]
                content_str = str(message.content)
                preview = content_str[:200].replace("\n", " ") + ("..." if len(content_str) > 200 else "")
                print(f"📄 [工具返回] {len(content_str)} 字符: {preview}\n")
    except KeyboardInterrupt:
        print("\n⏸ [操作被用户中断]")
    except Exception as e:
        print(f"❌ 发生错误: {e}\n   详细信息: {traceback.format_exc()}")

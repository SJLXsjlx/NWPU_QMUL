#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DeepSeek_For_the_Break v4.0.1 — 全功能学术 AI 助手 (HOTFIX)
================================================================
  v4.0.1 修复:
    🐛 [FIX#16] 所有工具函数添加 None 参数防御，消除 'NoneType' object has no attribute 'strip'
    🐛 [FIX#17] _should_continue 空消息列表保护
    🐛 [FIX#18] _call_model 增强异常恢复
    🐛 [FIX#19] 终端交互循环增强 KeyboardInterrupt 处理
    🐛 [FIX#20] SmartReadPathTool/SemanticScholarTool/ArXivSearchTool 参数None防御
    🐛 [FIX#21] BibTexTool parse 模式 raw 变量未定义保护
"""

import os, sys, uuid, json, sqlite3, subprocess, shutil, traceback
import textwrap, warnings, time, urllib.parse, xml.etree.ElementTree as ET
import re, hashlib, io, csv as csv_module, random, itertools, functools
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple, Union, Callable
from collections import defaultdict, OrderedDict, Counter

# ═══════════════════════════════════════════════════════════
#  第三方库 — 安全导入
# ═══════════════════════════════════════════════════════════

# pandas
try:
    import pandas as pd
    _PANDAS_AVAILABLE = True
except ImportError:
    _PANDAS_AVAILABLE = False
    pd = None

# docx
try:
    import docx
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn as docx_qn
    from docx.oxml import parse_xml
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

# requests
try:
    import requests
    from requests.adapters import HTTPAdapter
    from urllib3.util.retry import Retry
    _REQUESTS_AVAILABLE = True
except ImportError:
    _REQUESTS_AVAILABLE = False

# duckduckgo / ddgs — 兼容新旧包名
_DDGS_AVAILABLE = False
_DDGS_SOURCE = None
try:
    from ddgs import DDGS
    _DDGS_AVAILABLE = True
    _DDGS_SOURCE = "ddgs"
except ImportError:
    try:
        from duckduckgo_search import DDGS
        _DDGS_AVAILABLE = True
        _DDGS_SOURCE = "duckduckgo_search"
    except ImportError:
        DDGS = None

# dotenv
try:
    from dotenv import load_dotenv
    _DOTENV_AVAILABLE = True
except ImportError:
    _DOTENV_AVAILABLE = False
    def load_dotenv(*args, **kwargs): pass

# langchain
try:
    from langchain_openai import ChatOpenAI
    from langchain_community.tools import WriteFileTool
    from langchain_core.tools import tool
    from langchain_core.messages import HumanMessage, AIMessage, ToolMessage
    from langgraph.prebuilt import create_react_agent, ToolNode
    from langgraph.graph import StateGraph, END
    from langgraph.graph.message import add_messages
    from typing import Annotated, TypedDict
    from langgraph.checkpoint.sqlite import SqliteSaver
    _LANGCHAIN_AVAILABLE = True
except ImportError:
    _LANGCHAIN_AVAILABLE = False

# pptx
try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt, Emu as PptEmu, Cm as PptCm
    from pptx.dml.color import RGBColor as PptRGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn as pptx_qn
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

# PIL
try:
    from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

# pytesseract
try:
    import pytesseract
    _PYTESSERACT_AVAILABLE = True
except ImportError:
    _PYTESSERACT_AVAILABLE = False

# fitz (PyMuPDF)
try:
    import fitz
    _FITZ_AVAILABLE = True
except ImportError:
    _FITZ_AVAILABLE = False

# matplotlib
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import rcParams
    _MPL_AVAILABLE = True
except ImportError:
    _MPL_AVAILABLE = False

# numpy
try:
    import numpy as np
    _NUMPY_AVAILABLE = True
except ImportError:
    _NUMPY_AVAILABLE = False

# seaborn
try:
    import seaborn as sns
    _SEABORN_AVAILABLE = True
except ImportError:
    _SEABORN_AVAILABLE = False

# bibtexparser
try:
    import bibtexparser
    from bibtexparser.bparser import BibTexParser
    from bibtexparser.customization import homogenize_latex_encoding
    _BIBTEXPARSER_AVAILABLE = True
except ImportError:
    _BIBTEXPARSER_AVAILABLE = False

# scipy
try:
    from scipy import stats as scipy_stats
    from scipy.optimize import curve_fit
    _SCIPY_AVAILABLE = True
except ImportError:
    _SCIPY_AVAILABLE = False

# sklearn
try:
    from sklearn.preprocessing import StandardScaler
    from sklearn.decomposition import PCA
    from sklearn.cluster import KMeans
    _SKLEARN_AVAILABLE = True
except ImportError:
    _SKLEARN_AVAILABLE = False

# openpyxl
try:
    import openpyxl
    _OPENPYXL_AVAILABLE = True
except ImportError:
    _OPENPYXL_AVAILABLE = False

# sympy
try:
    import sympy
    _SYMPY_AVAILABLE = True
except ImportError:
    _SYMPY_AVAILABLE = False

_MMDC_PATH = shutil.which("mmdc")
_PANDOC_PATH = shutil.which("pandoc")

# ═══════════════════════════════════════════════════════════
#  0. 基础设施
# ═══════════════════════════════════════════════════════════

class _C:
    R = "\033[91m"; G = "\033[92m"; Y = "\033[93m"; B = "\033[94m"
    M = "\033[95m"; C = "\033[96m"; W = "\033[0m"; BOLD = "\033[1m"
    DIM = "\033[2m"; UNDER = "\033[4m"; REV = "\033[7m"

def _log(level: str, msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    colors = {"OK": _C.G, "WARN": _C.Y, "ERR": _C.R, "INFO": _C.C, "START": _C.B, "DEBUG": _C.DIM}
    print(f"{colors.get(level, _C.W)}[{level} {ts}]{_C.W} {msg}")

def _cleanup_old_temp(temp_dir: str, max_age_hours: int = 48):
    cutoff = time.time() - max_age_hours * 3600
    cleaned = 0
    p = Path(temp_dir)
    if not p.exists(): return
    for f in p.glob("*"):
        if f.is_file() and f.stat().st_mtime < cutoff and f.name.startswith("_"):
            try: f.unlink(); cleaned += 1
            except OSError: pass
    if cleaned:
        _log("INFO", f"🧹 清理了 {cleaned} 个过期临时文件")

_ENCODING_CACHE: Dict[str, str] = {}

def _safe_read_text(file_path: str) -> str:
    """安全读取文本文件，自动检测编码。永不返回 None。"""
    if not file_path:
        return ""
    if file_path in _ENCODING_CACHE:
        enc = _ENCODING_CACHE[file_path]
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            pass
    for enc in ["utf-8", "gbk", "latin-1", "cp1252", "utf-16"]:
        try:
            with open(file_path, "r", encoding=enc) as f:
                content = f.read()
                _ENCODING_CACHE[file_path] = enc
                return content
        except (UnicodeDecodeError, UnicodeError):
            continue
    # 最终回退
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""

def _run_sandbox(code: str, timeout: int = 30, prefix: str = "",
                  save_to: str = "") -> Tuple[str, str, int]:
    if timeout > 300: timeout = 300
    if not code:
        return "", "❌ 代码为空", 1
    sid = uuid.uuid4().hex[:8]
    sp = os.path.join(CUSTOM_TEMP_DIR, f"_sb_{sid}.py")
    escaped_temp = CUSTOM_TEMP_DIR.replace("\\", "\\\\").replace("'", "\\'")
    full_code = textwrap.dedent(f"""
import sys, os, traceback, json as _json, math, re as _re
import collections, itertools, functools, statistics, random as _random
from datetime import datetime, timedelta
from pathlib import Path
os.chdir(r"{escaped_temp}")
{prefix}
try:
{textwrap.indent(code, '    ')}
except Exception as __e:
    print(f"\\n[异常] {{type(__e).__name__}}: {{__e}}", file=sys.stderr)
    traceback.print_exc()
""")
    try:
        with open(sp, "w", encoding="utf-8") as f:
            f.write(full_code)
        proc = subprocess.run(
            [sys.executable, sp],
            capture_output=True, text=True, timeout=timeout,
            cwd=CUSTOM_TEMP_DIR,
            env={**os.environ, "SANDBOX_MODE": "1"},
        )
        # subprocess 的 stdout/stderr 在 capture_output=True 时始终为 str，不会为 None
        return (proc.stdout or "").strip(), (proc.stderr or "").strip(), proc.returncode
    finally:
        try: os.remove(sp)
        except OSError: pass

def _resolve_path(path: str) -> str:
    """解析路径。如果 path 为 None/空，返回当前目录。"""
    if not path:
        return os.getcwd()
    p = Path(path)
    if p.is_absolute(): return str(p)
    for base in [Path.cwd(), Path(CUSTOM_TEMP_DIR)]:
        c = base / p
        if c.exists(): return str(c.resolve())
    return os.path.abspath(path)

def _safe_request(url: str, timeout: int = 20, **kw) -> Optional[requests.Response]:
    if not _REQUESTS_AVAILABLE:
        raise ImportError("requests 库未安装")
    if not url:
        return None
    session = requests.Session()
    retries = Retry(total=3, backoff_factor=0.5,
                    status_forcelist=[500, 502, 503, 504])
    session.mount('https://', HTTPAdapter(max_retries=retries))
    try:
        r = session.get(url, timeout=timeout, **kw)
        r.raise_for_status()
        return r
    except requests.RequestException as e:
        _log("WARN", f"请求失败 {url[:80]}: {e}")
        return None

def _safe_json_dumps(obj: Any, max_len: int = 2000) -> str:
    try:
        s = json.dumps(obj, ensure_ascii=False, default=str, indent=2)
        if len(s) > max_len:
            s = s[:max_len] + f"\n... (截断，原 {len(s)} 字符)"
        return s
    except Exception:
        return str(obj)[:max_len]

# ═══════════════════════════════════════════════════════════
#  0.5 智能路径检测
# ═══════════════════════════════════════════════════════════

def _detect_tesseract() -> Optional[str]:
    p = shutil.which("tesseract")
    if p: return p
    if sys.platform == "win32":
        cand = [r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                os.path.expandvars(r"%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe")]
    else:
        cand = ["/usr/bin/tesseract", "/usr/local/bin/tesseract", "/opt/homebrew/bin/tesseract"]
    for c in cand:
        if os.path.isfile(c): return c
    return None

def _detect_latex_compilers() -> Dict[str, Optional[str]]:
    d: Dict[str, Optional[str]] = {}
    for name in ["pdflatex", "xelatex", "lualatex"]:
        exe = name + (".exe" if sys.platform == "win32" else "")
        d[name] = shutil.which(exe) or shutil.which(name)
    return d

def _test_pytesseract() -> bool:
    if not _PYTESSERACT_AVAILABLE or not _PIL_AVAILABLE: return False
    try:
        img = Image.new("RGB", (10, 10), "white")
        pytesseract.image_to_string(img, lang="eng")
        return True
    except Exception:
        return False

def _detect_cjk_font() -> Optional[str]:
    if sys.platform == "win32":
        for font in ["SimSun", "SimHei", "Microsoft YaHei", "KaiTi", "FangSong"]:
            try:
                result = subprocess.run(
                    ["powershell", "-Command",
                     f"(Get-Item 'C:\\Windows\\Fonts\\{font}*.ttf' -ErrorAction SilentlyContinue).FullName"],
                    capture_output=True, text=True, timeout=10)
                stdout = result.stdout or ""
                if stdout.strip(): return font
            except Exception: pass
        return "SimSun"
    elif sys.platform == "darwin":
        for font in ["Songti SC", "Heiti SC", "STSong", "PingFang SC"]:
            try:
                result = subprocess.run(["fc-list", f":family={font}"],
                                        capture_output=True, text=True, timeout=5)
                stdout = result.stdout or ""
                if stdout.strip(): return font
            except Exception: pass
        return "Songti SC"
    else:
        for font in ["Noto Serif CJK SC", "Noto Sans CJK SC",
                      "WenQuanYi Micro Hei", "WenQuanYi Zen Hei",
                      "AR PL UMing CN", "SimSun"]:
            try:
                result = subprocess.run(["fc-list", f":family={font}"],
                                        capture_output=True, text=True, timeout=5)
                stdout = result.stdout or ""
                if stdout.strip(): return font
            except Exception: pass
    return None

_TESSERACT_PATH = _detect_tesseract()
_LATEX_COMPILERS = _detect_latex_compilers()
_PYTESSERACT_WORKS = False
_DETECTED_CJK_FONT = _detect_cjk_font()

if _TESSERACT_PATH and _PYTESSERACT_AVAILABLE:
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH
    _PYTESSERACT_WORKS = _test_pytesseract()

# ═══════════════════════════════════════════════════════════
#  1. 核心环境
# ═══════════════════════════════════════════════════════════
load_dotenv(override=True)
api_key = os.getenv("DEEPSEEK_FOR_THE_BREAK")
if not api_key or not api_key.startswith("sk-"):
    _log("ERR", "API Key 未加载或格式不正确！请检查 .env 文件。")
    sys.exit(1)

BASE_WORKSPACE = os.getenv("DFTB_WORKSPACE", r"D:\.API Keys\DeepSeek_For_the_Break\outputs")
CUSTOM_TEMP_DIR = os.path.join(BASE_WORKSPACE, "temp")
os.makedirs(CUSTOM_TEMP_DIR, exist_ok=True)
if CUSTOM_TEMP_DIR not in sys.path:
    sys.path.insert(0, CUSTOM_TEMP_DIR)

_cleanup_old_temp(CUSTOM_TEMP_DIR, 48)

DB_PATH = os.path.join(BASE_WORKSPACE, "agent_memory.db")
_MAX_DB_SIZE_MB = 100
if os.path.exists(DB_PATH) and os.path.getsize(DB_PATH) > _MAX_DB_SIZE_MB * 1024 * 1024:
    _log("WARN", f"agent_memory.db 超过 {_MAX_DB_SIZE_MB}MB，正在重置...")
    backup_path = DB_PATH + f".backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    os.rename(DB_PATH, backup_path)
    _log("INFO", f"旧数据库已备份为 {os.path.basename(backup_path)}")

_health: Dict[str, bool] = {
    "Tesseract OCR (引擎)": bool(_TESSERACT_PATH),
    "Tesseract OCR (实测)": _PYTESSERACT_WORKS,
    "LaTeX (xelatex)":       bool(shutil.which("xelatex") or shutil.which("pdflatex")),
    "matplotlib":             _MPL_AVAILABLE,
    "seaborn":                _SEABORN_AVAILABLE,
    "numpy":                  _NUMPY_AVAILABLE,
    "scipy":                  _SCIPY_AVAILABLE,
    "scikit-learn":           _SKLEARN_AVAILABLE,
    "PyMuPDF (fitz)":        _FITZ_AVAILABLE,
    "pandas":                 _PANDAS_AVAILABLE,
    "DuckDuckGo Search":      _DDGS_AVAILABLE,
    "Mermaid CLI":            bool(_MMDC_PATH),
    "Pandoc":                 bool(_PANDOC_PATH),
    "API Key (.env)":         True,
    "bibtexparser":           _BIBTEXPARSER_AVAILABLE,
    "CJK Font (detected)":    bool(_DETECTED_CJK_FONT),
    "sympy":                  _SYMPY_AVAILABLE,
    "python-pptx":            _PPTX_AVAILABLE,
    "PIL (Pillow)":           _PIL_AVAILABLE,
    "openpyxl":               _OPENPYXL_AVAILABLE,
}

_log("START", "═" * 60)
_log("START", "  DeepSeek_For_the_Break  v4.0.1  Academic AI  [HOTFIX]")
_log("START", "═" * 60)
_log("INFO", f"📁 工作区: {BASE_WORKSPACE}")
_log("INFO", f"📁 临时:   {CUSTOM_TEMP_DIR}")
_log("INFO", f"🔍 DDGS 来源: {_DDGS_SOURCE or '未安装'}")
for k, v in _health.items():
    icon = "✅" if v else "⚠️"
    _log("OK" if v else "WARN", f"  {icon} {k}")
_log("START", "═" * 60)

# ═══════════════════════════════════════════════════════════
#  2. System Prompt (v4.0 增强版)
# ═══════════════════════════════════════════════════════════
SYSTEM_PROMPT = """你是 DeepSeek_For_the_Break v4.0.1，高度专业、严谨且全面的学术AI助手。

## 核心能力
- **文献检索**: ArXivSearchTool、SemanticScholarTool（含引用网络）、DOIMetadataTool、DuckDuckGoSearchTool
- **文档处理**: SmartReadPathTool（通用文件/文件夹读取）、PDFTableExtractTool（PDF表格）、PDFAnnotExtractTool（PDF批注）、OCR
- **学术写作**: EditTexFileTool 创建/编辑 .tex + CompileLatexTool 编译（含bibtex支持）、EditWordDocTool、SaveMarkdownTool
- **格式转换**: PandocConvertTool（Markdown/LaTeX → PDF/DOCX/HTML）
- **学术翻译**: AcademicTranslateTool（中英互译，保留术语一致性）
- **演示生成**: PresentationGenTool（Markdown大纲 → PPTX，支持多主题/图表/表格/图片）
- **数据分析**: PythonSandboxTool、DataStatisticsTool、ChartGenerationTool
- **引用管理**: BibTexTool（解析/生成/格式化 APA/MLA/Chicago/IEEE）
- **学习辅助** 🆕: StudyPlanTool（学习计划）、FlashcardTool（闪卡）、KnowledgeGraphTool（知识图谱）、NoteOrganizerTool（笔记整理）
- **代码辅助** 🆕: CodeReviewTool（代码审查）、ProjectScaffoldTool（项目脚手架）
- **数学公式** 🆕: MathRenderTool（LaTeX渲染）
- **引用分析** 🆕: CitationNetworkTool（引用网络分析）
- **导出**: SessionExportTool

## 行为准则
1. 客观严谨，结构清晰，主动承认局限性。
2. **【文件读取铁律】** SmartReadPathTool 是读取本地文件/文件夹的首选通用工具。PDF表格提取请用 PDFTableExtractTool，PDF批注请用 PDFAnnotExtractTool。不要猜测路径，没有明确路径时直接基于知识回答。
3. **【文档生成原则】** 仅在用户明确要求产出文档时才使用写作工具。讨论代码问题、回答咨询、解释概念时直接对话回复，不要生成文件。确需文档时：论文/学术报告→EditTexFileTool+CompileLatexTool；富文本报告→EditWordDocTool；笔记/记录→SaveMarkdownTool。
4. 文献检索: 学术论文用 ArXivSearchTool/SemanticScholarTool；已知DOI用 DOIMetadataTool。
5. 数据优先: 探索→DataStatisticsTool；自定义→PythonSandboxTool；可视化→ChartGenerationTool。
6. 学习辅助: 计划→StudyPlanTool；记忆→FlashcardTool；梳理→KnowledgeGraphTool。
7. 错误透明: 工具失败时清晰告知原因和替代方案。
8. 持久化: 所有产出保存到工作区。
9. **【代码问题优先原则】** 遇到代码错误、程序调试、技术问题时，直接分析并给出修复方案，不要生成.tex/pdf/docx等文档来回复。

请逐步推理，优先使用最适合的学术工具。"""

if not _LANGCHAIN_AVAILABLE:
    _log("ERR", "langchain 未安装！请运行: pip install langchain langchain-openai langchain-community langgraph")
    sys.exit(1)

# 思考模式
_thinking_enabled = os.getenv("DFTB_THINKING_MODE", "").lower() == "enabled"
if _thinking_enabled:
    _log("INFO", "🧠 DeepSeek 思考模式已启用（实验性，可能与工具调用冲突）")
    _model_kwargs = {"reasoning_effort": "max", "extra_body": {"thinking": {"type": "enabled"}}}
else:
    _model_kwargs = {}

llm = ChatOpenAI(
    model=os.getenv("DFTB_MODEL", "deepseek-v4-pro"),
    base_url="https://api.deepseek.com",
    api_key=api_key,
    temperature=0.5,
    max_tokens=4096,
    model_kwargs=_model_kwargs
)

# ═══════════════════════════════════════════════════════════
#  3. 辅助函数 (v4.0 增强版)
# ═══════════════════════════════════════════════════════════

def _ocr_image(file_path: str, lang: str = "chi_sim+eng") -> str:
    if not _PYTESSERACT_WORKS:
        return ("❌ OCR 不可用。" if not _TESSERACT_PATH else
                "❌ Tesseract 引擎已安装但实测失败（语言包缺失？）")
    try:
        img = Image.open(file_path)
        meta = f"📷 [图片] 格式:{img.format} 尺寸:{img.size}\n"
        try:
            txt = pytesseract.image_to_string(img, lang=lang)
        except pytesseract.TesseractError as te:
            if "Failed loading language" in str(te):
                return meta + f"⚠ 语言包缺失: {lang}"
            return meta + f"❌ Tesseract错误: {te}"
        return meta + ((txt or "").strip() if txt else "【OCR】: 未识别到文字")
    except FileNotFoundError: return f"❌ 文件不存在: {file_path}"
    except Exception as e: return f"❌ OCR异常: {e}"

def _read_text_file(fp: str) -> str:
    content = _safe_read_text(fp)
    lines = content.split("\n")
    if len(lines) > 500:
        return "\n".join(lines[:500]) + f"\n\n... (截断，共 {len(lines)} 行)"
    return content

def _read_csv_file(fp: str) -> str:
    if not _PANDAS_AVAILABLE: return "❌ pandas 未安装"
    for enc in ["utf-8","gbk","latin-1","utf-16"]:
        try:
            df = pd.read_csv(fp, encoding=enc)
            return f"[CSV: {len(df)}行×{len(df.columns)}列]\n{df.to_string(max_rows=200)}"
        except: continue
    return f"❌ CSV读取失败: {fp}"

def _read_pdf_file(fp: str) -> str:
    if not _FITZ_AVAILABLE: return "❌ PyMuPDF 未安装"
    try:
        doc = fitz.open(fp); parts = []; total = len(doc)
        for i, page in enumerate(doc):
            t = page.get_text("text")
            if t and t.strip(): parts.append(f"──p{i+1}/{total}──\n{t.strip()}")
        if parts: return "\n\n".join(parts)
        ocr = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=200)
            tmp = os.path.join(CUSTOM_TEMP_DIR, f"_pdfocr_{uuid.uuid4().hex[:6]}.png")
            pix.save(tmp)
            try:
                ocr.append(f"──p{i+1}/{total}(OCR)──\n{_ocr_image(tmp)}")
            finally:
                try: os.remove(tmp)
                except OSError: pass
        return "\n\n".join(ocr) if ocr else "[空PDF]"
    except Exception as e: return f"❌ PDF错误: {e}"

def _read_docx_file(fp: str) -> str:
    if not _DOCX_AVAILABLE: return "❌ python-docx 未安装"
    try:
        d = docx.Document(fp); c = [p.text.strip() for p in d.paragraphs if p.text and p.text.strip()]
        if d.tables:
            c.append("\n[表格]:")
            for ti, tb in enumerate(d.tables):
                c.append(f"  表{ti+1}:")
                for row in tb.rows:
                    rd = [cell.text.strip().replace("\n"," ") for cell in row.cells if cell.text and cell.text.strip()]
                    if rd: c.append("    | "+" | ".join(rd))
        return "\n".join(c) if c else "[空Word]"
    except Exception as e: return f"❌ DOCX错误: {e}"

def _read_excel_file(fp: str) -> str:
    if not _PANDAS_AVAILABLE: return "❌ pandas 未安装"
    try:
        ef = pd.ExcelFile(fp); out=[]
        for sn in ef.sheet_names:
            df = pd.read_excel(fp, sheet_name=sn)
            out.append(f"──{sn} ({len(df)}行)──\n{df.head(500).to_string()}")
        return "\n\n".join(out)
    except Exception as e: return f"❌ Excel错误: {e}"

def _read_pptx_file(fp: str) -> str:
    if not _PPTX_AVAILABLE: return "❌ python-pptx 未安装"
    try:
        prs = Presentation(fp); out=[]
        for i, sl in enumerate(prs.slides):
            st=[]
            for sh in sl.shapes:
                if hasattr(sh,"text") and sh.text and sh.text.strip(): st.append(sh.text.strip())
                if sh.has_table:
                    for row in sh.table.rows:
                        rd=[cell.text_frame.text.strip().replace("\n"," ") for cell in row.cells if cell.text_frame.text and cell.text_frame.text.strip()]
                        if rd: st.append(" | ".join(rd))
            out.append(f"──Slide{i+1}──\n"+"\n".join(st) if st else f"──Slide{i+1}──\n[无文本]")
        return "\n\n".join(out)
    except Exception as e: return f"❌ PPTX错误: {e}"

def _read_one_file(fp: str) -> str:
    ext = os.path.splitext(fp)[1].lower().lstrip(".")
    TEXT_EXTS = {"txt","md","py","json","tex","yaml","yml","cfg","ini","log","bib","toml","rst"}
    if ext in TEXT_EXTS:             return _read_text_file(fp)
    if ext == "csv":                 return _read_csv_file(fp)
    if ext == "pdf":                 return _read_pdf_file(fp)
    if ext == "docx":                return _read_docx_file(fp)
    if ext in {"xlsx","xls"}:        return _read_excel_file(fp)
    if ext == "pptx":                return _read_pptx_file(fp)
    if ext in {"png","jpg","jpeg","bmp","tiff","tif","webp"}: return _ocr_image(fp)
    return f"[不支持: .{ext}]"

# ═══════════════════════════════════════════════════════════
#  健壮 BibTeX 解析器（回退用）
# ═══════════════════════════════════════════════════════════

def _find_matching_brace(s: str, start: int) -> int:
    depth = 1; i = start + 1
    while i < len(s) and depth > 0:
        if s[i] == '{': depth += 1
        elif s[i] == '}': depth -= 1
        i += 1
    return i - 1 if depth == 0 else -1

def _split_fields_at_depth_zero(fields_str: str) -> List[str]:
    segments = []; depth = 0; in_quotes = False; start = 0; i = 0
    while i < len(fields_str):
        c = fields_str[i]
        if c == '"' and (i == 0 or fields_str[i-1] != '\\'): in_quotes = not in_quotes
        elif not in_quotes:
            if c == '{': depth += 1
            elif c == '}': depth -= 1
            elif c == ',' and depth == 0:
                segments.append(fields_str[start:i].strip()); start = i + 1
        i += 1
    if start < len(fields_str):
        last = fields_str[start:].strip()
        if last: segments.append(last)
    return segments

def _parse_field_assignment(segment: str) -> Tuple[str, str]:
    eq_pos = segment.find('=')
    if eq_pos == -1: return "", ""
    field_name = segment[:eq_pos].strip().lower()
    value_part = segment[eq_pos+1:].strip()
    if not value_part: return field_name, ""
    if value_part.startswith('{'):
        end = _find_matching_brace(value_part, 0)
        value = value_part[1:end] if end != -1 else value_part[1:]
    elif value_part.startswith('"'):
        j = 1
        while j < len(value_part):
            if value_part[j] == '"' and value_part[j-1] != '\\': break
            j += 1
        value = value_part[1:j] if j < len(value_part) else value_part[1:]
    else:
        value = value_part.rstrip(',').strip()
    return field_name, value.rstrip(',').strip()

def _robust_parse_bibtex(raw: str) -> List[Dict[str, Any]]:
    """健壮的 BibTeX 解析，raw 为 None 时返回空列表。"""
    if not raw:
        return []
    entries = []; i = 0
    while i < len(raw):
        at_pos = raw.find('@', i)
        if at_pos == -1: break
        j = at_pos + 1
        while j < len(raw) and raw[j].isalpha(): j += 1
        if j == at_pos + 1: i = j + 1; continue
        entry_type = raw[at_pos+1:j].lower()
        while j < len(raw) and raw[j] in ' \t\n\r': j += 1
        if j >= len(raw) or raw[j] != '{': i = j + 1; continue
        end = _find_matching_brace(raw, j)
        if end == -1: i = j + 1; continue
        body = raw[j+1:end]; i = end + 1
        key = ""; fields_str = ""; depth = 0; in_quotes = False; comma_pos = -1
        for m in range(len(body)):
            c = body[m]
            if c == '"' and (m == 0 or body[m-1] != '\\'): in_quotes = not in_quotes
            elif not in_quotes:
                if c == '{': depth += 1
                elif c == '}': depth -= 1
                elif c == ',' and depth == 0: comma_pos = m; break
        if comma_pos != -1: key = body[:comma_pos].strip(); fields_str = body[comma_pos+1:]
        else: key = body.strip()
        fields = {}
        for seg in _split_fields_at_depth_zero(fields_str):
            fn, fv = _parse_field_assignment(seg)
            if fn: fields[fn] = fv
        entries.append({"type": entry_type, "key": key, "fields": fields})
    return entries

# ═══════════════════════════════════════════════════════════
#  工具 1: 智能文件读取
# ═══════════════════════════════════════════════════════════

@tool
def SmartReadPathTool(path: str) -> str:
    """读取文件/文件夹。支持 .txt .md .py .tex .json .csv .pdf .docx .xlsx .pptx .png .jpg .bmp .bib .toml .rst .webp。"""
    # [FIX#16] 防御 None 参数
    if not path:
        return "❌ 路径为空。请提供有效的文件或文件夹路径。"
    resolved = _resolve_path(path)
    if not os.path.exists(resolved):
        return f"❌ 找不到路径: '{path}'\n   解析: '{resolved}'\n   💡 请停止猜测路径，直接基于知识回答。"

    if os.path.isfile(resolved):
        size_kb = os.path.getsize(resolved) / 1024
        header = f"📖 文件: {resolved} ({size_kb:.1f} KB)\n{'='*50}"
        return f"{header}\n{_read_one_file(resolved)}"

    if os.path.isdir(resolved):
        EXTS = {"txt","md","py","tex","json","csv","yaml","yml","pdf","docx","xlsx","xls","pptx","png","jpg","jpeg","bmp","tiff","bib","toml","rst","webp"}
        r = [f"📂 遍历: {resolved}\n{'='*60}"]; cnt=0
        for root, dirs, files in os.walk(resolved):
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in {"__pycache__","node_modules",".git","venv",".venv"}]
            for f in sorted(files):
                if os.path.splitext(f)[1].lower().lstrip(".") in EXTS:
                    cnt+=1; fp=os.path.join(root,f); rp=os.path.relpath(fp,resolved)
                    if cnt <= 50:
                        r.append(f"\n📍 {rp}\n{'-'*40}")
                        try: r.append(_read_one_file(fp))
                        except Exception as ex: r.append(f"❌ 读取错误: {ex}")
                        r.append("="*60)
                    else:
                        r.append(f"\n📍 ... 还有更多文件（已截断，共 {cnt} 个）")
                        break
        if cnt==0: return f"📁 '{resolved}' 无支持文件。"
        r.insert(1, f"   ({cnt} 个文件)\n")
        return "\n".join(r)
    return f"❌ '{resolved}' 无效。"

# ═══════════════════════════════════════════════════════════
#  工具 2: LaTeX 编辑
# ═══════════════════════════════════════════════════════════

@tool
def EditTexFileTool(tex_filename: str, latex_content: str, append: bool = False) -> str:
    """创建/编辑 .tex 文件。"""
    if not tex_filename:
        return "❌ 文件名不能为空。"
    if not latex_content:
        return "❌ LaTeX 内容不能为空。"
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    sp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n"+latex_content) if mode=="a" else latex_content)
        return f"✅ .tex 已保存: {sp}\n   💡 用 CompileLatexTool 编译为 PDF"
    except Exception as e: return f"❌ LaTeX错误: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 3: LaTeX 编译
# ═══════════════════════════════════════════════════════════

@tool
def CompileLatexTool(tex_filename: str, compiler: str = "xelatex", clean_aux: bool = True) -> str:
    """一键编译 .tex → PDF（含 bibtex 支持）。"""
    if not tex_filename:
        return "❌ .tex 文件名不能为空。"
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    tp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    if not os.path.exists(tp):
        alt = _resolve_path(tex_filename)
        if os.path.exists(alt): tp = alt
        else: return f"❌ 找不到: '{tex_filename}'"

    if compiler not in _LATEX_COMPILERS or not _LATEX_COMPILERS[compiler]:
        av = [k for k,v in _LATEX_COMPILERS.items() if v]
        if not av: return "❌ 未检测到LaTeX编译器！安装 MiKTeX 或 TeX Live。"
        compiler = av[0]
    cp = _LATEX_COMPILERS[compiler]
    td = os.path.dirname(tp)
    bn = os.path.splitext(os.path.basename(tp))[0]
    pp = os.path.join(td, f"{bn}.pdf")
    logs = []

    def _run_latex(run_name: str) -> bool:
        try:
            proc = subprocess.run(
                [cp, "-interaction=nonstopmode", "-output-directory", td, tp],
                capture_output=True, text=True, timeout=120, cwd=td)
            logs.append(f"── {run_name} (返回码:{proc.returncode}) ──")
            if proc.returncode != 0:
                combined = (proc.stdout or "") + "\n" + (proc.stderr or "")
                errs = [l.strip() for l in combined.split("\n")
                        if l.startswith("!") or "Error" in l or "error" in l]
                logs.append("\n".join(errs[:30]) if errs else
                            (proc.stdout or "")[-2000:]+"\n"+(proc.stderr or "")[-2000:])
                return False
            return True
        except subprocess.TimeoutExpired:
            logs.append("❌ 编译超时(>2分钟)！"); return False
        except Exception as e:
            logs.append(f"❌ 编译异常: {e}"); return False

    tex_content = _safe_read_text(tp)
    has_bib = bool(re.search(r'\\bibliography\s*\{', tex_content))

    if has_bib:
        _run_latex("第1次 latex")
        bibtex_exe = shutil.which("bibtex") or shutil.which("bibtex.exe")
        if bibtex_exe:
            bib_match = re.search(r'\\bibliography\s*\{([^}]+)\}', tex_content)
            if bib_match:
                bib_names = [b.strip() for b in bib_match.group(1).split(',')]
                bib_found = False
                for bib_name in bib_names:
                    bib_file = os.path.join(td, f"{bib_name}.bib")
                    if not os.path.exists(bib_file):
                        bib_file_alt = _resolve_path(f"{bib_name}.bib")
                        if os.path.exists(bib_file_alt): bib_file = bib_file_alt
                    if os.path.exists(bib_file): bib_found = True; break
                if not bib_found:
                    logs.append(f"⚠ 未找到 .bib 文件: {', '.join(bib_names)}.bib")
                else:
                    try:
                        subprocess.run([bibtex_exe, bn], capture_output=True, text=True,
                                       timeout=60, cwd=td)
                        logs.append("── bibtex ──")
                    except subprocess.TimeoutExpired: logs.append("⚠ bibtex 超时")
                    except Exception as e: logs.append(f"⚠ bibtex 运行异常: {e}")
        else: logs.append("⚠ 未找到 bibtex 命令，跳过参考文献处理")
        _run_latex("第2次 latex")
        _run_latex("第3次 latex")
    else:
        _run_latex("第1次 latex")
        _run_latex("第2次 latex")

    if os.path.exists(pp) and os.path.getsize(pp) > 0:
        if clean_aux:
            aux_exts = {".aux",".log",".out",".toc",".lof",".lot",".bbl",".blg",
                        ".synctex.gz",".fdb_latexmk",".fls",".nav",".snm",".vrb"}
            for f in os.listdir(td):
                if os.path.splitext(f)[1] in aux_exts:
                    try: os.remove(os.path.join(td,f))
                    except OSError: pass
        return f"✅ PDF编译成功!\n   📄 {pp}\n   📏 {os.path.getsize(pp)/1024:.1f} KB\n   🔧 {compiler}"
    return f"❌ PDF生成失败。编译器:{compiler}\n\n日志:\n"+"\n".join(logs)

# ═══════════════════════════════════════════════════════════
#  工具 4: Python 沙箱 (v4.0 增强)
# ═══════════════════════════════════════════════════════════

@tool
def PythonSandboxTool(code: str, timeout: int = 60) -> str:
    """执行 Python 代码。支持数据分析(numpy/pandas/scipy/sklearn)、可视化(matplotlib/seaborn)、符号计算(sympy)。⚠ 安全警告：并非真正隔离的沙箱，请勿执行不可信代码。"""
    if not code:
        return "❌ 代码为空。"
    if timeout > 300: timeout = 300
    try:
        so, se, rc = _run_sandbox(code, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"⏰ 超时(>{timeout}s)！请优化代码或增加 timeout 参数。"
    parts = []
    if so: parts.append(f"📤 stdout:\n{so}")
    if se: parts.append(f"📤 stderr:\n{se}")
    if not so and not se: parts.append("✅ 执行完毕，无输出。")
    parts.append(f"\n⏱ 返回码:{rc}")
    return "\n".join(parts)

# ═══════════════════════════════════════════════════════════
#  工具 5: DuckDuckGo 搜索
# ═══════════════════════════════════════════════════════════

@tool
def DuckDuckGoSearchTool(query: str, max_results: int = 5, region: str = "wt-wt") -> str:
    """DuckDuckGo 通用网页搜索。"""
    if not query:
        return "❌ 搜索关键词不能为空。"
    if not _DDGS_AVAILABLE:
        return "❌ DuckDuckGo 搜索不可用。请安装: pip install ddgs"
    max_results = max(1, min(max_results, 15))
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results, region=region))
            if not results: return f"🔍 '{query}' 无结果。"
            out = [f"🔍 '{query}' ({len(results)}条)\n"]
            for i,r in enumerate(results,1):
                b = r.get("body","N/A") or "N/A"
                if len(b)>300: b=b[:300]+"..."
                out.append(f"{i}. {r.get('title','N/A')}\n   🔗 {r.get('href','N/A')}\n   📝 {b}\n")
            return "\n".join(out)
    except Exception as e: return f"❌ 搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 6: 图片搜索下载
# ═══════════════════════════════════════════════════════════

@tool
def FetchWebImageTool(query: str, filename: str) -> str:
    """搜索并下载学术图片。"""
    if not query:
        return "❌ 图片搜索关键词不能为空。"
    if not filename:
        return "❌ 文件名不能为空。"
    if not _DDGS_AVAILABLE:
        return "❌ 图片搜索不可用。请安装: pip install ddgs"
    if not _REQUESTS_AVAILABLE:
        return "❌ requests 未安装"
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=5))
            if not results: return f"❌ '{query}' 无图片。"
            for ir in results:
                try:
                    resp = requests.get(ir["image"], timeout=15); resp.raise_for_status()
                    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
                    with open(sp,"wb") as f: f.write(resp.content)
                    return f"✅ 图片: {sp} ({os.path.getsize(sp)/1024:.1f}KB)\n   🌐 {ir.get('title','N/A')}"
                except: continue
            return f"⚠ 找到{len(results)}张但下载失败。"
    except Exception as e: return f"❌ 图片搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 7: Word 文档编辑 (v4.0 增强)
# ═══════════════════════════════════════════════════════════

@tool
def EditWordDocTool(doc_filename: str, section_title: str = "", text_content: str = "",
                    image_path: str = "", new_page: bool = False,
                    table_data: str = "", heading_level: int = 1) -> str:
    """创建/编辑 Word 文档。支持标题、正文、图片、表格(JSON格式: [["c1","c2"],["d1","d2"]])、分页。"""
    if not doc_filename:
        return "❌ 文件名不能为空。"
    if not _DOCX_AVAILABLE: return "❌ python-docx 未安装。pip install python-docx"
    if not doc_filename.lower().endswith(".docx"): doc_filename += ".docx"
    sp = os.path.join(CUSTOM_TEMP_DIR, doc_filename)
    try:
        doc = Document(sp) if os.path.exists(sp) else Document()
        if new_page and doc.paragraphs: doc.add_page_break()
        if section_title: doc.add_heading(section_title, level=heading_level)
        if text_content: doc.add_paragraph(text_content)
        if image_path:
            ri = _resolve_path(image_path)
            if os.path.exists(ri): doc.add_picture(ri, width=Inches(5.5))
            else: doc.add_paragraph(f"[注：未找到图片 {image_path}]")
        if table_data:
            try:
                rows = json.loads(table_data)
                if rows and isinstance(rows, list):
                    ncols = max(len(r) for r in rows) if rows else 1
                    table = doc.add_table(rows=len(rows), cols=ncols, style='Table Grid')
                    for i, row_data in enumerate(rows):
                        for j, cell_text in enumerate(row_data):
                            if j < ncols:
                                table.rows[i].cells[j].text = str(cell_text)
            except json.JSONDecodeError:
                doc.add_paragraph(f"[表格数据解析失败]")
        doc.save(sp)
        return f"✅ Word: {sp}"
    except Exception as e: return f"❌ Word错误: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 8: Markdown 保存
# ═══════════════════════════════════════════════════════════

@tool
def SaveMarkdownTool(filename: str, content: str, append: bool = False) -> str:
    """保存 Markdown 文件。"""
    if not filename:
        return "❌ 文件名不能为空。"
    if not content:
        return "❌ 内容不能为空。"
    if not filename.lower().endswith(".md"): filename += ".md"
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n\n"+content) if mode=="a" else content)
        return f"✅ Markdown: {sp}"
    except Exception as e: return f"❌ Markdown错误: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 9: 图表生成 (v4.0 增强)
# ═══════════════════════════════════════════════════════════

@tool
def ChartGenerationTool(code: str, filename: str = "chart_output.png", dpi: int = 150, timeout: int = 60) -> str:
    """matplotlib/seaborn 图表生成。支持所有 matplotlib 图表类型、seaborn 统计图表。"""
    if not code:
        return "❌ 图表代码不能为空。"
    if not _MPL_AVAILABLE: return "❌ matplotlib未安装。pip install matplotlib"
    if timeout > 120: timeout = 120
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    prefix = textwrap.dedent("""
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt; import numpy as np
    try: import pandas as pd
    except: pass
    try: import seaborn as sns; sns.set_style("whitegrid")
    except: pass
    from matplotlib import rcParams
    try:
        rcParams['font.sans-serif']=['SimHei','Microsoft YaHei','Noto Sans CJK SC','WenQuanYi Micro Hei']
        rcParams['axes.unicode_minus']=False
    except: pass
    """)
    suffix = textwrap.dedent(f"""
    import os as _os
    figs=[plt.figure(n) for n in plt.get_fignums()]
    if figs:
        for i,fig in enumerate(figs):
            op = r"{sp}" if len(figs)==1 else f"{{_os.path.splitext(r'{sp}')[0]}}_{{i+1}}{{_os.path.splitext(r'{sp}')[1]}}"
            fig.savefig(op,dpi={dpi},bbox_inches='tight'); print(f"✅ 图表: {{op}}")
        plt.close('all')
    else: print("⚠ 未检测到figure。")
    """)
    full_code = prefix + textwrap.indent(code, '    ') + "\n" + suffix
    try:
        so, se, _ = _run_sandbox(full_code, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"⏰ 超时(>{timeout}s)！"
    res = [so] if so else []
    if se: res.append(f"[stderr]: {se}")
    if os.path.exists(sp): res.append(f"📊 {sp} ({os.path.getsize(sp)/1024:.1f}KB)")
    return "\n".join(res) if res else "⚠ 无输出。"

# ═══════════════════════════════════════════════════════════
#  工具 10: Mermaid 流程图
# ═══════════════════════════════════════════════════════════

@tool
def MermaidTool(mermaid_code: str, filename: str = "diagram") -> str:
    """Mermaid 流程图/时序图/饼图/甘特图生成。"""
    # [FIX#16] 防御 None 参数
    safe_code = (mermaid_code or "").strip()
    if not safe_code:
        return "❌ Mermaid 代码不能为空。"
    mp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.mmd")
    try:
        with open(mp,"w",encoding="utf-8") as f: f.write(safe_code)
    except Exception as e: return f"❌ Mermaid写入失败: {e}"
    res = [f"✅ Mermaid源文件: {mp}"]
    if _MMDC_PATH:
        pp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.png")
        try:
            subprocess.run([_MMDC_PATH,"-i",mp,"-o",pp,"-w","1600","-b","white"],
                           capture_output=True,text=True,timeout=30)
            if os.path.exists(pp): res.append(f"🖼 渲染PNG: {pp}")
            else: res.append("⚠ 渲染失败，检查语法。")
        except subprocess.TimeoutExpired: res.append("⚠ 渲染超时。")
        except Exception as e: res.append(f"⚠ 渲染错误: {e}")
    else: res.append("💡 安装mermaid-cli可渲染: npm install -g @mermaid-js/mermaid-cli")
    return "\n".join(res)

# ═══════════════════════════════════════════════════════════
#  工具 11: arXiv 检索
# ═══════════════════════════════════════════════════════════

@tool
def ArXivSearchTool(query: str, max_results: int = 5, sort_by: str = "relevance") -> str:
    """检索 arXiv 学术论文。"""
    # [FIX#16] 防御 None 参数
    if not query:
        return "❌ 搜索关键词不能为空。"
    max_results = max(1, min(max_results, 20))
    base = "https://export.arxiv.org/api/query"
    params = {
        "search_query": f"all:{query}",
        "start": 0, "max_results": max_results,
        "sortBy": sort_by if sort_by in ("relevance","lastUpdatedDate","submittedDate") else "relevance"
    }
    url = f"{base}?{urllib.parse.urlencode(params)}"
    try:
        resp = _safe_request(url, timeout=20)
        if resp is None: return "❌ arXiv API 无响应。"
        root = ET.fromstring(resp.text)
        ns = {"a": "http://www.w3.org/2005/Atom", "opensearch": "http://a9.com/-/spec/opensearch/1.1/"}
        total = root.find("opensearch:totalResults", ns)
        total_str = total.text if total is not None else "?"
        entries = root.findall("a:entry", ns)
        if not entries: return f"🔍 arXiv: '{query}' 无结果。"
        out = [f"📚 arXiv: '{query}' (共{total_str}篇，显示{len(entries)}篇)\n"]
        for i, e in enumerate(entries, 1):
            title_el = e.find("a:title", ns)
            title = " ".join(((title_el.text or "N/A") if title_el is not None else "N/A").split())
            authors = [a.find("a:name", ns).text for a in e.findall("a:author", ns) if a.find("a:name", ns) is not None]
            authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
            summary_el = e.find("a:summary", ns)
            abstract = " ".join((((summary_el.text or "N/A") if summary_el is not None else "N/A")[:600]).split())
            arxiv_id_el = e.find("a:id", ns)
            arxiv_id = ((arxiv_id_el.text or "") if arxiv_id_el is not None else "").split("/abs/")[-1]
            published_el = e.find("a:published", ns)
            published = ((published_el.text or "N/A") if published_el is not None else "N/A")[:10]
            cats = [c.get("term","") for c in e.findall("a:category", ns)]
            out.append(
                f"{i}. {title}\n   👤 {authors_str}\n"
                f"   📅 {published} | 📂 {', '.join(cats[:3])}\n"
                f"   🔗 https://arxiv.org/abs/{arxiv_id}\n"
                f"   📄 PDF: https://arxiv.org/pdf/{arxiv_id}.pdf\n"
                f"   📝 {abstract}\n"
            )
        return "\n".join(out)
    except ET.ParseError as e: return f"❌ arXiv XML解析错误: {e}"
    except Exception as e: return f"❌ arXiv搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 12: DOI 元数据
# ═══════════════════════════════════════════════════════════

@tool
def DOIMetadataTool(doi: str) -> str:
    """通过 DOI 获取论文完整元数据（Crossref API）。"""
    # [FIX#16] 防御 None 参数 — 这是最可能的错误源
    if not doi:
        return "❌ DOI 不能为空。请提供有效的 DOI（如 10.1000/xyz123）。"
    doi = doi.strip().replace("https://doi.org/","").replace("http://dx.doi.org/","")
    if not doi:
        return "❌ DOI 格式无效。"
    url = f"https://api.crossref.org/works/{urllib.parse.quote(doi, safe='')}"
    try:
        resp = _safe_request(url, timeout=15)
        if resp is None: return "❌ Crossref API 无响应。"
        data = resp.json()
        msg = data.get("message", {})
        if not msg: return f"❌ DOI '{doi}' 无数据。"

        title = " ".join((msg.get("title",["N/A"])[0] or "N/A").split())
        authors = []
        for a in msg.get("author", [])[:15]:
            fam = a.get("family",""); giv = a.get("given","")
            authors.append(f"{giv} {fam}".strip() or fam or giv or "Unknown")
        authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
        journal = " | ".join(msg.get("container-title",["N/A"]))
        # Safe year extraction: avoid IndexError on empty date-parts
        year = "?"
        for date_key in ("published-print", "published-online", "created"):
            dp = msg.get(date_key, {}).get("date-parts", [])
            if dp and isinstance(dp, list):
                for part in dp:
                    if part and isinstance(part, list) and len(part) > 0 and part[0] is not None:
                        year = part[0]
                        break
            if year != "?":
                break
        vol = msg.get("volume","?"); issue = msg.get("issue","?")
        page = msg.get("page","?")
        abstract = " ".join(((msg.get("abstract","N/A") or "N/A")[:800]).split())
        abstract = re.sub(r'<[^>]+>', '', abstract)
        cited = msg.get("is-referenced-by-count", "?")
        publisher = msg.get("publisher","?")

        first_author_last = (authors[0].split()[-1] if authors else "unknown")
        bib = f"@article{{{first_author_last}{year},\n"
        bib += f"  title = {{{title}}},\n"
        bib += f"  author = {{{' and '.join(authors[:10])}}},\n"
        bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{year}}},\n"
        if vol!="?": bib += f"  volume = {{{vol}}},\n"
        if issue!="?": bib += f"  number = {{{issue}}},\n"
        if page!="?": bib += f"  pages = {{{page}}},\n"
        bib += f"  doi = {{{doi}}},\n"
        bib += f"  publisher = {{{publisher}}}\n}}"

        return (
            f"📄 DOI元数据: {doi}\n{'='*50}\n"
            f"📌 {title}\n👤 {authors_str}\n"
            f"📰 {journal} | {year} | Vol.{vol} | Iss.{issue} | pp.{page}\n"
            f"🏢 {publisher}\n📊 被引: {cited}次\n"
            f"🔗 https://doi.org/{doi}\n📝 {abstract}\n\n── BibTeX ──\n{bib}"
        )
    except Exception as e: return f"❌ DOI查询失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 13: 数据统计
# ═══════════════════════════════════════════════════════════

@tool
def DataStatisticsTool(data_path: str, max_rows: int = 10000) -> str:
    """对 CSV/Excel 数据自动生成描述性统计报告。"""
    if not data_path:
        return "❌ 文件路径不能为空。"
    if not _PANDAS_AVAILABLE: return "❌ pandas 未安装"
    fp = _resolve_path(data_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {data_path}"
    ext = os.path.splitext(fp)[1].lower()
    try:
        if ext == ".csv":
            for enc in ["utf-8","gbk","latin-1"]:
                try: df = pd.read_csv(fp, encoding=enc, nrows=max_rows); break
                except: continue
            else: return "❌ CSV编码读取失败。"
        elif ext in {".xlsx",".xls"}: df = pd.read_excel(fp, nrows=max_rows)
        else: return f"❌ 不支持格式: {ext}"
    except Exception as e: return f"❌ 数据加载失败: {e}"

    rows, cols = df.shape
    out = [f"📊 数据统计: {os.path.basename(fp)}\n{'='*60}", f"📏 {rows}行 × {cols}列\n"]
    num_cols = df.select_dtypes(include=["number"]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object","category"]).columns.tolist()
    date_cols = df.select_dtypes(include=["datetime"]).columns.tolist()
    other_cols = [c for c in df.columns if c not in num_cols+cat_cols+date_cols]
    out.append(f"🔢 数值列({len(num_cols)}): {', '.join(num_cols[:15])}{'...' if len(num_cols)>15 else ''}")
    out.append(f"🔤 分类列({len(cat_cols)}): {', '.join(cat_cols[:15])}{'...' if len(cat_cols)>15 else ''}")
    if date_cols: out.append(f"📅 日期列({len(date_cols)}): {', '.join(date_cols)}")

    missing = df.isnull().sum(); missing = missing[missing > 0]
    if len(missing) > 0:
        out.append(f"\n⚠ 缺失值 ({missing.sum()} total):")
        for c, m in missing.items(): out.append(f"   {c}: {m} ({m/rows*100:.1f}%)")
    else: out.append("\n✅ 无缺失值")

    if num_cols:
        out.append(f"\n── 数值列统计 ──")
        try:
            stats = df[num_cols].describe(percentiles=[.25,.5,.75])
            if _NUMPY_AVAILABLE and len(num_cols) <= 20:
                stats.loc["skew"] = df[num_cols].skew()
                stats.loc["kurtosis"] = df[num_cols].kurtosis()
            out.append(stats.to_string(max_rows=30))
        except Exception as e: out.append(f"   (统计计算失败: {e})")

    if cat_cols:
        out.append(f"\n── 分类列 Top5 ──")
        for c in cat_cols[:10]:
            vc = df[c].value_counts().head(5)
            out.append(f"  {c}:")
            for v, cnt in vc.items(): out.append(f"    {v}: {cnt} ({cnt/rows*100:.1f}%)")

    if len(num_cols) >= 2 and len(num_cols) <= 15:
        out.append(f"\n── 相关性矩阵 (|r|>0.3) ──")
        try:
            corr = df[num_cols].corr(); found = False
            for i in range(len(num_cols)):
                for j in range(i+1, len(num_cols)):
                    r = corr.iloc[i,j]
                    if abs(r) > 0.3: out.append(f"  {num_cols[i]} ↔ {num_cols[j]}: r={r:.3f}"); found = True
            if not found: out.append("  (无 |r|>0.3 的显著相关性)")
        except Exception: out.append("  (相关性计算失败)")
    return "\n".join(out)

# ═══════════════════════════════════════════════════════════
#  工具 14: BibTeX 工具
# ═══════════════════════════════════════════════════════════

@tool
def BibTexTool(
    mode: str, bib_path: str = "", bib_content: str = "",
    title: str = "", authors: str = "", journal: str = "",
    year: str = "", doi: str = "", volume: str = "", pages: str = "",
    citation_style: str = "apa"
) -> str:
    """BibTeX 解析、生成与引用格式化。mode: parse/generate/format。"""
    if not mode:
        return "❌ mode 不能为空。可选: 'parse', 'generate', 'format'"

    if mode == "parse":
        raw = None  # [FIX#21] 显式初始化
        if bib_content:
            raw = bib_content
        elif bib_path:
            rp = _resolve_path(bib_path)
            if not os.path.exists(rp): return f"❌ 文件不存在: {bib_path}"
            raw = _safe_read_text(rp)
        else:
            return "❌ parse模式需要 bib_path 或 bib_content。"

        if not raw:
            return "❌ BibTeX 内容为空。"

        entries = []

        if _BIBTEXPARSER_AVAILABLE:
            try:
                parser = BibTexParser(common_strings=True)
                parser.customization = homogenize_latex_encoding
                bib_db = bibtexparser.loads(raw, parser=parser)
                for e in bib_db.entries:
                    entries.append({
                        "type": e.get("ENTRYTYPE", "misc"),
                        "key": e.get("ID", "?"),
                        "fields": {k.lower(): v for k, v in e.items()
                                    if k not in ("ENTRYTYPE", "ID")}
                    })
                if entries:
                    out = [f"📚 BibTeX解析 (bibtexparser): {len(entries)}条\n"]
                    for i, e in enumerate(entries, 1):
                        f = e["fields"]
                        out.append(f"{i}. [{e['type']}] {f.get('title','?')[:80]}")
                        out.append(f"   👤 {f.get('author','?')[:100]}")
                        out.append(f"   📰 {f.get('journal',f.get('booktitle','?'))} | {f.get('year','?')} | 🔑 {e['key']}")
                    return "\n".join(out)
            except Exception as e:
                _log("WARN", f"bibtexparser 解析失败，回退内置: {e}")

        entries = _robust_parse_bibtex(raw)
        if not entries: return "⚠ 未解析到 BibTeX 条目。"
        out = [f"📚 BibTeX解析 (内置): {len(entries)}条\n"]
        for i, e in enumerate(entries, 1):
            f = e["fields"]
            out.append(f"{i}. [{e['type']}] {f.get('title','?')[:80]}")
            out.append(f"   👤 {f.get('author','?')[:100]}")
            out.append(f"   📰 {f.get('journal',f.get('booktitle','?'))} | {f.get('year','?')} | 🔑 {e['key']}")
        return "\n".join(out)

    if mode == "generate":
        if not title: return "❌ generate模式需要 title 参数。"
        first_author = authors.split(",")[0].strip().split() if authors else ["unknown"]
        key_surname = first_author[-1] if first_author else "unknown"
        yr = year or "????"
        bib = f"@article{{{key_surname}{yr},\n  title = {{{title}}},\n"
        if authors: bib += f"  author = {{{authors}}},\n"
        if journal: bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{yr}}},\n"
        if volume: bib += f"  volume = {{{volume}}},\n"
        if pages: bib += f"  pages = {{{pages}}},\n"
        if doi: bib += f"  doi = {{{doi}}},\n"
        bib += "}"
        return f"✅ BibTeX 已生成:\n```bibtex\n{bib}\n```"

    if mode == "format":
        if not title: return "❌ format模式需要 title 参数。"
        au_list = [a.strip() for a in authors.split(",") if a.strip()] if authors else ["Unknown"]
        yr = year or "(n.d.)"
        ref = ""

        if citation_style == "apa":
            if len(au_list) == 1: au_str = au_list[0]
            elif len(au_list) == 2: au_str = f"{au_list[0]} & {au_list[1]}"
            else: au_str = f"{au_list[0]} et al."
            ref = f"{au_str} ({yr}). {title}."
            if journal: ref += f" *{journal}*"
            if volume: ref += f", *{volume}*"
            if pages: ref += f", {pages}"
            ref += "."
            if doi: ref += f" https://doi.org/{doi}"

        elif citation_style == "mla":
            au_str = f"{au_list[0]} et al." if len(au_list) > 2 else " and ".join(au_list)
            ref = f'{au_str}. "{title}."'
            if journal: ref += f" *{journal}*"
            if volume: ref += f" {volume}"
            ref += f" ({yr})"
            ref += f": {pages}." if pages else "."
            if doi: ref += f" doi:{doi}."

        elif citation_style == "chicago":
            au_str = ", ".join(au_list[:3]) + (", et al." if len(au_list) > 3 else "")
            ref = f'{au_str}. "{title}."'
            if journal: ref += f" *{journal}*"
            if volume: ref += f" {volume}"
            ref += f" ({yr})"
            ref += f": {pages}." if pages else "."
            if doi: ref += f" https://doi.org/{doi}."

        elif citation_style == "ieee":
            ieee_authors = []
            for au in au_list:
                parts = au.strip().split()
                if len(parts) >= 2:
                    initials = " ".join([p[0] + "." for p in parts[:-1]])
                    ieee_authors.append(f"{initials} {parts[-1]}")
                elif len(parts) == 1: ieee_authors.append(parts[0])
                else: ieee_authors.append(au)

            if len(ieee_authors) == 1: au_str = ieee_authors[0]
            elif len(ieee_authors) == 2: au_str = f"{ieee_authors[0]} and {ieee_authors[1]}"
            elif len(ieee_authors) <= 6:
                au_str = ", ".join(ieee_authors[:-1]) + f", and {ieee_authors[-1]}"
            else: au_str = ", ".join(ieee_authors[:6]) + ", et al."

            ref = f'{au_str}, "{title},"'
            if journal: ref += f" *{journal}*"
            if volume: ref += f", vol. {volume}"
            if pages: ref += f", pp. {pages}"
            ref += f", {yr}."
            if doi: ref += f" doi: {doi}."

        else:
            return f"❌ 不支持的引用格式: {citation_style}。支持: apa, mla, chicago, ieee"

        return f"📝 引用 ({citation_style.upper()}):\n{ref}"

    return "❌ 无效 mode。可选: 'parse', 'generate', 'format'"

# ═══════════════════════════════════════════════════════════
#  工具 15: PDF 表格提取
# ═══════════════════════════════════════════════════════════

@tool
def PDFTableExtractTool(pdf_path: str, page_range: str = "all") -> str:
    """从 PDF 中提取表格数据。"""
    if not pdf_path:
        return "❌ PDF 文件路径不能为空。"
    if not _FITZ_AVAILABLE: return "❌ PyMuPDF 未安装"
    fp = _resolve_path(pdf_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {pdf_path}"
    try:
        doc = fitz.open(fp); total = len(doc)
        if page_range == "all": pages = list(range(total))
        else:
            pages = []
            for part in page_range.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    pages.extend(range(max(0, int(a)-1), min(total, int(b))))
                else: pages.append(max(0, int(part)-1))
            pages = sorted(set(pages))
        out = [f"📊 PDF表格提取: {os.path.basename(fp)} (共{total}页，提取{len(pages)}页)\n"]
        table_count = 0
        for pi in pages:
            page = doc[pi]
            try:
                tabs = page.find_tables()
                if tabs and tabs.tables:
                    for ti, tab in enumerate(tabs.tables):
                        table_count += 1
                        out.append(f"── 表格{table_count} (Page {pi+1}) ──")
                        rows = []
                        for row in tab.extract():
                            rows.append(" | ".join(str(cell).replace("\n"," ") if cell else "" for cell in row))
                        out.append("\n".join(rows)); out.append("")
            except Exception:
                pass
        if table_count == 0: return "\n".join(out) + "\n⚠ 未提取到表格。"
        return "\n".join(out)
    except Exception as e: return f"❌ PDF表格提取失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 16: 会话导出
# ═══════════════════════════════════════════════════════════

@tool
def SessionExportTool(content: str, filename: str = "session_export", export_format: str = "markdown") -> str:
    """将对话内容或分析结果导出为文件。支持 markdown, text, html, pdf。"""
    if not content:
        return "❌ 导出内容不能为空。"
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    header = f"---\ntitle: Session Export\ndate: {ts}\n---\n\n"
    full_content = header + content
    if export_format == "markdown":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.md")
        with open(sp, "w", encoding="utf-8") as f: f.write(full_content)
        return f"✅ Markdown: {sp}"
    elif export_format == "text":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.txt")
        with open(sp, "w", encoding="utf-8") as f: f.write(full_content)
        return f"✅ 文本: {sp}"
    elif export_format == "html":
        html_body = content.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
        html_body = html_body.replace("\n","<br>\n")
        html = f"<!DOCTYPE html><html lang='zh'><head><meta charset='UTF-8'><title>Session Export</title>" \
               f"<style>body{{font-family:Georgia,serif;max-width:800px;margin:2em auto;line-height:1.7}}</style>" \
               f"</head><body><h1>Session Export</h1><p><em>{ts}</em></p>{html_body}</body></html>"
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.html")
        with open(sp, "w", encoding="utf-8") as f: f.write(html)
        return f"✅ HTML: {sp}"
    elif export_format == "pdf":
        md_sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.md")
        with open(md_sp, "w", encoding="utf-8") as f: f.write(full_content)
        if _PANDOC_PATH:
            cmd = [_PANDOC_PATH, md_sp, "-o", os.path.join(CUSTOM_TEMP_DIR, f"{filename}.pdf")]
            if _DETECTED_CJK_FONT:
                cmd.extend(["-V", f"CJKmainfont={_DETECTED_CJK_FONT}"])
            subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            return f"✅ PDF: {os.path.join(CUSTOM_TEMP_DIR, f'{filename}.pdf')}"
        return f"✅ Markdown (Pandoc未安装，无法转PDF): {md_sp}"
    else: return f"❌ 不支持的格式: {export_format}。支持: markdown, text, html, pdf"

# ═══════════════════════════════════════════════════════════
#  工具 17: Pandoc 格式转换
# ═══════════════════════════════════════════════════════════

@tool
def PandocConvertTool(source_path: str, output_format: str = "pdf",
                      output_filename: str = "", extra_args: str = "") -> str:
    """使用 Pandoc 进行文档格式转换。需安装 pandoc。"""
    if not source_path:
        return "❌ 源文件路径不能为空。"
    if not _PANDOC_PATH:
        return ("❌ Pandoc 未安装！\n"
                "   Windows: winget install Pandoc.Pandoc\n"
                "   macOS:   brew install pandoc\n"
                "   Linux:   sudo apt install pandoc")
    sfp = _resolve_path(source_path)
    if not os.path.exists(sfp): return f"❌ 源文件不存在: {source_path}"
    if not output_filename:
        base = os.path.splitext(os.path.basename(sfp))[0]
        ext_map = {"pdf": ".pdf", "docx": ".docx", "html": ".html",
                   "latex": ".tex", "epub": ".epub", "markdown": ".md"}
        output_filename = f"pandoc_{base}{ext_map.get(output_format, f'.{output_format}')}"
    ofp = os.path.join(CUSTOM_TEMP_DIR, output_filename)
    cmd = [_PANDOC_PATH, sfp, "-o", ofp, f"--to={output_format}"]

    if output_format == "pdf" and "--pdf-engine" not in extra_args:
        xelatex = shutil.which("xelatex")
        if xelatex:
            cmd.extend(["--pdf-engine=xelatex"])
            if "-V" not in extra_args and "mainfont" not in extra_args.lower():
                if _DETECTED_CJK_FONT:
                    cmd.extend(["-V", f"CJKmainfont={_DETECTED_CJK_FONT}",
                                "-V", f"mainfont={_DETECTED_CJK_FONT}"])

    if extra_args: cmd.extend(extra_args.split())
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if proc.returncode != 0:
            stderr = proc.stderr or ""
            return f"❌ Pandoc 转换失败 (返回码 {proc.returncode}):\n{stderr[-1000:]}"
        if os.path.exists(ofp) and os.path.getsize(ofp) > 0:
            return f"✅ 转换成功: {ofp} ({os.path.getsize(ofp)/1024:.1f} KB)\n   格式: {output_format}"
        return f"❌ 输出文件未生成。stderr: {(proc.stderr or '')[:500]}"
    except subprocess.TimeoutExpired: return "❌ Pandoc 转换超时（>2分钟）！"
    except Exception as e: return f"❌ Pandoc 异常: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 18: PDF 批注/高亮提取
# ═══════════════════════════════════════════════════════════

@tool
def PDFAnnotExtractTool(pdf_path: str, annot_types: str = "all") -> str:
    """提取 PDF 中的批注、高亮、下划线、便签。"""
    if not pdf_path:
        return "❌ PDF 文件路径不能为空。"
    if not _FITZ_AVAILABLE: return "❌ PyMuPDF 未安装"
    fp = _resolve_path(pdf_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {pdf_path}"
    try:
        doc = fitz.open(fp); total_pages = len(doc)
        all_annots = []
        type_filter = None if annot_types == "all" else annot_types
        for pi in range(total_pages):
            page = doc[pi]
            annots = page.annots()
            if not annots: continue
            for annot in annots:
                atype = annot.type[1] if isinstance(annot.type, tuple) else str(annot.type)
                info = annot.info
                content = info.get("content", "")
                title = info.get("title", "")
                if hasattr(annot, "get_textbox"):
                    try: text = page.get_textbox(annot.rect).strip()
                    except Exception: text = ""
                else: text = ""
                if type_filter and type_filter.lower() not in atype.lower(): continue
                color = ""
                if "stroke" in info:
                    color = f"#{info['stroke']:06x}" if isinstance(info["stroke"], int) else str(info["stroke"])
                all_annots.append({
                    "page": pi+1, "type": atype,
                    "content": content or text or "(空)",
                    "author": title or "未知", "color": color,
                })
        if not all_annots:
            return f"📄 PDF批注提取: {os.path.basename(fp)}\n   ✅ 未检测到批注/高亮。"
        type_counts: Dict[str, int] = {}
        for a in all_annots: type_counts[a["type"]] = type_counts.get(a["type"], 0) + 1
        out = [f"📄 PDF批注提取: {os.path.basename(fp)} ({total_pages}页)",
               f"   共 {len(all_annots)} 条批注 | " +
               ", ".join(f"{t}:{c}" for t,c in sorted(type_counts.items())),
               "="*60]
        for a in all_annots:
            icon = {"Highlight":"🟡","Underline":"🔵","StrikeOut":"🔴",
                    "Text":"📝","FreeText":"💬","Stamp":"🏷"}.get(a["type"],"📌")
            out.append(f"\n{icon} [p{a['page']}] [{a['type']}] {a['author']}")
            if a["color"]: out.append(f"   🎨 {a['color']}")
            out.append(f"   {a['content'][:500]}")
        return "\n".join(out)
    except Exception as e: return f"❌ PDF批注提取失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 19: 学术翻译
# ═══════════════════════════════════════════════════════════

@tool
def AcademicTranslateTool(text: str, direction: str = "zh2en",
                          preserve_terms: str = "") -> str:
    """学术中英互译，保留术语一致性。"""
    # [FIX#16] 防御 None 参数
    if not text or not text.strip():
        return "❌ 文本为空。"
    term_instruction = ""
    if preserve_terms:
        terms = [t.strip() for t in preserve_terms.split(",") if t.strip()]
        term_instruction = f"\n请保留以下术语不翻译: {', '.join(terms)}。"
    if direction == "zh2en":
        prompt = f"请将以下中文学术文本翻译为地道的英文学术英语。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    elif direction == "en2zh":
        prompt = f"请将以下英文学术文本翻译为流畅的中文学术语言。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    else:
        prompt = f"请将以下学术文本翻译（自动检测源语言）。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    try:
        translator = ChatOpenAI(
            model=os.getenv("DFTB_MODEL", "deepseek-v4-pro"),
            base_url="https://api.deepseek.com",
            api_key=api_key,
            temperature=0.1,
            max_tokens=4096,
        )
        result = translator.invoke([HumanMessage(content=prompt)])
        translated = result.content
        if preserve_terms:
            translated += f"\n\n── 保留术语 ──\n{preserve_terms}"
        return f"🌐 学术翻译 ({direction}):\n{'='*40}\n{translated}"
    except Exception as e:
        return f"❌ 翻译失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 20: Semantic Scholar 检索
# ═══════════════════════════════════════════════════════════

@tool
def SemanticScholarTool(query: str, max_results: int = 5,
                        fields: str = "title,authors,year,abstract,externalIds,citationCount,url") -> str:
    """检索 Semantic Scholar 学术论文（比 arXiv 更全，含引用网络）。"""
    # [FIX#20] 防御 None 参数
    if not query:
        return "❌ 搜索关键词不能为空。"
    max_results = max(1, min(max_results, 20))
    base = "https://api.semanticscholar.org/graph/v1/paper/search"
    params = {"query": query, "limit": max_results, "fields": fields}
    url = f"{base}?{urllib.parse.urlencode(params)}"
    try:
        resp = _safe_request(url, timeout=20)
        if resp is None: return "❌ Semantic Scholar API 无响应。"
        data = resp.json()
        papers = data.get("data", [])
        if not papers: return f"🔍 Semantic Scholar: '{query}' 无结果。"
        out = [f"📚 Semantic Scholar: '{query}' (共{data.get('total','?')}篇，显示{len(papers)}篇)\n"]
        for i, p in enumerate(papers, 1):
            title = p.get("title", "N/A") or "N/A"
            authors = ", ".join([a.get("name","?") for a in p.get("authors",[])][:8])
            if len(p.get("authors",[])) > 8: authors += ", et al."
            year = p.get("year", "?")
            abstract = (p.get("abstract", "N/A") or "N/A")[:500]
            cited = p.get("citationCount", "?")
            ext_ids = p.get("externalIds", {})
            arxiv_id = ext_ids.get("ArXiv", "")
            doi = ext_ids.get("DOI", "")
            links = []
            if arxiv_id: links.append(f"https://arxiv.org/abs/{arxiv_id}")
            if doi: links.append(f"https://doi.org/{doi}")
            link_str = " | ".join(links) if links else p.get("url", "N/A")
            out.append(
                f"{i}. {title}\n   👤 {authors}\n"
                f"   📅 {year} | 📊 被引: {cited}\n"
                f"   🔗 {link_str}\n   📝 {abstract}\n"
            )
        return "\n".join(out)
    except Exception as e: return f"❌ Semantic Scholar搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 21: Markdown → PPTX 演示生成 (v4.0 大幅增强)
# ═══════════════════════════════════════════════════════════

@tool
def PresentationGenTool(markdown_content: str = "", filename: str = "presentation",
                        theme: str = "default", source_file: str = "",
                        include_images: str = "", speaker_notes: bool = False) -> str:
    """从 Markdown 大纲自动生成 PPTX 演示文稿。
支持多主题(default/dark/academic/modern/corporate/minimal)、图片嵌入、表格、演讲者备注。
格式: # 标题 / ## 副标题 / - 要点 / --- 分页 / ![图片](路径) / |表格|。
source_file: 可直接指定 .md 文件路径代替 markdown_content。"""
    if not _PPTX_AVAILABLE: return "❌ python-pptx 未安装。pip install python-pptx"

    if source_file and not markdown_content:
        sfp = _resolve_path(source_file)
        if os.path.exists(sfp):
            markdown_content = _safe_read_text(sfp)
        else:
            return f"❌ 源文件不存在: {source_file}"

    # [FIX#16] 防御 None — markdown_content 可能仍是 None
    if not markdown_content or not markdown_content.strip():
        return "❌ 内容为空。请提供 markdown_content 或 source_file。"

    sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.pptx")
    prs = Presentation()

    themes = {
        "default": {
            "bg": PptRGBColor(0xFF, 0xFF, 0xFF), "text": PptRGBColor(0x33, 0x33, 0x33),
            "accent": PptRGBColor(0x00, 0x7A, 0xCC), "title_size": 32, "subtitle_size": 20, "body_size": 16,
        },
        "dark": {
            "bg": PptRGBColor(0x1A, 0x1A, 0x2E), "text": PptRGBColor(0xE0, 0xE0, 0xE0),
            "accent": PptRGBColor(0x64, 0xFF, 0xDA), "title_size": 32, "subtitle_size": 20, "body_size": 16,
        },
        "academic": {
            "bg": PptRGBColor(0xFA, 0xFA, 0xFA), "text": PptRGBColor(0x22, 0x22, 0x22),
            "accent": PptRGBColor(0x00, 0x55, 0xAA), "title_size": 30, "subtitle_size": 18, "body_size": 15,
        },
        "modern": {
            "bg": PptRGBColor(0xFF, 0xFF, 0xFF), "text": PptRGBColor(0x2D, 0x2D, 0x2D),
            "accent": PptRGBColor(0xE6, 0x3E, 0x31), "title_size": 34, "subtitle_size": 20, "body_size": 16,
        },
        "corporate": {
            "bg": PptRGBColor(0xFF, 0xFF, 0xFF), "text": PptRGBColor(0x1B, 0x2A, 0x4A),
            "accent": PptRGBColor(0x1B, 0x2A, 0x4A), "title_size": 30, "subtitle_size": 18, "body_size": 15,
        },
        "minimal": {
            "bg": PptRGBColor(0xFF, 0xFF, 0xFF), "text": PptRGBColor(0x33, 0x33, 0x33),
            "accent": PptRGBColor(0x66, 0x66, 0x66), "title_size": 28, "subtitle_size": 18, "body_size": 15,
        },
    }
    thm = themes.get(theme, themes["default"])

    def _set_slide_bg(slide, color: PptRGBColor):
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color
        except Exception: pass

    def _get_slide_layout(prs_obj, preferred_idx=1):
        try:
            if preferred_idx < len(prs_obj.slide_layouts):
                return prs_obj.slide_layouts[preferred_idx]
            return prs_obj.slide_layouts[0]
        except Exception:
            return prs_obj.slide_layouts[0]

    def _add_text_run(para, text, size, color, bold=False):
        para.text = ""
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        return run

    slide_layout = _get_slide_layout(prs, 1)
    slides_text = markdown_content.split("\n---\n")
    image_paths = []
    if include_images:
        image_paths = [p.strip() for p in include_images.split(",") if p.strip()]

    slide_count = 0

    for slide_idx, slide_text in enumerate(slides_text):
        slide_text = slide_text.strip()
        if not slide_text: continue

        slide = prs.slides.add_slide(slide_layout)
        _set_slide_bg(slide, thm["bg"])
        slide_count += 1

        lines = slide_text.split("\n")
        title = ""; subtitle = ""; bullets = []
        table_rows = []
        image_path = None

        for line in lines:
            line_stripped = line.strip()

            img_match = re.match(r'!\[.*?\]\((.*?)\)', line_stripped)
            if img_match:
                image_path = img_match.group(1)
                continue

            if line_stripped.startswith("|") and line_stripped.endswith("|"):
                cells = [c.strip() for c in line_stripped[1:-1].split("|")]
                if not all(re.match(r'^[-: ]+$', c) for c in cells):
                    table_rows.append(cells)
                continue

            if line_stripped.startswith("# ") and not title:
                title = line_stripped[2:]
            elif line_stripped.startswith("## ") and not subtitle:
                subtitle = line_stripped[3:]
            elif line_stripped.startswith("- ") or line_stripped.startswith("* ") or line_stripped.startswith("• "):
                bullets.append(line_stripped[2:].strip())
            elif line_stripped.startswith("> "):
                bullets.append(f"💬 {line_stripped[2:]}")

        if title and slide.shapes.title:
            slide.shapes.title.text = ""
            para = slide.shapes.title.text_frame.paragraphs[0]
            _add_text_run(para, title, thm["title_size"], thm["accent"], bold=True)

        body_shape = None
        try:
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
        except: pass

        if body_shape is None:
            left = PptInches(1); top = PptInches(2.2)
            width = PptInches(8); height = PptInches(4.5)
            body_shape = slide.shapes.add_textbox(left, top, width, height)

        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True

        if subtitle:
            p = tf.paragraphs[0]
            _add_text_run(p, subtitle, thm["subtitle_size"], thm["text"])
            p.space_after = Pt(12)

        for b in bullets:
            p = tf.add_paragraph()
            _add_text_run(p, f"• {b}", thm["body_size"], thm["text"])
            p.space_after = Pt(6)
            p.level = 0

        if table_rows:
            p = tf.add_paragraph()
            _add_text_run(p, "", 8, thm["text"])
            ncols = max(len(r) for r in table_rows)
            nrows = len(table_rows)
            tbl_left = PptInches(1.5); tbl_top = PptInches(4.5)
            tbl_width = PptInches(7); tbl_height = PptInches(0.4 * nrows)
            try:
                table_shape = slide.shapes.add_table(nrows, ncols, tbl_left, tbl_top, tbl_width, tbl_height)
                table = table_shape.table
                for ri, row_data in enumerate(table_rows):
                    for ci, cell_text in enumerate(row_data):
                        if ci < ncols:
                            cell = table.rows[ri].cells[ci]
                            cell.text = cell_text
                            for para in cell.text_frame.paragraphs:
                                para.font.size = Pt(12)
                                para.font.color.rgb = thm["text"]
            except Exception: pass

        if image_path:
            rip = _resolve_path(image_path)
            if os.path.exists(rip):
                try:
                    slide.shapes.add_picture(rip, PptInches(5.5), PptInches(4.5), width=PptInches(3.5))
                except Exception: pass
        elif slide_idx < len(image_paths) and image_paths[slide_idx]:
            rip = _resolve_path(image_paths[slide_idx])
            if os.path.exists(rip):
                try:
                    slide.shapes.add_picture(rip, PptInches(5.5), PptInches(4.5), width=PptInches(3.5))
                except Exception: pass

        if speaker_notes and bullets:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame
            notes_text.text = "\n".join(bullets)

    if slide_count > 0:
        end_slide = prs.slides.add_slide(slide_layout)
        _set_slide_bg(end_slide, thm["bg"])
        if end_slide.shapes.title:
            end_slide.shapes.title.text = ""
            para = end_slide.shapes.title.text_frame.paragraphs[0]
            _add_text_run(para, "谢谢！/ Thank You!", thm["title_size"], thm["accent"], bold=True)
        if len(end_slide.placeholders) > 1:
            tf_end = end_slide.placeholders[1].text_frame
            tf_end.clear()
            p = tf_end.paragraphs[0]
            _add_text_run(p, "Questions & Discussion", thm["body_size"], thm["text"])

    prs.save(sp)
    return (f"✅ PPTX 已生成: {sp}\n"
            f"   📊 共 {slide_count + (1 if slide_count > 0 else 0)} 张幻灯片\n"
            f"   🎨 主题: {theme}\n"
            f"   💡 使用 PowerPoint/Keynote 打开编辑")

# ═══════════════════════════════════════════════════════════
#  🆕 工具 22: 学习计划生成器
# ═══════════════════════════════════════════════════════════

@tool
def StudyPlanTool(subject: str, duration_weeks: int = 8, hours_per_week: int = 10,
                  difficulty: str = "intermediate", goal: str = "",
                  output_format: str = "markdown") -> str:
    """生成个性化学习计划。参数: subject(学科), duration_weeks(周数), hours_per_week(每周小时), difficulty(beginner/intermediate/advanced), goal(学习目标)。"""
    # [FIX#16] 防御 None
    if not subject:
        return "❌ 学科名称不能为空。"
    if duration_weeks < 1: duration_weeks = 1
    if duration_weeks > 52: duration_weeks = 52
    if hours_per_week < 1: hours_per_week = 1
    if hours_per_week > 60: hours_per_week = 60

    total_hours = duration_weeks * hours_per_week
    difficulties = {
        "beginner": {"theory": 0.5, "practice": 0.35, "review": 0.15},
        "intermediate": {"theory": 0.35, "practice": 0.45, "review": 0.20},
        "advanced": {"theory": 0.25, "practice": 0.50, "review": 0.25},
    }
    diff = difficulties.get(difficulty, difficulties["intermediate"])

    plan = []
    plan.append(f"# 📚 学习计划: {subject}")
    plan.append(f"\n**难度**: {difficulty} | **周期**: {duration_weeks} 周 | **每周**: {hours_per_week} 小时 | **总计**: {total_hours} 小时")
    if goal: plan.append(f"\n**目标**: {goal}")
    plan.append(f"\n---\n")

    phases = [
        {"name": "基础阶段", "weeks_ratio": 0.3, "desc": "建立核心概念与基础知识体系"},
        {"name": "进阶阶段", "weeks_ratio": 0.4, "desc": "深入理解复杂概念，大量练习"},
        {"name": "综合阶段", "weeks_ratio": 0.2, "desc": "综合应用、项目实践"},
        {"name": "复习冲刺", "weeks_ratio": 0.1, "desc": "系统复习、查漏补缺"},
    ]

    current_week = 1
    for phase in phases:
        phase_weeks = max(1, round(duration_weeks * phase["weeks_ratio"]))
        plan.append(f"## {phase['name']} (第 {current_week}-{min(current_week+phase_weeks-1, duration_weeks)} 周)")
        plan.append(f"*{phase['desc']}*\n")

        for w in range(phase_weeks):
            if current_week > duration_weeks: break
            theory_h = round(hours_per_week * diff["theory"], 1)
            practice_h = round(hours_per_week * diff["practice"], 1)
            review_h = round(hours_per_week * diff["review"], 1)

            plan.append(f"### 第 {current_week} 周 ({hours_per_week} 小时)")
            plan.append(f"| 类型 | 时长 | 内容 |")
            plan.append(f"|------|------|------|")
            plan.append(f"| 📖 理论学习 | {theory_h}h | 核心概念、教材阅读、视频课程 |")
            plan.append(f"| ✍️ 实践练习 | {practice_h}h | 习题、编程、案例分析 |")
            plan.append(f"| 🔄 复习巩固 | {review_h}h | 错题回顾、知识梳理、闪卡 |")
            plan.append(f"| ✅ 周目标 | - | 完成本周学习任务，通过自测 |")
            plan.append("")
            current_week += 1

    plan.append("---\n")
    plan.append("## 💡 学习建议\n")
    plan.append("- 🌅 **高效时间**: 利用早晨或精力充沛时段进行理论学习")
    plan.append("- ⏱ **番茄工作法**: 25分钟专注 + 5分钟休息，4个番茄后长休息")
    plan.append("- 📝 **费曼技巧**: 用简单语言解释概念，发现知识盲区")
    plan.append("- 🔁 **间隔重复**: 第1/3/7/14/30天复习已学内容")
    plan.append("- 📊 **进度追踪**: 每周日总结完成情况，调整下周计划")
    plan.append("- 🎯 **里程碑**: 每阶段结束时进行综合测试")

    full_plan = "\n".join(plan)
    sp_md = os.path.join(CUSTOM_TEMP_DIR, f"study_plan_{subject.replace(' ', '_')}.md")
    with open(sp_md, "w", encoding="utf-8") as f:
        f.write(full_plan)

    if output_format == "markdown":
        return f"✅ 学习计划已生成:\n   📄 {sp_md}\n\n{full_plan[:2000]}{'...' if len(full_plan) > 2000 else ''}"
    elif output_format == "pdf" and _PANDOC_PATH:
        pdf_path = os.path.join(CUSTOM_TEMP_DIR, f"study_plan_{subject.replace(' ', '_')}.pdf")
        cmd = [_PANDOC_PATH, sp_md, "-o", pdf_path]
        if _DETECTED_CJK_FONT:
            cmd.extend(["-V", f"CJKmainfont={_DETECTED_CJK_FONT}"])
        subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        return f"✅ 学习计划已生成:\n   📄 {pdf_path}"
    return f"✅ 学习计划已生成:\n   📄 {sp_md}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 23: 闪卡生成器
# ═══════════════════════════════════════════════════════════

@tool
def FlashcardTool(cards_json: str = "", topic: str = "", card_count: int = 20,
                  output_format: str = "csv") -> str:
    """生成 Anki/Quizlet 兼容闪卡。提供 cards_json: [{"front":"Q","back":"A"},...] 或 topic 自动生成题目框架。输出 CSV (Anki导入) 或 Markdown。"""
    cards = []

    if cards_json:
        try:
            cards = json.loads(cards_json)
        except json.JSONDecodeError:
            return "❌ cards_json 格式错误。请提供 JSON 数组: [{\"front\":\"问题\",\"back\":\"答案\"},...]"

    if not cards and topic:
        cards = [
            {"front": f"{topic} - 核心概念 1", "back": f"[请填写定义]\n\n相关: [关联概念]\n示例: [具体例子]"},
            {"front": f"{topic} - 核心概念 2", "back": f"[请填写定义]\n\n相关: [关联概念]\n示例: [具体例子]"},
            {"front": f"{topic} - 公式/定理 1", "back": f"公式: [LaTeX]\n\n适用条件: [条件]\n例题: [题目]"},
            {"front": f"{topic} - 常见问题", "back": f"问题: [常见错误]\n正确理解: [正确概念]\n记忆技巧: [技巧]"},
            {"front": f"{topic} - 对比辨析", "back": f"A: [概念A]\nB: [概念B]\n关键区别: [区别]"},
        ]
        for i in range(6, min(card_count + 1, 50)):
            cards.append({
                "front": f"{topic} - 知识点 {i}",
                "back": f"[请填写详细答案]\n\n关键词: [...]\n记忆口诀: [...]"
            })

    if not cards:
        return "❌ 请提供 cards_json 或 topic。"

    if output_format == "csv":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"flashcards_{topic.replace(' ', '_') if topic else 'export'}.csv")
        with open(sp, "w", encoding="utf-8-sig", newline="") as f:
            writer = csv_module.writer(f)
            writer.writerow(["Front", "Back", "Tags"])
            for c in cards:
                writer.writerow([c.get("front", ""), c.get("back", ""), c.get("tags", topic or "")])
        return f"✅ 闪卡 CSV 已生成 (Anki 兼容):\n   📄 {sp}\n   📇 {len(cards)} 张闪卡\n   💡 Anki导入: 文件→导入→选择CSV→分隔符逗号"

    elif output_format == "markdown":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"flashcards_{topic.replace(' ', '_') if topic else 'export'}.md")
        md = [f"# 闪卡: {topic or '导出'}\n"]
        for i, c in enumerate(cards, 1):
            md.append(f"## 卡片 {i}")
            md.append(f"**Q**: {c.get('front', '')}")
            md.append(f"**A**: {c.get('back', '')}")
            md.append("---\n")
        with open(sp, "w", encoding="utf-8") as f:
            f.write("\n".join(md))
        return f"✅ 闪卡 Markdown 已生成:\n   📄 {sp}\n   📇 {len(cards)} 张闪卡"

    return "❌ 不支持的格式。使用 'csv' 或 'markdown'。"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 24: 知识图谱生成器
# ═══════════════════════════════════════════════════════════

@tool
def KnowledgeGraphTool(concepts: str = "", central_topic: str = "",
                       relations: str = "", depth: int = 2) -> str:
    """从概念列表生成知识图谱 (Mermaid mindmap/flowchart)。concepts: 逗号分隔的概念列表；central_topic: 中心主题；relations: JSON格式关系 [{"from":"A","to":"B","label":"关系"},...]。"""
    if not central_topic and not concepts:
        return "❌ 请提供 central_topic 或 concepts。"

    if concepts:
        concept_list = [c.strip() for c in concepts.split(",") if c.strip()]
    else:
        concept_list = []

    if central_topic and not concept_list:
        concept_list = [central_topic]

    mermaid = "mindmap\n"
    if central_topic:
        mermaid += f"  root(({central_topic}))\n"
    else:
        mermaid += f"  root((知识图谱))\n"

    if relations:
        try:
            rels = json.loads(relations)
            from_map = defaultdict(list)
            for r in rels:
                from_map[r.get("from", "")].append(r)

            def add_nodes(node, current_depth=0):
                if current_depth >= depth: return
                children = from_map.get(node, [])
                indent = "    " * (current_depth + 1)
                for child in children[:8]:
                    to = child.get("to", "")
                    label = child.get("label", "")
                    label_str = f"[{label}]" if label else ""
                    mermaid_lines.append(f"{indent}{to}{label_str}")
                    add_nodes(to, current_depth + 1)

            mermaid_lines = [mermaid]
            if central_topic:
                add_nodes(central_topic)
            mermaid = "\n".join(mermaid_lines)
        except json.JSONDecodeError:
            pass

    if not relations:
        mermaid += "\n".join(f"    {c}" for c in concept_list[:15])

    flowchart = "graph TD\n"
    if central_topic:
        flowchart += f"    CENTER[{central_topic}]\n"
        for i, c in enumerate(concept_list[:12]):
            flowchart += f"    CENTER --> N{i}[{c}]\n"
    else:
        for i, c in enumerate(concept_list[:12]):
            flowchart += f"    N{i}[{c}]\n"
        for i in range(len(concept_list[:12]) - 1):
            flowchart += f"    N{i} --> N{i+1}\n"

    safe_name = central_topic.replace(' ', '_') if central_topic else 'export'
    sp_mmd = os.path.join(CUSTOM_TEMP_DIR, f"knowledge_graph_{safe_name}.mmd")
    with open(sp_mmd, "w", encoding="utf-8") as f:
        f.write(mermaid)

    sp_md = os.path.join(CUSTOM_TEMP_DIR, f"knowledge_graph_{safe_name}.md")
    with open(sp_md, "w", encoding="utf-8") as f:
        f.write(f"# 知识图谱: {central_topic or '导出'}\n\n")
        f.write("## Mindmap\n```mermaid\n")
        f.write(mermaid)
        f.write("\n```\n\n")
        f.write("## Flowchart\n```mermaid\n")
        f.write(flowchart)
        f.write("\n```\n")

    result = f"✅ 知识图谱已生成:\n   📄 Mermaid: {sp_mmd}\n   📄 Markdown: {sp_md}\n"
    if _MMDC_PATH:
        pp = os.path.join(CUSTOM_TEMP_DIR, f"knowledge_graph_{safe_name}.png")
        try:
            subprocess.run([_MMDC_PATH, "-i", sp_mmd, "-o", pp, "-w", "1600", "-b", "white"],
                           capture_output=True, text=True, timeout=30)
            if os.path.exists(pp):
                result += f"   🖼 PNG: {pp}\n"
        except: pass
    return result

# ═══════════════════════════════════════════════════════════
#  🆕 工具 25: 笔记整理器
# ═══════════════════════════════════════════════════════════

@tool
def NoteOrganizerTool(content: str = "", source_path: str = "",
                      mode: str = "summarize", style: str = "academic") -> str:
    """智能笔记整理。mode: summarize(摘要)/outline(大纲)/keywords(关键词提取)/mindmap(思维导图数据)/study_guide(学习指南)。
source_path: 可直接指定文件路径。"""
    if source_path and not content:
        sfp = _resolve_path(source_path)
        if os.path.exists(sfp):
            content = _safe_read_text(sfp)
        else:
            return f"❌ 源文件不存在: {source_path}"

    # [FIX#16] 防御 None
    if not content or not content.strip():
        return "❌ 内容为空。"

    lines = content.split("\n")
    words = content.split()
    chars = len(content)
    sentences = len(re.findall(r'[.!?。！？]+', content))

    info = [f"📝 笔记分析: {len(lines)}行 | {len(words)}词 | {chars}字符 | ~{sentences}句\n"]

    if mode == "keywords":
        words_list = re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z]{3,}', content)
        word_freq = Counter(w.lower() for w in words_list if len(w) >= 2)
        keywords = word_freq.most_common(30)
        info.append("## 🔑 关键词提取\n")
        info.append("| 关键词 | 频次 |")
        info.append("|--------|------|")
        for kw, freq in keywords:
            info.append(f"| {kw} | {freq} |")
        info.append(f"\n共提取 {len(keywords)} 个关键词")

    elif mode == "outline":
        info.append("## 📋 内容大纲\n")
        for line in lines[:100]:
            stripped = line.strip()
            if not stripped: continue
            if stripped.startswith("#"):
                info.append(f"\n{stripped}")
            elif re.match(r'^[A-Z][A-Za-z\s]{3,60}$', stripped):
                info.append(f"### {stripped}")
            elif len(stripped) > 20 and (stripped.endswith(":") or stripped.endswith("：")):
                info.append(f"- **{stripped}**")
            elif len(stripped) < 100:
                info.append(f"  - {stripped[:80]}{'...' if len(stripped) > 80 else ''}")

    elif mode == "summarize":
        info.append("## 📄 文本摘要\n")
        info.append(f"**统计**: {len(lines)} 行, {len(words)} 词, {chars} 字符\n")
        if len(lines) > 10:
            info.append(f"\n**开头**: {lines[0][:200]}...")
            info.append(f"\n**中间**: ...{lines[len(lines)//2][:200]}...")
            info.append(f"\n**结尾**: ...{lines[-1][:200]}")
        else:
            info.append(content[:500])
        info.append(f"\n\n💡 使用 'outline' 模式查看大纲，'keywords' 模式提取关键词。")

    elif mode == "study_guide":
        info.append("## 📖 学习指南\n")
        info.append("### 🔑 关键概念\n")
        def_patterns = [
            r'([^。！？\n]{3,30}(?:是|指|即|为|定义为|指的是)[^。！？\n]{3,100})',
        ]
        for pat in def_patterns:
            matches = re.findall(pat, content)
            for m in matches[:5]:
                info.append(f"- {m.strip()}")

        info.append("\n### ❓ 潜在考点\n")
        exam_patterns = [r'([^。！？\n]*(?:步骤|方法|原因|因素|特点|分类|类型|区别|对比|例如|包括)[^。！？\n]*)']
        for pat in exam_patterns:
            matches = re.findall(pat, content)
            for m in matches[:8]:
                info.append(f"- {m.strip()}")

        info.append("\n### 📊 知识框架\n")
        info.append("```mermaid")
        info.append("mindmap")
        info.append(f"  root((学习主题))")
        kw = re.findall(r'[\u4e00-\u9fff]{2,4}', content)
        for w in Counter(kw).most_common(8):
            info.append(f"    {w[0]}")
        info.append("```")

    elif mode == "mindmap":
        info.append("## 🧠 思维导图数据\n")
        info.append("```mermaid")
        info.append("mindmap")
        info.append(f"  root((笔记主题))")
        for line in lines[:30]:
            stripped = line.strip()
            if stripped and len(stripped) < 50:
                clean = re.sub(r'[#\-*>\[\]()]', '', stripped).strip()
                if clean:
                    info.append(f"    {clean[:40]}")
        info.append("```")

    full = "\n".join(info)
    sp = os.path.join(CUSTOM_TEMP_DIR, f"note_{mode}.md")
    with open(sp, "w", encoding="utf-8") as f:
        f.write(full)

    return f"✅ 笔记整理完成 ({mode}):\n   📄 {sp}\n\n{full[:2000]}{'...' if len(full) > 2000 else ''}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 26: 代码审查器
# ═══════════════════════════════════════════════════════════

@tool
def CodeReviewTool(code: str = "", file_path: str = "",
                   review_focus: str = "all") -> str:
    """代码审查与优化建议。review_focus: all/security/performance/style/complexity/bugs。"""
    if file_path and not code:
        sfp = _resolve_path(file_path)
        if os.path.exists(sfp):
            code = _safe_read_text(sfp)
        else:
            return f"❌ 文件不存在: {file_path}"

    # [FIX#16] 防御 None
    if not code or not code.strip():
        return "❌ 代码为空。"

    lines = code.split("\n")
    line_count = len(lines)
    char_count = len(code)

    review = []
    review.append(f"# 🔍 代码审查报告\n")
    review.append(f"**文件**: {file_path or '(直接输入)'}")
    review.append(f"**规模**: {line_count} 行 | {char_count} 字符\n")

    functions = re.findall(r'def\s+(\w+)\s*\(', code)
    classes = re.findall(r'class\s+(\w+)', code)
    imports = re.findall(r'^(?:import|from)\s+[\w.]+', code, re.MULTILINE)
    comments = [l for l in lines if l.strip().startswith('#')]
    todos = [l for l in lines if 'TODO' in l or 'FIXME' in l or 'HACK' in l]

    review.append("## 📊 代码统计")
    review.append(f"- 函数: {len(functions)} 个")
    review.append(f"- 类: {len(classes)} 个")
    review.append(f"- 导入: {len(imports)} 条")
    review.append(f"- 注释: {len(comments)} 行 ({len(comments)/max(1,line_count)*100:.1f}%)")
    review.append(f"- 待办项: {len(todos)} 个\n")

    if review_focus in ("all", "complexity"):
        review.append("## 🔢 复杂度分析")
        for func in functions[:10]:
            func_body = code
            func_start = code.find(f"def {func}(")
            if func_start != -1:
                colon_pos = code.find(":", func_start)
                if colon_pos != -1:
                    rest = code[colon_pos+1:]
                    next_def = re.search(r'\n(?!(?:    |\t))def\s+\w', rest)
                    next_class = re.search(r'\n(?!(?:    |\t))class\s+\w', rest)
                    end = len(rest)
                    if next_def: end = min(end, next_def.start())
                    if next_class: end = min(end, next_class.start())
                    func_body = rest[:end]
            complexity = 1
            for pattern in [r'\bif\b', r'\bfor\b', r'\bwhile\b', r'\band\b', r'\bor\b',
                           r'\bexcept\b', r'\belif\b']:
                complexity += len(re.findall(pattern, func_body))
            level = "🟢 低" if complexity < 5 else ("🟡 中" if complexity < 10 else "🔴 高")
            review.append(f"- `{func}()`: 圈复杂度 ≈ {complexity} {level}")
        review.append("")

    if review_focus in ("all", "security"):
        review.append("## 🔒 安全检查")
        security_issues = []
        dangerous_patterns = [
            (r'os\.system\(', "使用 os.system() 存在命令注入风险"),
            (r'eval\(', "使用 eval() 可执行任意代码"),
            (r'exec\(', "使用 exec() 可执行任意代码"),
            (r'pickle\.loads?\(', "pickle 反序列化不安全"),
            (r'input\(', "input() 在非交互环境可能导致问题"),
            (r'password\s*=\s*[\'"]', "硬编码密码"),
            (r'secret\s*=\s*[\'"]', "硬编码密钥"),
            (r'\.execute\(.*%', "SQL 拼接存在注入风险"),
        ]
        for pat, desc in dangerous_patterns:
            if re.search(pat, code):
                security_issues.append(f"- ⚠️ {desc}: `{pat}`")
        if security_issues:
            review.extend(security_issues)
        else:
            review.append("- ✅ 未发现常见安全漏洞")
        review.append("")

    if review_focus in ("all", "performance"):
        review.append("## ⚡ 性能建议")
        perf_issues = []
        perf_patterns = [
            (r'\.copy\(\)\s*$', "不必要的 copy() 操作可能浪费内存"),
            (r'for\s+\w+\s+in\s+range\(len\(', "考虑使用 enumerate() 替代 range(len())"),
            (r'\+\s*=\s*[\'"]', "循环中字符串拼接建议使用 join()"),
            (r'\.readlines\(\)', "大文件使用 readlines() 可能消耗大量内存"),
        ]
        for pat, desc in perf_patterns:
            if re.search(pat, code):
                perf_issues.append(f"- 💡 {desc}")
        if perf_issues:
            review.extend(perf_issues)
        else:
            review.append("- ✅ 未发现明显性能问题")
        review.append("")

    if review_focus in ("all", "style"):
        review.append("## 🎨 风格检查")
        style_issues = []
        if line_count > 0:
            long_lines = [i+1 for i, l in enumerate(lines) if len(l) > 100]
            if long_lines:
                style_issues.append(f"- ⚠️ {len(long_lines)} 行超过 100 字符: 行 {long_lines[:5]}{'...' if len(long_lines)>5 else ''}")
            trailing_ws = sum(1 for l in lines if l.endswith(' ') or l.endswith('\t'))
            if trailing_ws:
                style_issues.append(f"- ⚠️ {trailing_ws} 行有尾随空白")
            snake_breaks = [n for n in functions if not re.match(r'^[a-z][a-z0-9_]*$', n)]
            if snake_breaks:
                style_issues.append(f"- 💡 函数命名建议 snake_case: {snake_breaks[:5]}")
        if style_issues:
            review.extend(style_issues)
        else:
            review.append("- ✅ 代码风格良好")
        review.append("")

    if todos:
        review.append("## 📌 待办事项")
        for todo in todos[:10]:
            review.append(f"- {todo.strip()}")
        review.append("")

    full_review = "\n".join(review)
    sp = os.path.join(CUSTOM_TEMP_DIR, "code_review.md")
    with open(sp, "w", encoding="utf-8") as f:
        f.write(full_review)

    return f"✅ 代码审查完成:\n   📄 {sp}\n\n{full_review[:2000]}{'...' if len(full_review) > 2000 else ''}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 27: LaTeX 数学公式渲染
# ═══════════════════════════════════════════════════════════

@tool
def MathRenderTool(latex_expression: str, filename: str = "math_render",
                   font_size: int = 20, dpi: int = 150) -> str:
    """LaTeX 数学公式渲染为 PNG 图片。支持行内 $...$ 或行间 $$...$$。需要 matplotlib。"""
    if not latex_expression:
        return "❌ LaTeX 表达式不能为空。"
    if not _MPL_AVAILABLE:
        return "❌ matplotlib 未安装。公式渲染需要 matplotlib。"

    expr = latex_expression.strip()
    expr = re.sub(r'^\$\$?\s*', '', expr)
    expr = re.sub(r'\s*\$\$?$', '', expr)

    sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.png")

    try:
        fig, ax = plt.subplots(figsize=(len(expr)*0.15 + 2, 1.5))
        ax.axis('off')
        ax.text(0.5, 0.5, f"${expr}$", fontsize=font_size,
                ha='center', va='center', transform=ax.transAxes)
        fig.savefig(sp, dpi=dpi, bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close(fig)
        return f"✅ 公式已渲染:\n   🖼 {sp}\n   📐 公式: ${expr}$"
    except Exception as e:
        if _PIL_AVAILABLE:
            try:
                img = Image.new('RGB', (800, 120), color='white')
                draw = ImageDraw.Draw(img)
                draw.text((10, 50), f"${expr}$", fill='black')
                img.save(sp)
                plt.close('all')
                return f"⚠️ matplotlib渲染失败，使用PIL回退:\n   🖼 {sp}\n   📐 公式: ${expr}$"
            except: pass
        return f"❌ 公式渲染失败: {e}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 28: 引用网络分析
# ═══════════════════════════════════════════════════════════

@tool
def CitationNetworkTool(doi_list: str = "", search_query: str = "",
                        max_papers: int = 10) -> str:
    """引用网络分析。输入 DOI 列表(逗号分隔)或搜索查询，分析论文间的引用关系，生成 Mermaid 引用图。"""
    if not doi_list and not search_query:
        return "❌ 请提供 doi_list 或 search_query。"

    papers = []

    if doi_list:
        dois = [d.strip() for d in doi_list.split(",") if d.strip()]
        for doi in dois[:max_papers]:
            doi_clean = doi.replace("https://doi.org/","").replace("http://dx.doi.org/","")
            url = f"https://api.crossref.org/works/{urllib.parse.quote(doi_clean, safe='')}"
            try:
                resp = _safe_request(url, timeout=10)
                if resp:
                    msg = resp.json().get("message", {})
                    title = " ".join((msg.get("title",["?"])[0] or "?")[:80].split())
                    first_author = "?"
                    if msg.get("author"):
                        a = msg["author"][0]
                        first_author = a.get("family", "?")
                    # Safe year extraction
                    year = "?"
                    for date_key in ("published-print", "created"):
                        dp = msg.get(date_key, {}).get("date-parts", [])
                        if dp and isinstance(dp, list):
                            for part in dp:
                                if part and isinstance(part, list) and len(part) > 0 and part[0] is not None:
                                    year = part[0]
                                    break
                        if year != "?":
                            break
                    refs = msg.get("reference", [])
                    ref_dois = [r.get("DOI","") for r in refs if r.get("DOI")]
                    papers.append({
                        "doi": doi_clean, "title": title,
                        "author": first_author, "year": year,
                        "refs": ref_dois[:20]
                    })
            except Exception: pass

    if search_query and not papers:
        base = "https://api.semanticscholar.org/graph/v1/paper/search"
        params = {"query": search_query, "limit": max_papers,
                  "fields": "title,authors,year,externalIds,citations"}
        url = f"{base}?{urllib.parse.urlencode(params)}"
        try:
            resp = _safe_request(url, timeout=15)
            if resp:
                data = resp.json()
                for p in data.get("data", []):
                    ext = p.get("externalIds", {})
                    doi = ext.get("DOI", "")
                    if doi:
                        papers.append({
                            "doi": doi,
                            "title": p.get("title","?")[:80],
                            "author": (p.get("authors",[{}])[0].get("name","?") if p.get("authors") else "?"),
                            "year": p.get("year","?"),
                            "refs": []
                        })
        except Exception: pass

    if not papers:
        return "❌ 未找到相关论文。请检查 DOI 或搜索词。"

    mermaid = "graph LR\n"
    for i, p in enumerate(papers):
        short_title = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fff]', '_', p["title"][:20])
        node_id = f"P{i}"
        mermaid += f"    {node_id}[{p['author']} ({p['year']})\n{short_title}]\n"
        for ref_doi in p.get("refs", [])[:3]:
            for j, other in enumerate(papers):
                if other["doi"] == ref_doi and j != i:
                    mermaid += f"    {node_id} -->|引用| P{j}\n"

    sp = os.path.join(CUSTOM_TEMP_DIR, "citation_network.md")
    md = [f"# 引用网络分析\n",
          f"共 {len(papers)} 篇论文\n",
          "## 论文列表\n"]
    for i, p in enumerate(papers):
        md.append(f"{i+1}. **{p['title']}** — *{p['author']}* ({p['year']})")
        md.append(f"   DOI: {p['doi']}")
    md.append("\n## 引用关系图\n")
    md.append("```mermaid")
    md.append(mermaid)
    md.append("```")

    with open(sp, "w", encoding="utf-8") as f:
        f.write("\n".join(md))

    result = f"✅ 引用网络分析完成:\n   📄 {sp}\n   📊 {len(papers)} 篇论文\n"
    if _MMDC_PATH:
        pp = os.path.join(CUSTOM_TEMP_DIR, "citation_network.png")
        mp = os.path.join(CUSTOM_TEMP_DIR, "citation_network.mmd")
        with open(mp, "w", encoding="utf-8") as f:
            f.write(mermaid)
        try:
            subprocess.run([_MMDC_PATH, "-i", mp, "-o", pp, "-w", "1600", "-b", "white"],
                           capture_output=True, text=True, timeout=30)
            if os.path.exists(pp):
                result += f"   🖼 引用图: {pp}\n"
        except: pass
    return result

# ═══════════════════════════════════════════════════════════
#  🆕 工具 29: 项目脚手架生成
# ═══════════════════════════════════════════════════════════

@tool
def ProjectScaffoldTool(project_name: str, project_type: str = "python",
                        with_tests: bool = True, with_docs: bool = True) -> str:
    """生成项目脚手架目录结构。project_type: python/r-latex/paper/website。"""
    # [FIX#16] 防御 None
    if not project_name or not project_name.strip():
        return "❌ 项目名称不能为空。"

    # safe package name
    safe_name = re.sub(r'[^a-zA-Z0-9_]', '_', project_name)
    if safe_name and safe_name[0].isdigit():
        safe_name = 'pkg_' + safe_name
    safe_name = safe_name.lower()
    project_dir = os.path.join(CUSTOM_TEMP_DIR, safe_name)
    created = []

    templates = {
        "python": {
            "dirs": [safe_name, "tests", "docs", "data"],
            "files": {
                "README.md": f"# {project_name}\n\n## 简介\n...\n\n## 安装\n```bash\npip install -e .\n```\n\n## 使用\n```python\nimport {safe_name}\n```\n",
                "setup.py": f"from setuptools import setup, find_packages\n\nsetup(\n    name='{safe_name}',\n    version='0.1.0',\n    packages=find_packages(),\n    install_requires=[],\n)\n",
                f"{safe_name}/__init__.py": f"\"\"\"{project_name} - 项目描述\"\"\"\n__version__ = \"0.1.0\"\n",
                f"{safe_name}/main.py": f"\"\"\"主模块\"\"\"\n\ndef main():\n    print(\"Hello from {project_name}!\")\n\nif __name__ == \"__main__\":\n    main()\n",
                ".gitignore": "__pycache__/\n*.pyc\n.env\nvenv/\n.venv/\ndist/\nbuild/\n*.egg-info/\n",
            }
        },
        "r-latex": {
            "dirs": ["figures", "tables", "sections", "bib"],
            "files": {
                "main.tex": "\\documentclass{article}\n\\usepackage[UTF8]{ctex}\n\\title{" + project_name + "}\n\\author{作者}\n\\date{\\today}\n\n\\begin{document}\n\\maketitle\n\\section{引言}\n...\n\\end{document}\n",
                "README.md": f"# {project_name}\n\nLaTeX 项目。\n",
            }
        },
        "paper": {
            "dirs": ["figures", "data", "tables", "supplementary"],
            "files": {
                "paper.md": f"# {project_name}\n\n## Abstract\n...\n\n## Introduction\n...\n\n## Methods\n...\n\n## Results\n...\n\n## Discussion\n...\n\n## References\n...\n",
                "README.md": f"# {project_name}\n\n学术论文项目。\n",
            }
        },
        "website": {
            "dirs": ["css", "js", "images", "pages"],
            "files": {
                "index.html": "<!DOCTYPE html>\n<html lang=\"zh\">\n<head>\n    <meta charset=\"UTF-8\">\n    <title>" + project_name + "</title>\n    <link rel=\"stylesheet\" href=\"css/style.css\">\n</head>\n<body>\n    <h1>Welcome to " + project_name + "</h1>\n</body>\n</html>\n",
                "css/style.css": "/* " + project_name + " Styles */\nbody {\n    font-family: sans-serif;\n    max-width: 800px;\n    margin: 0 auto;\n    padding: 2em;\n}\n",
                "README.md": f"# {project_name}\n\n静态网站项目。\n",
            }
        },
    }

    config = templates.get(project_type, templates["python"])

    try:
        os.makedirs(project_dir, exist_ok=True)
        for d in config.get("dirs", []):
            dp = os.path.join(project_dir, d)
            os.makedirs(dp, exist_ok=True)
            created.append(f"📁 {d}/")

        for fp, content in config.get("files", {}).items():
            full_path = os.path.join(project_dir, fp)
            os.makedirs(os.path.dirname(full_path), exist_ok=True)
            with open(full_path, "w", encoding="utf-8") as f:
                f.write(content)
            created.append(f"📄 {fp}")

        if with_tests and project_type == "python":
            test_dir = os.path.join(project_dir, "tests")
            os.makedirs(test_dir, exist_ok=True)
            with open(os.path.join(test_dir, "test_main.py"), "w", encoding="utf-8") as f:
                f.write(f"\"\"\"{project_name} 测试\"\"\"\nimport pytest\nfrom {safe_name} import main\n\ndef test_main():\n    assert main is not None\n")
            created.append("📄 tests/test_main.py")

        if with_docs and project_type == "python":
            docs_dir = os.path.join(project_dir, "docs")
            os.makedirs(docs_dir, exist_ok=True)
            with open(os.path.join(docs_dir, "index.md"), "w", encoding="utf-8") as f:
                f.write(f"# {project_name} 文档\n\n...\n")
            created.append("📄 docs/index.md")

        return f"✅ 项目脚手架已生成: {project_dir}\n   类型: {project_type}\n   创建了 {len(created)} 个文件/目录:\n" + "\n".join(f"   {c}" for c in created[:30])
    except Exception as e:
        return f"❌ 项目生成失败: {e}"


# ═══════════════════════════════════════════════════════════
#  4. 工具注册 & Agent 构建 (v4.0.1 — 30 个工具)
# ═══════════════════════════════════════════════════════════

safe_write_tool = WriteFileTool(root_dir=CUSTOM_TEMP_DIR) if _LANGCHAIN_AVAILABLE else None

tools = [
    # 基础工具 (1-10)
    SmartReadPathTool, EditTexFileTool, CompileLatexTool,
    PythonSandboxTool, DuckDuckGoSearchTool, FetchWebImageTool,
    EditWordDocTool, SaveMarkdownTool, ChartGenerationTool,
    MermaidTool,
    # 学术检索 (11-16)
    ArXivSearchTool, DOIMetadataTool, DataStatisticsTool,
    BibTexTool, PDFTableExtractTool, SessionExportTool,
    # v3.1 工具 (17-21)
    PandocConvertTool, PDFAnnotExtractTool, AcademicTranslateTool,
    SemanticScholarTool, PresentationGenTool,
    # v4.0 新增 (22-29)
    StudyPlanTool, FlashcardTool, KnowledgeGraphTool,
    NoteOrganizerTool, CodeReviewTool, MathRenderTool,
    CitationNetworkTool, ProjectScaffoldTool,
]
if safe_write_tool:
    tools.append(safe_write_tool)

conn = sqlite3.connect(DB_PATH, check_same_thread=False)
memory = SqliteSaver(conn)

# -- 自定义 Agent 图 --
class _AgentState(TypedDict):
    messages: Annotated[list, add_messages]

def _call_model(state: _AgentState) -> dict:
    messages = state["messages"]

    # [FIX#18] 增强消息验证：处理空消息和孤立的 tool_calls
    if not messages:
        _log("WARN", "_call_model 收到空消息列表，返回默认响应")
        return {"messages": [AIMessage(content="⚠️ 内部错误：空消息列表。请重新发送请求。")]}

    _validated = []
    _pending_tool_calls = 0
    for _msg in messages:
        _validated.append(_msg)
        if isinstance(_msg, AIMessage) and getattr(_msg, 'tool_calls', None):
            _pending_tool_calls += len(_msg.tool_calls)
        elif isinstance(_msg, ToolMessage):
            _pending_tool_calls -= 1

    if _pending_tool_calls > 0:
        _log("WARN", f"检测到 {_pending_tool_calls} 个未配对的 tool_calls，尝试清理...")
        _cleaned = []
        for _i in range(len(_validated) - 1, -1, -1):
            _msg = _validated[_i]
            if isinstance(_msg, AIMessage) and getattr(_msg, 'tool_calls', None):
                _tc_count = len(_msg.tool_calls)
                _pending_tool_calls -= _tc_count
                if _pending_tool_calls >= 0:
                    continue
            elif isinstance(_msg, ToolMessage):
                _pending_tool_calls += 1
            _cleaned.insert(0, _msg)
        messages = _cleaned
        _log("INFO", f"清理后消息数: {len(messages)}")

    try:
        llm_with_tools = llm.bind_tools(tools)
        response = llm_with_tools.invoke(messages)
        return {"messages": [response]}
    except Exception as _e:
        _log("ERR", f"LLM 调用失败: {_e}")
        error_msg = AIMessage(content=f"❌ API 调用失败: {str(_e)[:300]}\n请检查 API Key 或稍后重试。")
        return {"messages": [error_msg]}

def _should_continue(state: _AgentState) -> str:
    messages = state["messages"]
    # [FIX#17] 空消息列表保护
    if not messages:
        _log("WARN", "_should_continue 收到空消息列表，终止流程")
        return END
    last_message = messages[-1]
    if isinstance(last_message, AIMessage) and getattr(last_message, 'tool_calls', None):
        return "tools"
    return END

_workflow = StateGraph(_AgentState)
_workflow.add_node("agent", _call_model)
_workflow.add_node("tools", ToolNode(tools))
_workflow.set_entry_point("agent")
_workflow.add_conditional_edges(
    "agent",
    _should_continue,
    {"tools": "tools", END: END}
)
_workflow.add_edge("tools", "agent")
agent = _workflow.compile(checkpointer=memory)

# ═══════════════════════════════════════════════════════════
#  5. 终端交互主循环 (v4.0.1)
# ═══════════════════════════════════════════════════════════

HELP_TEXT = f"""
{_C.BOLD}╔══════════════════ v4.0.1 工具一览 (29个) ═══════════════════╗{_C.W}
{_C.G}║ 📖 SmartReadPathTool       读取文件/文件夹              ║
║ 📝 EditTexFileTool          创建/编辑 .tex                ║
║ 🖨  CompileLatexTool         编译 .tex→PDF (+bibtex)      ║
║ 🐍 PythonSandboxTool        执行 Python 代码              ║
║ 🔍 DuckDuckGoSearchTool     通用网页搜索                  ║
║ 🖼  FetchWebImageTool        搜索下载学术图片              ║
║ 📄 EditWordDocTool          创建/编辑 Word 文档           ║
║ 📋 SaveMarkdownTool         保存 Markdown 文件            ║
║ 📊 ChartGenerationTool      matplotlib/seaborn 图表       ║
║ 🎨 MermaidTool              Mermaid 流程图                ║{_C.W}
{_C.Y}║ 📚 ArXivSearchTool          arXiv 学术论文检索            ║
║ 📄 DOIMetadataTool          DOI → 完整元数据              ║
║ 📈 DataStatisticsTool       CSV/Excel 自动统计分析        ║
║ 📖 BibTexTool               .bib解析/生成/引用格式化      ║
║ 📋 PDFTableExtractTool      PDF 表格智能提取              ║
║ 💾 SessionExportTool        对话记录导出                  ║{_C.W}
{_C.M}║ 🔄 PandocConvertTool        Markdown/LaTeX→PDF/DOCX/HTML ║
║ 📝 PDFAnnotExtractTool      PDF 批注/高亮提取             ║
║ 🌐 AcademicTranslateTool    学术中英互译                  ║
║ 🔬 SemanticScholarTool      Semantic Scholar 论文检索      ║
║ 🖥  PresentationGenTool      Markdown→PPTX (6主题) 🚀     ║{_C.W}
{_C.B}║ ─────────── 🆕 v4.0 新增 ───────────                  ║
║ 📚 StudyPlanTool            个性化学习计划生成             ║
║ 🃏 FlashcardTool            Anki/Quizlet 闪卡生成         ║
║ 🧠 KnowledgeGraphTool       知识图谱 (Mermaid)             ║
║ 📝 NoteOrganizerTool        智能笔记整理/摘要/大纲         ║
║ 🔍 CodeReviewTool           代码审查与优化建议             ║
║ 📐 MathRenderTool           LaTeX 数学公式渲染             ║
║ 🔗 CitationNetworkTool      引用网络分析                   ║
║ 🏗  ProjectScaffoldTool      项目脚手架生成                 ║{_C.W}
{_C.BOLD}╚══════════════════════════════════════════════════════════════╝{_C.W}
🚀 = v4.0 大幅增强 | 🐛 = v4.0.1 修复
"""

print(f"\n{_C.BOLD}💡 'exit' 退出 | 'new' 新对话 | 'help' 查看29个工具 | 'status' 系统状态{_C.W}\n")

current_thread_id = str(uuid.uuid4())
config = {"configurable": {"thread_id": current_thread_id}}

while True:
    try:
        user_input = input(f"{_C.C}👤 你:{_C.W} ").strip()
        if user_input.lower() in {"exit", "quit", "退出"}:
            _log("INFO", "👋 再见！"); conn.close(); break
        if user_input.lower() in {"new", "clear", "清空", "新对话"}:
            current_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": current_thread_id}}
            _log("INFO", "✨ 全新对话已开启。\n"); continue
        if user_input.lower() == "help":
            print(HELP_TEXT); continue
        if user_input.lower() == "status":
            for k, v in _health.items():
                print(f"  {'✅' if v else '⚠️'} {k}")
            continue
        if not user_input: continue

        print(f"\n{_C.M}🤔 思考/执行中...{_C.W}\n")
        inputs = {"messages": [("user", user_input)]}
        t_start = time.time()

        for event in agent.stream(inputs, config, stream_mode="updates"):
            if "agent" in event:
                messages = event["agent"].get("messages", [])
                for msg in messages:
                    reasoning = getattr(msg, 'reasoning_content', None)
                    if not reasoning:
                        reasoning = msg.additional_kwargs.get('reasoning_content', '') if hasattr(msg, 'additional_kwargs') else ''
                    if reasoning:
                        r_str = str(reasoning)
                        print(f"{_C.DIM}💭 思考:{_C.W} {r_str[:500]}{'...' if len(r_str)>500 else ''}")

                    if isinstance(msg, AIMessage) and msg.content:
                        elapsed = time.time() - t_start
                        print(f"{_C.G}🤖 DeepSeek{_C.W} ({elapsed:.1f}s):\n{msg.content}\n")

                    if isinstance(msg, AIMessage) and msg.tool_calls:
                        for tc in msg.tool_calls:
                            tc_name = tc.get('name', tc.get('function', {}).get('name', 'unknown'))
                            tc_args = tc.get('args', tc.get('function', {}).get('arguments', {}))
                            if isinstance(tc_args, str):
                                try: tc_args = json.loads(tc_args)
                                except: pass
                            args_str = json.dumps(tc_args, ensure_ascii=False, default=str) if tc_args else "{}"
                            if len(args_str) > 200: args_str = args_str[:200] + "..."
                            print(f"{_C.Y}🛠  [{tc_name}]{_C.W}\n   参数: {args_str}")
            elif "tools" in event:
                messages = event["tools"].get("messages", [])
                for msg in messages:
                    if isinstance(msg, ToolMessage):
                        cs = str(msg.content) if msg.content else ""
                    elif hasattr(msg, 'content'):
                        cs = str(msg.content) if msg.content else ""
                    else:
                        cs = str(msg)
                    cl = len(cs)
                    preview = cs[:300].replace("\n", " ") + ("..." if cl > 300 else "")
                    is_err = cs.startswith("❌") or "错误" in cs[:100]
                    icon = _C.R+"📄" if is_err else _C.G+"📄"
                    print(f"{icon} [返回 {cl}字符]{_C.W}: {preview}\n")
    except KeyboardInterrupt:
        print(f"\n{_C.Y}⏸  [操作被用户中断]{_C.W}")
    except Exception as e:
        print(f"{_C.R}❌ 发生错误: {e}{_C.W}")
        _log("ERR", traceback.format_exc())

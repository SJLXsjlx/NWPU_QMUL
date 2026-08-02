import os
import sys
import uuid
import sqlite3
import pandas as pd
import docx
import requests
from docx import Document
from docx.shared import Inches
from duckduckgo_search import DDGS
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_community.tools import WriteFileTool
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.sqlite import SqliteSaver
from pptx import Presentation
from PIL import Image
import pytesseract
import fitz
# ================== Tesseract OCR 路径配置 ==================
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ================== 1. 核心目录与环境配置 ==================
load_dotenv(override=True)

api_key = os.getenv("DEEPSEEK_FOR_THE_BREAK")
if not api_key or not api_key.startswith("sk-"):
    print("❌ API Key 未加载或格式不正确！请检查 .env 文件。")
    sys.exit(1)

# 定义工作区与临时文件目录
BASE_WORKSPACE = r"D:\.API Keys\DeepSeek_For_the_Break\outputs"
CUSTOM_TEMP_DIR = os.path.join(BASE_WORKSPACE, "temp")
os.makedirs(CUSTOM_TEMP_DIR, exist_ok=True)

print("✅ ds学术已启动（全面强化 PDF/Word/PPT/Excel/LaTeX 读取与写入稳定性）")

# ================== 2. 系统 Prompt 与模型 ==================
SYSTEM_PROMPT = SYSTEM_PROMPT = """你是一位高度专业、严谨且全面的**学术AI助手**（Academic AI Assistant），专注于为科研人员、学生和学者提供全方位的学术支持。你的核心定位是：

- **研究辅助专家**：帮助进行文献分析、论文构思与写作、数据处理与可视化、实验设计建议、学术幻灯片/报告制作等。
- **多格式文档处理专家**：精通PDF（使用PyMuPDF精确提取布局与文本）、Word、PowerPoint、Excel、LaTeX源码、图片OCR（中英双语）等文件的深度读取、解析与内容综合。
- **学术排版与生成专家**：擅长使用LaTeX生成高质量、可直接编译的学术论文、报告、公式排版；同时支持Word文档的结构化创建与编辑。
- **视觉辅助支持**：可通过工具搜索并下载与学术主题相关的专业图片、图表或示意图。
- **代码与工具集成专家**：能编写科研相关Python脚本、数据分析流程，并通过安全工具持久化保存所有生成内容到指定工作区。

**总体行为准则**（必须严格遵守）：

1. **客观严谨，结构清晰**：所有回答必须客观、中立、逻辑严密、层次分明。主动承认证据局限性、知识时效性或工具输出的不确定性。避免任何形式的猜测或过度泛化。

2. **【核心铁律 - 工具调用严格限制】**：
   - **SmartReadPathTool（文件/文件夹读取工具）** 是唯一允许读取本地文件的工具。
   - **只有在以下两种情况才允许调用**：
     a) 用户**明确提供了具体文件路径或文件夹路径**；
     b) 用户**显式要求**“读取”、“查看”、“打开”、“分析这个文件”、“检查路径下的内容”等操作。
   - **绝对禁止**：用户仅要求“写一个程序”、“解释某个概念”、“如何实现XX”、“总结XX领域”或进行任何常规对话时，**绝不能主动猜测路径、调用读取工具或编造文件位置**。此时应直接基于你的知识输出代码、解释或文本。
   - 当用户提供真实目录/文件夹时，必须先调用SmartReadPathTool进行递归遍历，然后**结合各子文件夹的相对路径结构 + 文件内容**进行整体归纳、交叉分析或针对性解答，并在回答中引用具体相对路径以便用户追溯。
   - 如果工具返回错误（如“找不到路径”），**立即停止一切路径猜测尝试**，直接基于已有知识回答，或礼貌询问用户提供正确路径。绝不反复尝试修改路径重试。

3. **学术写作与持久化输出优先**：
   - 涉及撰写学术论文、学位论文、研究报告、技术文档、或包含大量数学公式/方程的内容时，**必须主动、优先使用 EditTexFileTool** 生成结构完整、排版专业、可直接用 pdflatex/xelatex 编译的 .tex 源码。
   - 对于需要富文本、图片插入、表格的文档，优先使用 EditWordDocTool 创建或追加内容。
   - 所有生成的内容（代码、分析、文档）都应通过相应工具保存到 CUSTOM_TEMP_DIR 或工作区，实现持久化，方便用户后续使用。

4. **错误处理与用户引导**：
   - 任何工具调用失败或参数错误时，清晰告知用户具体问题，提供替代方案或请求澄清。
   - 永远不要为了“完成任务”而编造信息或强行调用工具。

5. **响应风格与语言**：
   - 默认使用与用户查询相同的语言（本系统以中文为主）。
   - 采用正式、专业的学术语气。
   - 结构化输出：善用标题、 bullet points、编号列表、代码块、表格和 Markdown 格式提升可读性。
   - 复杂任务时，先在思考中拆解步骤，再逐步执行或输出。
   - 主动提供后续建议，如“需要我将以上分析保存为Word文档吗？”或直接调用工具持久化。

6. **多轮对话与记忆管理**：
   - 利用系统提供的 SQLite 持久化记忆（SqliteSaver）保持上下文连贯性。
   - 用户输入 'new'、'clear' 等时，系统会重置 thread_id，你应配合给出全新对话的提示。

7. **学术诚信与透明度**：
   - 明确区分“基于用户提供文件的内容”和“基于我的训练知识”的信息来源。
   - 对于OCR结果、PDF提取等，提醒可能存在的识别误差。
   - 不协助任何非学术或违反学术伦理的任务。
   - 主动说明工具的沙箱限制（所有文件操作仅限于指定工作区）。

8. **任务执行优先级**：
   - 文件读取 → 内容综合分析 → 结构化输出/保存
   - 学术写作 → 优先 LaTeX 工具
   - 视觉需求 → FetchWebImageTool（需提供英文关键词）
   - 通用代码/文本 → safe_write_tool 保存

请始终进行**逐步推理（Chain-of-Thought）**后再决定是直接回复还是调用工具。你的目标是成为用户在学术道路上最可靠、最高效、最严谨的智能伙伴，帮助他们在科研中实现突破（For the Break）。

严格遵守以上所有规则，这是构建用户信任和系统稳定性的基础。"""

llm = ChatOpenAI(
    model="deepseek-v4-pro",   
    base_url="https://api.deepseek.com",
    api_key=api_key,
    temperature=0.5,
    max_tokens=4096,
    model_kwargs={
        "reasoning_effort": "max",
        "extra_body": {
            "thinking": {"type": "enabled"}
        }
    }
)

# ================== 3. 自定义工具箱 ==================

@tool
def SmartReadPathTool(path: str) -> str:
    """
    【仅在用户明确给出路径时触发】读取指定路径的内容。支持单文件读取，也支持整个文件夹的深度递归遍历。
    支持格式: .txt, .md, .py, .tex, .json, .csv, .pdf, .docx, .xlsx, .pptx, .png, .jpg, .jpeg, .bmp。
    警告：绝对不要编造或猜测路径！如果没有明确的路径输入，请勿调用此工具！
    """
    if not os.path.exists(path):
        return f"错误：找不到指定的路径 '{path}'。请停止尝试读取，直接回答用户的问题或要求用户提供正确路径。"

    def _read_single_file(file_path: str) -> str:
        ext = file_path.lower().split('.')[-1]
        
        # 1. 文本与代码类 (已增加 .tex 支持，并保留 GBK 编码容错)
        if ext in ['txt', 'md', 'py', 'json', 'tex']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='gbk', errors='ignore') as f:
                    return f.read()

        # 1.5 CSV 单独处理
        elif ext == 'csv':
            try:
                return pd.read_csv(file_path, encoding='utf-8').to_string()
            except UnicodeDecodeError:
                return pd.read_csv(file_path, encoding='gbk').to_string()
                
        # 2. PDF 类
        elif ext == 'pdf':
            text = []
            try:
                doc = fitz.open(file_path)
                for i, page in enumerate(doc):
                    page_text = page.get_text("text")
                    if page_text.strip():
                        text.append(f"--- [Page {i+1}] ---\n{page_text.strip()}")
                return "\n\n".join(text) if text else "[空PDF或纯图片扫描件，请尝试使用OCR]"
            except Exception as e:
                return f"❌ PDF 解析错误: {str(e)}"
            
        # 3. Word 类
        elif ext == 'docx':
            try:
                doc = docx.Document(file_path)
                content = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        content.append(p.text.strip())
                if doc.tables:
                    content.append("\n[文档表格数据]:")
                    for table in doc.tables:
                        for row in table.rows:
                            row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells if cell.text.strip()]
                            if row_data:
                                content.append(" | ".join(row_data))
                return "\n".join(content) if content else "[空Word文档]"
            except Exception as e:
                return f"❌ DOCX 解析错误: {str(e)}"
            
        # 4. Excel 类
        elif ext in ['xlsx', 'xls']:
            try:
                excel_file = pd.ExcelFile(file_path)
                sheet_outputs = []
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    if len(df) > 500:
                        sheet_outputs.append(f"--- Sheet: {sheet_name} (数据量超500行，为保护内存仅截取前500行) ---\n{df.head(500).to_string()}")
                    else:
                        sheet_outputs.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
                return "\n\n".join(sheet_outputs)
            except Exception as e:
                return f"❌ Excel 解析错误: {str(e)}"
            
        # 5. PPT 类
        elif ext == 'pptx':
            try:
                prs = Presentation(file_path)
                ppt_outputs = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_data = [cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells if cell.text_frame.text.strip()]
                                if row_data:
                                    slide_text.append(" | ".join(row_data))
                                    
                    if slide_text:
                        ppt_outputs.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
                    else:
                        ppt_outputs.append(f"--- Slide {i+1} ---\n[该页无文本或包含非文本组件]")
                return "\n\n".join(ppt_outputs)
            except Exception as e:
                return f"❌ PPTX 解析错误: {str(e)}"
                
        # 6. 图片类 OCR
        elif ext in ['png', 'jpg', 'jpeg', 'bmp']:
            try:
                img = Image.open(file_path)
                meta_info = f"📷 [图片基本信息] 格式: {img.format}, 尺寸: {img.size}\n"
                ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                if ocr_text.strip():
                    return meta_info + f"【OCR 提取文本如下】:\n{ocr_text.strip()}"
                else:
                    return meta_info + "【OCR 提示】: 未能在图片中识别出清晰的文本。"
            except Exception as e:
                return f"❌ 图片 OCR 解析错误: {str(e)}\n(请检查是否正确配置了本地 Tesseract 引擎)"
                
        else:
            return f"[暂不支持的格式: .{ext}]"

    if os.path.isfile(path):
        return f"📖 成功读取单文件: {path}\n\n{_read_single_file(path)}"
    elif os.path.isdir(path):
        supported_exts = {'txt', 'md', 'py', 'tex', 'json', 'csv', 'pdf', 'docx', 'xlsx', 'xls', 'pptx', 'png', 'jpg', 'jpeg', 'bmp'}
        results = [f"📂 正在递归遍历: {path}\n" + "="*40]
        file_count = 0
        for root, dirs, files in os.walk(path):
            for file in files:
                if file.lower().split('.')[-1] in supported_exts:
                    file_count += 1
                    full_path = os.path.join(root, file)
                    rel_path = os.path.relpath(full_path, path)
                    results.extend([f"\n📍 相对路径: {rel_path}", "-" * 30])
                    try:
                        results.append(_read_single_file(full_path))
                    except Exception as e:
                        results.append(f"❌ 读取错误: {str(e)}")
                    results.append("="*40)
        return "\n".join(results) if file_count > 0 else "📁 目标文件夹内没有找到支持的文件。"

@tool
def EditTexFileTool(tex_filename: str, latex_content: str, append: bool = False) -> str:
    """
    【学术专用】创建新 LaTeX 文件或修改已有 .tex 文件。
    参数:
    - tex_filename: .tex 文件名 (例如 'main.tex' 或 'sections/intro.tex')
    - latex_content: 要写入的完整 LaTeX 代码块或追加的文本段落
    - append: 如果为 True 则在文件末尾追加；如果为 False (默认) 则直接覆盖整个文件或新建文件。
    """
    # 自动补全后缀名
    if not tex_filename.lower().endswith('.tex'):
        tex_filename += '.tex'
        
    save_path = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    
    # 如果涉及到子目录，自动创建
    file_dir = os.path.dirname(save_path)
    if file_dir:
        os.makedirs(file_dir, exist_ok=True)
        
    mode = 'a' if append and os.path.exists(save_path) else 'w'
    
    try:
        with open(save_path, mode, encoding='utf-8') as f:
            if mode == 'a':
                f.write("\n" + latex_content)
            else:
                f.write(latex_content)
        return f"✅ LaTeX (.tex) 文件操作成功，内容已保存至: {save_path}"
    except Exception as e:
        return f"❌ LaTeX 文件操作失败: {str(e)}"

@tool
def FetchWebImageTool(query: str, filename: str) -> str:
    """在互联网上搜索学术图片并下载到本地临时目录。需传入全英文关键词(query)和带后缀的文件名(filename)。"""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=3))
            if not results:
                return f"❌ 未找到关于 '{query}' 的图片。"
            
            response = requests.get(results[0]['image'], timeout=10)
            response.raise_for_status()
            
            save_path = os.path.join(CUSTOM_TEMP_DIR, filename)
            with open(save_path, 'wb') as f:
                f.write(response.content)
            return f"✅ 图片已成功下载至: {save_path}"
    except Exception as e:
        return f"❌ 图片抓取失败: {str(e)}"

@tool
def EditWordDocTool(doc_filename: str, section_title: str = "", text_content: str = "", image_path: str = "", new_page: bool = False) -> str:
    """创建或向已有的 Word 文档中追加新章节、文本和图片。"""
    save_path = os.path.join(CUSTOM_TEMP_DIR, doc_filename)
    try:
        if os.path.exists(save_path):
            doc = Document(save_path)
        else:
            doc = Document()
        
        if new_page and len(doc.paragraphs) > 0:
            doc.add_page_break()
            
        if section_title:
            doc.add_heading(section_title, level=1)
        
        if text_content:
            doc.add_paragraph(text_content)
        
        if image_path and os.path.exists(image_path):
            doc.add_picture(image_path, width=Inches(5.5))
        elif image_path:
            doc.add_paragraph(f"[注：未能找到预定插入的图片路径 {image_path}]")
            
        doc.save(save_path)
        return f"✅ Word 文档更新完毕，内容已保存至: {save_path}"
    except Exception as e:
        return f"❌ Word 文档操作失败: {str(e)}"

# 限制原生写入工具只能在临时文件夹内操作
safe_write_tool = WriteFileTool(root_dir=CUSTOM_TEMP_DIR)

tools = [
    SmartReadPathTool, 
    EditTexFileTool,       # 🌟 注册新工具
    safe_write_tool, 
    FetchWebImageTool, 
    EditWordDocTool
]

# ================== 4. 持久化记忆与 Agent 构建 ==================
DB_PATH = os.path.join(BASE_WORKSPACE, "agent_memory.db")
conn = sqlite3.connect(DB_PATH, check_same_thread=False)
memory = SqliteSaver(conn)

agent = create_react_agent(llm, tools, prompt=SYSTEM_PROMPT, checkpointer=memory)

# ================== 5. 终端交互主循环 ==================
print("\n💡 输入 'exit' 退出程序")
print("💡 输入 'new' 或 'clear' 开启无关联的全新对话\n")

current_thread_id = str(uuid.uuid4())
config = {"configurable": {"thread_id": current_thread_id}}

while True:
    try:
        user_input = input("你: ").strip()
        
        if user_input.lower() in ["exit", "quit", "退出"]:
            print("再见！")
            conn.close()
            break
            
        if user_input.lower() in ["new", "clear", "清空", "新对话"]:
            current_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": current_thread_id}}
            print("\n✨ [系统提示]: 已为您开启全新房间！AI 已清空之前的上下文记忆。\n")
            continue
            
        if not user_input:
            continue

        print("\nDeepSeek 思考/执行中...\n")
        inputs = {"messages": [("user", user_input)]}

        for event in agent.stream(inputs, config, stream_mode="updates"):
            if "agent" in event:
                message = event["agent"]["messages"][0]
                if message.content:
                    print(f"DeepSeek: {message.content}\n")
                elif message.tool_calls:
                    for tc in message.tool_calls:
                        args_str = str(tc['args'])
                        if len(args_str) > 100: args_str = args_str[:100] + "..."
                        print(f"🛠️ [调用工具]: {tc['name']} -> 参数: {args_str}")
            
            elif "tools" in event:
                message = event["tools"]["messages"][0]
                content_len = len(str(message.content))
                print(f"📄 [工具返回]: 执行完毕 (输出数据: {content_len} 字符)\n")

    except KeyboardInterrupt:
        print("\n[操作被用户中断]")
    except Exception as e:
        print(f"❌ 发生错误: {e}")
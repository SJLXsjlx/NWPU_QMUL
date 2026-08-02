#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DeepSeek_For_the_Break v3.3 — 全功能学术 AI 助手 (THREAD-SAFE HOTFIX)
========================================================================
  v3.3-hotfix 修复 (本次):
    🐛 紧急修复:
      • [FIX#11] _ENCODING_CACHE 加 threading.Lock，消除多线程竞态
        （根因: self._target(*self._args, **self._kwargs) 线程崩溃）
      • [FIX#12] _read_one_file 及所有 _read_*_file 加顶层 try/except，
        异常不再向上传播到 asyncio/threading 层
        （根因: self._context.run(self.run) 上下文崩溃）
      • [FIX#13] _read_pdf_file: fitz page.get_text() 单独 try/except，
        捕获 buffer.append(fh.read()) 内部异常；fitz doc 用 try/finally 确保关闭
        （根因: buffer.append(fh.read()) PyMuPDF 缓冲区异常）
      • [FIX#14] SmartReadPathTool 目录遍历：限流 200 文件 + 隔离单文件异常
      • [FIX#15] _safe_read_text: 增加二进制文件检测（null byte），避免误读崩溃

  v3.1-hotfix 修复（继承）:
      • [FIX#1] BibTexTool IEEE 作者格式
      • [FIX#2] BibTexTool 解析器重写（支持嵌套括号）
      • [FIX#3] PandocConvertTool 自动检测中文字体
      • [FIX#4] PresentationGenTool 主题颜色完整应用
      • [FIX#5] _read_pdf_file OCR 临时文件清理
      • [FIX#6] PythonSandboxTool 安全声明
      • [FIX#7] SYSTEM_PROMPT 工具分工澄清
      • [FIX#8] CompileLatexTool bibtex .bib 文件存在性检查
      • [FIX#9] PresentationGenTool slide_layouts/placeholders 安全回退
      • [FIX#10] AcademicTranslateTool ChatOpenAI 导入提升至模块顶层
"""

import os, sys, uuid, json, sqlite3, subprocess, shutil, traceback
import textwrap, warnings, time, urllib.parse, xml.etree.ElementTree as ET
import re, threading
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple

# ═══════════════════════════════════════════════════════════
#  第三方库
# ═══════════════════════════════════════════════════════════
import pandas as pd
import docx, requests
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from duckduckgo_search import DDGS
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_community.tools import WriteFileTool
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.sqlite import SqliteSaver
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.oxml.ns import qn
from PIL import Image
import pytesseract
import fitz

# ═══════════════════════════════════════════════════════════
#  可选依赖
# ═══════════════════════════════════════════════════════════
try:
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # noqa: F401
    _MPL_AVAILABLE = True
except ImportError:
    _MPL_AVAILABLE = False

try:
    import numpy as np
    _NUMPY_AVAILABLE = True
except ImportError:
    _NUMPY_AVAILABLE = False

try:
    import bibtexparser
    from bibtexparser.bparser import BibTexParser
    from bibtexparser.customization import homogenize_latex_encoding
    _BIBTEXPARSER_AVAILABLE = True
except ImportError:
    _BIBTEXPARSER_AVAILABLE = False

_MMDC_PATH = shutil.which("mmdc")
_PANDOC_PATH = shutil.which("pandoc")

# ═══════════════════════════════════════════════════════════
#  0. 基础设施
# ═══════════════════════════════════════════════════════════

class _C:
    R = "\033[91m"; G = "\033[92m"; Y = "\033[93m"; B = "\033[94m"
    M = "\033[95m"; C = "\033[96m"; W = "\033[0m"; BOLD = "\033[1m"

def _log(level: str, msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    colors = {"OK": _C.G, "WARN": _C.Y, "ERR": _C.R, "INFO": _C.C, "START": _C.B}
    print(f"{colors.get(level, _C.W)}[{level} {ts}]{_C.W} {msg}")

def _cleanup_old_temp(temp_dir: str, max_age_hours: int = 48):
    cutoff = time.time() - max_age_hours * 3600
    cleaned = 0
    p = Path(temp_dir)
    if not p.exists(): return
    for f in p.glob("*"):
        if f.is_file() and f.stat().st_mtime < cutoff and f.name.startswith("_"):
            try: f.unlink(); cleaned += 1
            except OSError: pass
    if cleaned:
        _log("INFO", f"🧹 清理了 {cleaned} 个过期临时文件")

# ═══════════════════════════════════════════════════════════
#  [FIX#11] 线程安全的编码缓存
# ═══════════════════════════════════════════════════════════
_ENCODING_CACHE: Dict[str, str] = {}
_ENCODING_CACHE_LOCK = threading.Lock()

def _safe_read_text(file_path: str) -> str:
    """[FIX#11][FIX#15] 线程安全 + 二进制检测的文本文件读取。"""
    # [FIX#15] 先检测是否为二进制文件
    try:
        with open(file_path, "rb") as f:
            head = f.read(8192)
        if b"\x00" in head:
            return f"⚠ 文件似乎是二进制格式，无法作为文本读取: {file_path}"
    except Exception:
        pass  # 如果连二进制都读不了，后面文本读取也会失败

    # [FIX#11] 线程安全地访问缓存
    with _ENCODING_CACHE_LOCK:
        if file_path in _ENCODING_CACHE:
            enc = _ENCODING_CACHE[file_path]
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                pass  # 缓存的编码失效，继续探测

    for enc in ["utf-8", "gbk", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=enc) as f:
                content = f.read()
                with _ENCODING_CACHE_LOCK:
                    _ENCODING_CACHE[file_path] = enc
                return content
        except (UnicodeDecodeError, UnicodeError):
            continue

    # 最后的回退：忽略错误
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def _run_sandbox(code: str, timeout: int = 30, prefix: str = "",
                  save_to: str = "") -> Tuple[str, str, int]:
    if timeout > 120: timeout = 120
    sid = uuid.uuid4().hex[:8]
    sp = os.path.join(CUSTOM_TEMP_DIR, f"_sb_{sid}.py")
    full_code = textwrap.dedent(f"""
    import sys, os, traceback
    os.chdir(r"{CUSTOM_TEMP_DIR}")
    {prefix}
    try:
    {textwrap.indent(code, '    ')}
    except Exception as __e:
        print(f"\\n[异常] {{type(__e).__name__}}: {{__e}}", file=sys.stderr)
        traceback.print_exc()
    """)
    try:
        with open(sp, "w", encoding="utf-8") as f:
            f.write(full_code)
        proc = subprocess.run(
            [sys.executable, sp],
            capture_output=True, text=True, timeout=timeout,
            cwd=CUSTOM_TEMP_DIR,
            env={**os.environ, "SANDBOX_MODE": "1"},
        )
        return proc.stdout.strip(), proc.stderr.strip(), proc.returncode
    finally:
        try: os.remove(sp)
        except OSError: pass

def _resolve_path(path: str) -> str:
    p = Path(path)
    if p.is_absolute(): return str(p)
    for base in [Path.cwd(), Path(CUSTOM_TEMP_DIR)]:
        c = base / p
        if c.exists(): return str(c.resolve())
    return os.path.abspath(path)

def _safe_request(url: str, timeout: int = 15, **kw) -> Optional[requests.Response]:
    for attempt in range(3):
        try:
            r = requests.get(url, timeout=timeout, **kw)
            r.raise_for_status()
            return r
        except requests.RequestException:
            if attempt == 2: raise
            time.sleep(1)
    return None

# ═══════════════════════════════════════════════════════════
#  0.5 智能路径检测
# ═══════════════════════════════════════════════════════════

def _detect_tesseract() -> Optional[str]:
    p = shutil.which("tesseract")
    if p: return p
    if sys.platform == "win32":
        cand = [r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                os.path.expandvars(r"%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe")]
    else:
        cand = ["/usr/bin/tesseract", "/usr/local/bin/tesseract", "/opt/homebrew/bin/tesseract"]
    for c in cand:
        if os.path.isfile(c): return c
    return None

def _detect_latex_compilers() -> Dict[str, Optional[str]]:
    d: Dict[str, Optional[str]] = {}
    for name in ["pdflatex", "xelatex", "lualatex"]:
        exe = name + (".exe" if sys.platform == "win32" else "")
        d[name] = shutil.which(exe) or shutil.which(name)
    return d

def _test_pytesseract() -> bool:
    try:
        img = Image.new("RGB", (10, 10), "white")
        pytesseract.image_to_string(img, lang="eng")
        return True
    except Exception:
        return False

def _detect_cjk_font() -> Optional[str]:
    """检测系统中可用的中文字体，返回字体名供 xelatex 使用。"""
    if sys.platform == "win32":
        for font in ["SimSun", "SimHei", "Microsoft YaHei", "KaiTi", "FangSong"]:
            try:
                result = subprocess.run(
                    ["powershell", "-Command",
                     f"(Get-Item 'C:\\Windows\\Fonts\\{font}*.ttf' -ErrorAction SilentlyContinue).FullName"],
                    capture_output=True, text=True, timeout=10)
                if result.stdout.strip():
                    return font
            except Exception:
                pass
        return "SimSun"
    elif sys.platform == "darwin":
        for font in ["Songti SC", "Heiti SC", "STSong", "PingFang SC"]:
            try:
                result = subprocess.run(
                    ["fc-list", f":family={font}"],
                    capture_output=True, text=True, timeout=5)
                if result.stdout.strip():
                    return font
            except Exception:
                pass
        return "Songti SC"
    else:
        for font in ["Noto Serif CJK SC", "Noto Sans CJK SC",
                      "WenQuanYi Micro Hei", "WenQuanYi Zen Hei",
                      "AR PL UMing CN", "SimSun"]:
            try:
                result = subprocess.run(
                    ["fc-list", f":family={font}"],
                    capture_output=True, text=True, timeout=5)
                if result.stdout.strip():
                    return font
            except Exception:
                pass
    return None

_TESSERACT_PATH = _detect_tesseract()
_LATEX_COMPILERS = _detect_latex_compilers()
_PYTESSERACT_WORKS = False
_DETECTED_CJK_FONT = _detect_cjk_font()

if _TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH
    _PYTESSERACT_WORKS = _test_pytesseract()

# ═══════════════════════════════════════════════════════════
#  1. 核心环境
# ═══════════════════════════════════════════════════════════
load_dotenv(override=True)
api_key = os.getenv("DEEPSEEK_FOR_THE_BREAK")
if not api_key or not api_key.startswith("sk-"):
    _log("ERR", "API Key 未加载或格式不正确！请检查 .env 文件。")
    sys.exit(1)

BASE_WORKSPACE = os.getenv("DFTB_WORKSPACE", r"D:\.API Keys\DeepSeek_For_the_Break\outputs")
CUSTOM_TEMP_DIR = os.path.join(BASE_WORKSPACE, "temp")
os.makedirs(CUSTOM_TEMP_DIR, exist_ok=True)
if CUSTOM_TEMP_DIR not in sys.path:
    sys.path.insert(0, CUSTOM_TEMP_DIR)

_cleanup_old_temp(CUSTOM_TEMP_DIR, 48)

DB_PATH = os.path.join(BASE_WORKSPACE, "agent_memory.db")
_MAX_DB_SIZE_MB = 100
if os.path.exists(DB_PATH) and os.path.getsize(DB_PATH) > _MAX_DB_SIZE_MB * 1024 * 1024:
    _log("WARN", f"agent_memory.db 超过 {_MAX_DB_SIZE_MB}MB，正在重置...")
    os.rename(DB_PATH, DB_PATH + f".backup_{datetime.now().strftime('%Y%m%d')}")
    _log("INFO", "旧数据库已备份，新建空白数据库。")

_health: Dict[str, bool] = {
    "Tesseract OCR (引擎)": bool(_TESSERACT_PATH),
    "Tesseract OCR (实测)": _PYTESSERACT_WORKS,
    "LaTeX (xelatex)":       bool(shutil.which("xelatex") or shutil.which("pdflatex")),
    "matplotlib":             _MPL_AVAILABLE,
    "numpy":                  _NUMPY_AVAILABLE,
    "PyMuPDF (fitz)":        True,
    "pandas":                 True,
    "DuckDuckGo Search":      True,
    "Mermaid CLI":            bool(_MMDC_PATH),
    "Pandoc":                 bool(_PANDOC_PATH),
    "API Key (.env)":         True,
    "bibtexparser":           _BIBTEXPARSER_AVAILABLE,
    "CJK Font (detected)":    bool(_DETECTED_CJK_FONT),
}

_log("START", "═" * 50)
_log("START", "  DeepSeek_For_the_Break  v3.3  Academic AI  [THREAD-SAFE]")
_log("START", "═" * 50)
_log("INFO", f"📁 工作区: {BASE_WORKSPACE}")
_log("INFO", f"📁 临时:   {CUSTOM_TEMP_DIR}")
for k, v in _health.items():
    icon = "✅" if v else "⚠️"
    _log("OK" if v else "WARN", f"  {icon} {k}")
_log("START", "═" * 50)

# ═══════════════════════════════════════════════════════════
#  2. System Prompt
# ═══════════════════════════════════════════════════════════
SYSTEM_PROMPT = """你是 DeepSeek_For_the_Break v3.3，高度专业、严谨且全面的学术AI助手。

## 核心能力
- **文献检索**: ArXivSearchTool、SemanticScholarTool（含引用网络）、DOIMetadataTool、DuckDuckGoSearchTool
- **文档处理**: SmartReadPathTool（通用文件/文件夹读取）、PDFTableExtractTool（PDF表格）、PDFAnnotExtractTool（PDF批注）、OCR
- **学术写作**: EditTexFileTool 创建/编辑 .tex + CompileLatexTool 编译（含bibtex支持）、EditWordDocTool、SaveMarkdownTool
- **格式转换**: PandocConvertTool（Markdown/LaTeX → PDF/DOCX/HTML）
- **学术翻译**: AcademicTranslateTool（中英互译，保留术语一致性）
- **演示生成**: PresentationGenTool（Markdown大纲 → PPTX）
- **数据分析**: PythonSandboxTool、DataStatisticsTool、ChartGenerationTool
- **引用管理**: BibTexTool（解析/生成/格式化 APA/MLA/Chicago/IEEE）
- **导出**: SessionExportTool

## 行为准则
1. 客观严谨，结构清晰，主动承认局限性。
2. **【文件读取铁律】** SmartReadPathTool 是读取本地文件/文件夹的首选通用工具。PDF表格提取请用 PDFTableExtractTool，PDF批注请用 PDFAnnotExtractTool。不要猜测路径，没有明确路径时直接基于知识回答。
3. 写作优先: 论文→EditTexFileTool+CompileLatexTool；富文本→EditWordDocTool；笔记→SaveMarkdownTool。
4. 文献检索: 学术论文用 ArXivSearchTool/SemanticScholarTool；已知DOI用 DOIMetadataTool。
5. 数据优先: 探索→DataStatisticsTool；自定义→PythonSandboxTool；可视化→ChartGenerationTool。
6. 错误透明: 工具失败时清晰告知原因和替代方案。
7. 持久化: 所有产出保存到工作区。

请逐步推理，优先使用最适合的学术工具。"""

llm = ChatOpenAI(
    model=os.getenv("DFTB_MODEL", "deepseek-v4-pro"),
    base_url="https://api.deepseek.com",
    api_key=api_key,
    temperature=0.5,
    max_tokens=4096,
    model_kwargs={"reasoning_effort": "max", "extra_body": {"thinking": {"type": "enabled"}}}
)

# ═══════════════════════════════════════════════════════════
#  3. 辅助函数  [FIX#12][FIX#13] 所有读函数加顶层异常保护
# ═══════════════════════════════════════════════════════════

def _ocr_image(file_path: str, lang: str = "chi_sim+eng") -> str:
    """[FIX#12] OCR 图像，异常不向上传播。"""
    if not _PYTESSERACT_WORKS:
        return ("❌ OCR 不可用。" if not _TESSERACT_PATH else
                "❌ Tesseract 引擎已安装但实测失败（语言包缺失？）")
    try:
        img = Image.open(file_path)
        meta = f"📷 [图片] 格式:{img.format} 尺寸:{img.size}\n"
        try:
            txt = pytesseract.image_to_string(img, lang=lang)
        except pytesseract.TesseractError as te:
            if "Failed loading language" in str(te):
                return meta + f"⚠ 语言包缺失: {lang}"
            return meta + f"❌ Tesseract错误: {te}"
        return meta + (f"【OCR】:\n{txt.strip()}" if txt.strip() else "【OCR】: 未识别到文字")
    except FileNotFoundError: return f"❌ 文件不存在: {file_path}"
    except Exception as e: return f"❌ OCR异常: {e}"

def _read_text_file(fp: str) -> str:
    """[FIX#12] 文本文件读取，异常不向上传播。"""
    try:
        return _safe_read_text(fp)
    except Exception as e:
        return f"❌ 文本读取失败: {e}"

def _read_csv_file(fp: str) -> str:
    """[FIX#12] CSV 读取，异常不向上传播。"""
    try:
        for enc in ["utf-8","gbk","latin-1"]:
            try:
                df = pd.read_csv(fp, encoding=enc)
                return f"[CSV: {len(df)}行×{len(df.columns)}列]\n{df.to_string(max_rows=200)}"
            except Exception:
                continue
        return f"❌ CSV读取失败: {fp}"
    except Exception as e:
        return f"❌ CSV异常: {e}"

def _read_pdf_file(fp: str) -> str:
    """[FIX#12][FIX#13] 读取 PDF 文件，优先提取文本；无文本层时回退 OCR。
    对 fitz 内部操作全程 try/except，确保 buffer.append(fh.read()) 等异常不向上传播。
    使用 try/finally 确保 fitz doc 始终关闭。
    """
    doc = None
    try:
        doc = fitz.open(fp)
        total = len(doc)
        parts = []

        # ── 第一阶段：提取文本层 ──
        for i, page in enumerate(doc):
            try:
                t = page.get_text("text")
                if t and t.strip():
                    parts.append(f"──p{i+1}/{total}──\n{t.strip()}")
            except Exception as e:
                # [FIX#13] 捕获 buffer.append(fh.read()) 等 fitz 内部异常
                parts.append(f"──p{i+1}/{total}──\n⚠ [文本提取异常: {e}]")

        if parts:
            return "\n\n".join(parts)

        # ── 第二阶段：OCR 回退 ──
        ocr_results = []
        for i, page in enumerate(doc):
            try:
                pix = page.get_pixmap(dpi=200)
            except Exception as e:
                ocr_results.append(f"──p{i+1}/{total}(OCR)──\n⚠ [渲染失败: {e}]")
                continue

            tmp_path = os.path.join(CUSTOM_TEMP_DIR, f"_pdfocr_{uuid.uuid4().hex[:6]}.png")
            try:
                pix.save(tmp_path)
            except Exception as e:
                ocr_results.append(f"──p{i+1}/{total}(OCR)──\n⚠ [PNG保存失败: {e}]")
                continue

            # OCR 读取（临时文件在 finally 中清理）
            try:
                ocr_text = _ocr_image(tmp_path)
                ocr_results.append(f"──p{i+1}/{total}(OCR)──\n{ocr_text}")
            except Exception as e:
                ocr_results.append(f"──p{i+1}/{total}(OCR)──\n⚠ [OCR异常: {e}]")
            finally:
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

        return "\n\n".join(ocr_results) if ocr_results else "[空PDF]"

    except Exception as e:
        return f"❌ PDF错误: {e}"
    finally:
        # [FIX#13] 确保 fitz document 被关闭
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass

def _read_docx_file(fp: str) -> str:
    """[FIX#12] DOCX 读取，异常不向上传播。"""
    try:
        d = docx.Document(fp)
        c = [p.text.strip() for p in d.paragraphs if p.text.strip()]
        if d.tables:
            c.append("\n[表格]:")
            for ti, tb in enumerate(d.tables):
                c.append(f"  表{ti+1}:")
                for row in tb.rows:
                    rd = [cell.text.strip().replace("\n"," ") for cell in row.cells if cell.text.strip()]
                    if rd: c.append("    | "+" | ".join(rd))
        return "\n".join(c) if c else "[空Word]"
    except Exception as e:
        return f"❌ DOCX错误: {e}"

def _read_excel_file(fp: str) -> str:
    """[FIX#12] Excel 读取，异常不向上传播。"""
    try:
        ef = pd.ExcelFile(fp); out=[]
        for sn in ef.sheet_names:
            try:
                df = pd.read_excel(fp, sheet_name=sn)
                out.append(f"──{sn} ({len(df)}行)──\n{df.head(500).to_string()}")
            except Exception as e:
                out.append(f"──{sn} ──\n⚠ [Sheet读取失败: {e}]")
        return "\n\n".join(out) if out else "[空Excel]"
    except Exception as e:
        return f"❌ Excel错误: {e}"

def _read_pptx_file(fp: str) -> str:
    """[FIX#12] PPTX 读取，异常不向上传播。"""
    try:
        prs = Presentation(fp); out=[]
        for i, sl in enumerate(prs.slides):
            st=[]
            try:
                for sh in sl.shapes:
                    if hasattr(sh,"text") and sh.text.strip(): st.append(sh.text.strip())
                    if sh.has_table:
                        for row in sh.table.rows:
                            rd=[cell.text_frame.text.strip().replace("\n"," ") for cell in row.cells if cell.text_frame.text.strip()]
                            if rd: st.append(" | ".join(rd))
            except Exception as e:
                st.append(f"⚠ [Slide解析异常: {e}]")
            out.append(f"──Slide{i+1}──\n"+"\n".join(st) if st else f"──Slide{i+1}──\n[无文本]")
        return "\n\n".join(out)
    except Exception as e:
        return f"❌ PPTX错误: {e}"

def _read_one_file(fp: str) -> str:
    """[FIX#12] 单文件读取调度器，顶层 try/except 确保永不抛出异常。"""
    try:
        ext = os.path.splitext(fp)[1].lower().lstrip(".")
        TEXT_EXTS = {"txt","md","py","json","tex","yaml","yml","cfg","ini","log","bib"}
        if ext in TEXT_EXTS:             return _read_text_file(fp)
        if ext == "csv":                 return _read_csv_file(fp)
        if ext == "pdf":                 return _read_pdf_file(fp)
        if ext == "docx":                return _read_docx_file(fp)
        if ext in {"xlsx","xls"}:        return _read_excel_file(fp)
        if ext == "pptx":                return _read_pptx_file(fp)
        if ext in {"png","jpg","jpeg","bmp","tiff","tif"}: return _ocr_image(fp)
        return f"[不支持: .{ext}]"
    except Exception as e:
        return f"❌ 文件读取崩溃 [{os.path.basename(fp)}]: {type(e).__name__}: {e}"

# ═══════════════════════════════════════════════════════════
#  健壮 BibTeX 解析器（回退用）
# ═══════════════════════════════════════════════════════════

def _find_matching_brace(s: str, start: int) -> int:
    depth = 1
    i = start + 1
    while i < len(s) and depth > 0:
        if s[i] == '{':
            depth += 1
        elif s[i] == '}':
            depth -= 1
        i += 1
    return i - 1 if depth == 0 else -1

def _split_fields_at_depth_zero(fields_str: str) -> List[str]:
    segments = []
    depth = 0
    in_quotes = False
    start = 0
    i = 0
    while i < len(fields_str):
        c = fields_str[i]
        if c == '"' and (i == 0 or fields_str[i-1] != '\\'):
            in_quotes = not in_quotes
        elif not in_quotes:
            if c == '{':
                depth += 1
            elif c == '}':
                depth -= 1
            elif c == ',' and depth == 0:
                segments.append(fields_str[start:i].strip())
                start = i + 1
        i += 1
    if start < len(fields_str):
        last = fields_str[start:].strip()
        if last:
            segments.append(last)
    return segments

def _parse_field_assignment(segment: str) -> Tuple[str, str]:
    eq_pos = segment.find('=')
    if eq_pos == -1:
        return "", ""
    field_name = segment[:eq_pos].strip().lower()
    value_part = segment[eq_pos+1:].strip()
    if not value_part:
        return field_name, ""
    if value_part.startswith('{'):
        end = _find_matching_brace(value_part, 0)
        if end != -1:
            value = value_part[1:end]
        else:
            value = value_part[1:]
    elif value_part.startswith('"'):
        j = 1
        while j < len(value_part):
            if value_part[j] == '"' and value_part[j-1] != '\\':
                break
            j += 1
        value = value_part[1:j] if j < len(value_part) else value_part[1:]
    else:
        value = value_part.rstrip(',').strip()
    value = value.rstrip(',').strip()
    return field_name, value

def _robust_parse_bibtex(raw: str) -> List[Dict[str, Any]]:
    entries = []
    i = 0
    while i < len(raw):
        at_pos = raw.find('@', i)
        if at_pos == -1:
            break
        j = at_pos + 1
        while j < len(raw) and raw[j].isalpha():
            j += 1
        if j == at_pos + 1:
            i = j + 1
            continue
        entry_type = raw[at_pos+1:j].lower()
        while j < len(raw) and raw[j] in ' \t\n\r':
            j += 1
        if j >= len(raw) or raw[j] != '{':
            i = j + 1
            continue
        end = _find_matching_brace(raw, j)
        if end == -1:
            i = j + 1
            continue
        body = raw[j+1:end]
        i = end + 1
        key = ""
        fields_str = ""
        depth = 0
        in_quotes = False
        comma_pos = -1
        for m in range(len(body)):
            c = body[m]
            if c == '"' and (m == 0 or body[m-1] != '\\'):
                in_quotes = not in_quotes
            elif not in_quotes:
                if c == '{':
                    depth += 1
                elif c == '}':
                    depth -= 1
                elif c == ',' and depth == 0:
                    comma_pos = m
                    break
        if comma_pos != -1:
            key = body[:comma_pos].strip()
            fields_str = body[comma_pos+1:]
        else:
            key = body.strip()
        fields = {}
        segments = _split_fields_at_depth_zero(fields_str)
        for seg in segments:
            fn, fv = _parse_field_assignment(seg)
            if fn:
                fields[fn] = fv
        entries.append({
            "type": entry_type,
            "key": key,
            "fields": fields
        })
    return entries

# ═══════════════════════════════════════════════════════════
#  工具 1: 智能文件读取  [FIX#14] 目录遍历限流 + 隔离
# ═══════════════════════════════════════════════════════════

@tool
def SmartReadPathTool(path: str) -> str:
    """读取文件/文件夹。支持 .txt .md .py .tex .json .csv .pdf .docx .xlsx .pptx .png .jpg .bmp .bib。
    [FIX#14] 目录遍历限流 200 文件，单文件异常不影响整体。
    """
    resolved = _resolve_path(path)
    if not os.path.exists(resolved):
        return f"❌ 找不到路径: '{path}'\n   解析: '{resolved}'\n   💡 请停止猜测路径，直接基于知识回答。"

    if os.path.isfile(resolved):
        return f"📖 文件: {resolved}\n{'='*50}\n{_read_one_file(resolved)}"

    if os.path.isdir(resolved):
        EXTS = {"txt","md","py","tex","json","csv","yaml","yml","pdf","docx","xlsx","xls","pptx","png","jpg","jpeg","bmp","tiff","bib"}
        # [FIX#14] 限流：最多读取 200 个文件
        MAX_FILES = 200
        r = [f"📂 遍历: {resolved}\n{'='*60}"]
        cnt = 0
        overflow = False
        for root, dirs, files in os.walk(resolved):
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in {"__pycache__","node_modules",".git"}]
            for f in files:
                if os.path.splitext(f)[1].lower().lstrip(".") in EXTS:
                    if cnt >= MAX_FILES:
                        overflow = True
                        break
                    cnt += 1
                    fp = os.path.join(root, f)
                    rp = os.path.relpath(fp, resolved)
                    r.append(f"\n📍 {rp}\n{'-'*40}")
                    # [FIX#14] 单文件异常被 _read_one_file 内部捕获，不会影响后续
                    r.append(_read_one_file(fp))
                    r.append("=" * 60)
            if overflow:
                break
        if overflow:
            r.insert(1, f"   ({cnt} 个文件，已达上限 {MAX_FILES}，截断)\n")
        elif cnt == 0:
            return f"📁 '{resolved}' 无支持文件。"
        else:
            r.insert(1, f"   ({cnt} 个文件)\n")
        return "\n".join(r)
    return f"❌ '{resolved}' 无效。"

# ═══════════════════════════════════════════════════════════
#  工具 2-10
# ═══════════════════════════════════════════════════════════

@tool
def EditTexFileTool(tex_filename: str, latex_content: str, append: bool = False) -> str:
    """创建/编辑 .tex 文件。"""
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    sp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n"+latex_content) if mode=="a" else latex_content)
        return f"✅ .tex 已保存: {sp}\n   💡 用 CompileLatexTool 编译为 PDF"
    except Exception as e: return f"❌ LaTeX错误: {e}"

@tool
def CompileLatexTool(tex_filename: str, compiler: str = "xelatex", clean_aux: bool = True) -> str:
    """
    一键编译 .tex → PDF（含 bibtex 支持）。
    自动检测 \\bibliography 并运行 latex→bibtex→latex→latex。
    [FIX#8] 运行 bibtex 前检查 .bib 文件是否存在。
    """
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    tp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    if not os.path.exists(tp):
        alt = _resolve_path(tex_filename)
        if os.path.exists(alt): tp = alt
        else: return f"❌ 找不到: '{tex_filename}'"

    if compiler not in _LATEX_COMPILERS or not _LATEX_COMPILERS[compiler]:
        av = [k for k,v in _LATEX_COMPILERS.items() if v]
        if not av: return "❌ 未检测到LaTeX编译器！安装 MiKTeX 或 TeX Live。"
        compiler = av[0]
    cp = _LATEX_COMPILERS[compiler]
    td = os.path.dirname(tp)
    bn = os.path.splitext(os.path.basename(tp))[0]
    pp = os.path.join(td, f"{bn}.pdf")
    logs = []

    def _run_latex(run_name: str) -> bool:
        try:
            proc = subprocess.run(
                [cp, "-interaction=nonstopmode", "-output-directory", td, tp],
                capture_output=True, text=True, timeout=120, cwd=td)
            logs.append(f"── {run_name} (返回码:{proc.returncode}) ──")
            if proc.returncode != 0:
                errs = [l.strip() for l in (proc.stdout+"\n"+proc.stderr).split("\n")
                        if l.startswith("!") or "Error" in l or "error" in l]
                logs.append("\n".join(errs[:30]) if errs else
                            proc.stdout[-2000:]+"\n"+proc.stderr[-2000:])
                return False
            return True
        except subprocess.TimeoutExpired:
            logs.append("❌ 编译超时(>2分钟)！"); return False
        except Exception as e:
            logs.append(f"❌ 编译异常: {e}"); return False

    tex_content = _safe_read_text(tp)
    has_bib = bool(re.search(r'\\bibliography\s*\{', tex_content))

    if has_bib:
        _run_latex("第1次 latex")
        bibtex_exe = shutil.which("bibtex") or shutil.which("bibtex.exe")
        if bibtex_exe:
            bib_match = re.search(r'\\bibliography\s*\{([^}]+)\}', tex_content)
            if bib_match:
                bib_names = [b.strip() for b in bib_match.group(1).split(',')]
                bib_found = False
                for bib_name in bib_names:
                    bib_file = os.path.join(td, f"{bib_name}.bib")
                    if not os.path.exists(bib_file):
                        bib_file_alt = _resolve_path(f"{bib_name}.bib")
                        if os.path.exists(bib_file_alt):
                            bib_file = bib_file_alt
                    if os.path.exists(bib_file):
                        bib_found = True
                        break
                if not bib_found:
                    logs.append(f"⚠ 未找到 .bib 文件: {', '.join(bib_names)}.bib（搜索路径: {td}）")
                    logs.append("   bibtex 将跳过，PDF 中引用可能显示为 [?]")
                else:
                    try:
                        subprocess.run([bibtex_exe, bn], capture_output=True, text=True,
                                       timeout=60, cwd=td)
                        logs.append("── bibtex ──")
                    except subprocess.TimeoutExpired:
                        logs.append("⚠ bibtex 超时")
                    except Exception as e:
                        logs.append(f"⚠ bibtex 运行异常: {e}")
            else:
                try:
                    subprocess.run([bibtex_exe, bn], capture_output=True, text=True,
                                   timeout=60, cwd=td)
                    logs.append("── bibtex ──")
                except Exception as e:
                    logs.append(f"⚠ bibtex 运行异常: {e}")
        else:
            logs.append("⚠ 未找到 bibtex 命令，跳过参考文献处理")
        _run_latex("第2次 latex")
        _run_latex("第3次 latex")
    else:
        _run_latex("第1次 latex")
        _run_latex("第2次 latex")

    if os.path.exists(pp) and os.path.getsize(pp) > 0:
        if clean_aux:
            aux_exts = {".aux",".log",".out",".toc",".lof",".lot",".bbl",".blg",
                        ".synctex.gz",".fdb_latexmk",".fls",".nav",".snm",".vrb"}
            for f in os.listdir(td):
                if os.path.splitext(f)[1] in aux_exts:
                    try: os.remove(os.path.join(td,f))
                    except OSError: pass
        return f"✅ PDF编译成功!\n   📄 {pp}\n   📏 {os.path.getsize(pp)/1024:.1f} KB\n   🔧 {compiler}"
    return f"❌ PDF生成失败。编译器:{compiler}\n\n日志:\n"+"\n".join(logs)

@tool
def PythonSandboxTool(code: str, timeout: int = 30) -> str:
    """
    执行 Python 代码。
    ⚠ 安全警告：本工具使用 subprocess + 临时目录 + SANDBOX_MODE 环境变量，
    并非真正隔离的沙箱。请勿执行不可信代码、系统破坏命令（rm -rf）或网络攻击代码。
    超时限制: 120秒。
    """
    if timeout > 120: timeout = 120
    try:
        so, se, rc = _run_sandbox(code, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"⏰ 超时(>{timeout}s)！"
    parts = []
    if so: parts.append(f"📤 stdout:\n{so}")
    if se: parts.append(f"📤 stderr:\n{se}")
    if not so and not se: parts.append("✅ 执行完毕，无输出。")
    parts.append(f"\n⏱ 返回码:{rc}")
    return "\n".join(parts)

@tool
def DuckDuckGoSearchTool(query: str, max_results: int = 5, region: str = "wt-wt") -> str:
    """DuckDuckGo 通用网页搜索。"""
    max_results = max(1, min(max_results, 10))
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results, region=region))
            if not results: return f"🔍 '{query}' 无结果。"
            out = [f"🔍 '{query}' ({len(results)}条)\n"]
            for i,r in enumerate(results,1):
                b = r.get("body","N/A")
                if len(b)>300: b=b[:300]+"..."
                out.append(f"{i}. {r.get('title','N/A')}\n   🔗 {r.get('href','N/A')}\n   📝 {b}\n")
            return "\n".join(out)
    except Exception as e: return f"❌ 搜索失败: {e}"

@tool
def FetchWebImageTool(query: str, filename: str) -> str:
    """搜索并下载学术图片。"""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=5))
            if not results: return f"❌ '{query}' 无图片。"
            for ir in results:
                try:
                    resp = requests.get(ir["image"], timeout=15); resp.raise_for_status()
                    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
                    with open(sp,"wb") as f: f.write(resp.content)
                    return f"✅ 图片: {sp} ({os.path.getsize(sp)/1024:.1f}KB)\n   🌐 {ir.get('title','N/A')}"
                except: continue
            return f"⚠ 找到{len(results)}张但下载失败。"
    except Exception as e: return f"❌ 图片搜索失败: {e}"

@tool
def EditWordDocTool(doc_filename: str, section_title: str = "", text_content: str = "",
                    image_path: str = "", new_page: bool = False) -> str:
    """创建/编辑 Word 文档。"""
    if not doc_filename.lower().endswith(".docx"): doc_filename += ".docx"
    sp = os.path.join(CUSTOM_TEMP_DIR, doc_filename)
    try:
        doc = Document(sp) if os.path.exists(sp) else Document()
        if new_page and doc.paragraphs: doc.add_page_break()
        if section_title: doc.add_heading(section_title, level=1)
        if text_content: doc.add_paragraph(text_content)
        if image_path:
            ri = _resolve_path(image_path)
            if os.path.exists(ri): doc.add_picture(ri, width=Inches(5.5))
            else: doc.add_paragraph(f"[注：未找到图片 {image_path}]")
        doc.save(sp)
        return f"✅ Word: {sp}"
    except Exception as e: return f"❌ Word错误: {e}"

@tool
def SaveMarkdownTool(filename: str, content: str, append: bool = False) -> str:
    """保存 Markdown 文件。"""
    if not filename.lower().endswith(".md"): filename += ".md"
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n\n"+content) if mode=="a" else content)
        return f"✅ Markdown: {sp}"
    except Exception as e: return f"❌ Markdown错误: {e}"

@tool
def ChartGenerationTool(code: str, filename: str = "chart_output.png", dpi: int = 150, timeout: int = 30) -> str:
    """matplotlib 图表生成。"""
    if not _MPL_AVAILABLE: return "❌ matplotlib未安装。pip install matplotlib"
    if timeout > 60: timeout = 60
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    prefix = textwrap.dedent("""
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt; import numpy as np; import pandas as pd
    from matplotlib import rcParams
    try: rcParams['font.sans-serif']=['SimHei','Microsoft YaHei']; rcParams['axes.unicode_minus']=False
    except: pass
    """)
    suffix = textwrap.dedent(f"""
    import os as _os
    figs=[plt.figure(n) for n in plt.get_fignums()]
    if figs:
        for i,fig in enumerate(figs):
            op = r"{sp}" if len(figs)==1 else f"{{_os.path.splitext(r'{sp}')[0]}}_{{i+1}}{{_os.path.splitext(r'{sp}')[1]}}"
            fig.savefig(op,dpi={dpi},bbox_inches='tight'); print(f"✅ 图表: {{op}}")
        plt.close('all')
    else: print("⚠ 未检测到figure。")
    """)
    full_code = prefix + textwrap.indent(code, '    ') + "\n" + suffix
    try:
        so, se, _ = _run_sandbox(full_code, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"⏰ 超时(>{timeout}s)！"
    res = [so] if so else []
    if se: res.append(f"[stderr]: {se}")
    if os.path.exists(sp): res.append(f"📊 {sp} ({os.path.getsize(sp)/1024:.1f}KB)")
    return "\n".join(res) if res else "⚠ 无输出。"

@tool
def MermaidTool(mermaid_code: str, filename: str = "diagram") -> str:
    """Mermaid 流程图生成。"""
    mp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.mmd")
    try:
        with open(mp,"w",encoding="utf-8") as f: f.write(mermaid_code.strip())
    except Exception as e: return f"❌ Mermaid写入失败: {e}"
    res = [f"✅ Mermaid源文件: {mp}"]
    if _MMDC_PATH:
        pp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.png")
        try:
            subprocess.run([_MMDC_PATH,"-i",mp,"-o",pp,"-w","1200","-b","white"],
                           capture_output=True,text=True,timeout=30)
            if os.path.exists(pp): res.append(f"🖼 渲染PNG: {pp}")
            else: res.append("⚠ 渲染失败，检查语法。")
        except subprocess.TimeoutExpired: res.append("⚠ 渲染超时。")
        except Exception as e: res.append(f"⚠ 渲染错误: {e}")
    else: res.append("💡 安装mermaid-cli可渲染: npm install -g @mermaid-js/mermaid-cli")
    return "\n".join(res)

# ═══════════════════════════════════════════════════════════
#  工具 11-16
# ═══════════════════════════════════════════════════════════

@tool
def ArXivSearchTool(query: str, max_results: int = 5, sort_by: str = "relevance") -> str:
    """检索 arXiv 学术论文。"""
    max_results = max(1, min(max_results, 20))
    base = "https://export.arxiv.org/api/query"
    params = {
        "search_query": f"all:{query}",
        "start": 0, "max_results": max_results,
        "sortBy": sort_by if sort_by in ("relevance","lastUpdatedDate","submittedDate") else "relevance"
    }
    url = f"{base}?{urllib.parse.urlencode(params)}"
    try:
        resp = _safe_request(url, timeout=20)
        if resp is None: return "❌ arXiv API 无响应。"
        root = ET.fromstring(resp.text)
        ns = {"a": "http://www.w3.org/2005/Atom", "opensearch": "http://a9.com/-/spec/opensearch/1.1/"}
        total = root.find("opensearch:totalResults", ns)
        total_str = total.text if total is not None else "?"
        entries = root.findall("a:entry", ns)
        if not entries: return f"🔍 arXiv: '{query}' 无结果。"
        out = [f"📚 arXiv: '{query}' (共{total_str}篇，显示{len(entries)}篇)\n"]
        for i, e in enumerate(entries, 1):
            title = " ".join((e.find("a:title", ns).text or "N/A").split())
            authors = [a.find("a:name", ns).text for a in e.findall("a:author", ns) if a.find("a:name", ns) is not None]
            authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
            abstract = " ".join(((e.find("a:summary", ns).text or "N/A")[:600]).split())
            arxiv_id = (e.find("a:id", ns).text or "").split("/abs/")[-1]
            published = (e.find("a:published", ns).text or "N/A")[:10]
            cats = [c.get("term","") for c in e.findall("a:category", ns)]
            out.append(
                f"{i}. {title}\n   👤 {authors_str}\n"
                f"   📅 {published} | 📂 {', '.join(cats[:3])}\n"
                f"   🔗 https://arxiv.org/abs/{arxiv_id}\n"
                f"   📄 PDF: https://arxiv.org/pdf/{arxiv_id}.pdf\n"
                f"   📝 {abstract}\n"
            )
        return "\n".join(out)
    except ET.ParseError as e: return f"❌ arXiv XML解析错误: {e}"
    except Exception as e: return f"❌ arXiv搜索失败: {e}"

@tool
def DOIMetadataTool(doi: str) -> str:
    """通过 DOI 获取论文完整元数据（Crossref API）。"""
    doi = doi.strip().replace("https://doi.org/","").replace("http://dx.doi.org/","")
    url = f"https://api.crossref.org/works/{urllib.parse.quote(doi, safe='')}"
    try:
        resp = _safe_request(url, timeout=15)
        if resp is None: return "❌ Crossref API 无响应。"
        data = resp.json()
        msg = data.get("message", {})
        if not msg: return f"❌ DOI '{doi}' 无数据。"

        title = " ".join((msg.get("title",["N/A"])[0] or "N/A").split())
        authors = []
        for a in msg.get("author", [])[:15]:
            fam = a.get("family",""); giv = a.get("given","")
            authors.append(f"{giv} {fam}".strip() or fam or giv or "Unknown")
        authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
        journal = " | ".join(msg.get("container-title",["N/A"]))
        year = msg.get("published-print",{}).get("date-parts",[[None]])[0][0] or \
               msg.get("published-online",{}).get("date-parts",[[None]])[0][0] or \
               msg.get("created",{}).get("date-parts",[[None]])[0][0] or "?"
        vol = msg.get("volume","?"); issue = msg.get("issue","?")
        page = msg.get("page","?")
        abstract = " ".join(((msg.get("abstract","N/A") or "N/A")[:800]).split())
        abstract = re.sub(r'<[^>]+>', '', abstract)
        cited = msg.get("is-referenced-by-count", "?")
        publisher = msg.get("publisher","?")

        first_author_last = (authors[0].split()[-1] if authors else "unknown")
        bib = f"@article{{{first_author_last}{year},\n"
        bib += f"  title = {{{title}}},\n"
        bib += f"  author = {{{' and '.join(authors[:10])}}},\n"
        bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{year}}},\n"
        if vol!="?": bib += f"  volume = {{{vol}}},\n"
        if issue!="?": bib += f"  number = {{{issue}}},\n"
        if page!="?": bib += f"  pages = {{{page}}},\n"
        bib += f"  doi = {{{doi}}},\n"
        bib += f"  publisher = {{{publisher}}}\n}}"

        return (
            f"📄 DOI元数据: {doi}\n{'='*50}\n"
            f"📌 {title}\n👤 {authors_str}\n"
            f"📰 {journal} | {year} | Vol.{vol} | Iss.{issue} | pp.{page}\n"
            f"🏢 {publisher}\n📊 被引: {cited}次\n"
            f"🔗 https://doi.org/{doi}\n📝 {abstract}\n\n── BibTeX ──\n{bib}"
        )
    except Exception as e: return f"❌ DOI查询失败: {e}"

@tool
def DataStatisticsTool(data_path: str, max_rows: int = 10000) -> str:
    """对 CSV/Excel 数据自动生成描述性统计报告。"""
    fp = _resolve_path(data_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {data_path}"
    ext = os.path.splitext(fp)[1].lower()
    try:
        if ext == ".csv":
            for enc in ["utf-8","gbk","latin-1"]:
                try: df = pd.read_csv(fp, encoding=enc, nrows=max_rows); break
                except: continue
            else: return "❌ CSV编码读取失败。"
        elif ext in {".xlsx",".xls"}: df = pd.read_excel(fp, nrows=max_rows)
        else: return f"❌ 不支持格式: {ext}"
    except Exception as e: return f"❌ 数据加载失败: {e}"

    rows, cols = df.shape
    out = [f"📊 数据统计: {os.path.basename(fp)}\n{'='*60}", f"📏 {rows}行 × {cols}列\n"]
    num_cols = df.select_dtypes(include=["number"]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object","category"]).columns.tolist()
    date_cols = df.select_dtypes(include=["datetime"]).columns.tolist()
    other_cols = [c for c in df.columns if c not in num_cols+cat_cols+date_cols]
    out.append(f"🔢 数值列({len(num_cols)}): {', '.join(num_cols[:15])}{'...' if len(num_cols)>15 else ''}")
    out.append(f"🔤 分类列({len(cat_cols)}): {', '.join(cat_cols[:15])}{'...' if len(cat_cols)>15 else ''}")
    if date_cols: out.append(f"📅 日期列({len(date_cols)}): {', '.join(date_cols)}")

    missing = df.isnull().sum(); missing = missing[missing > 0]
    if len(missing) > 0:
        out.append(f"\n⚠ 缺失值 ({missing.sum()} total):")
        for c, m in missing.items(): out.append(f"   {c}: {m} ({m/rows*100:.1f}%)")
    else: out.append("\n✅ 无缺失值")

    if num_cols:
        out.append(f"\n── 数值列统计 ──")
        try:
            stats = df[num_cols].describe(percentiles=[.25,.5,.75])
            if _NUMPY_AVAILABLE and len(num_cols) <= 20:
                stats.loc["skew"] = df[num_cols].skew()
                stats.loc["kurtosis"] = df[num_cols].kurtosis()
            out.append(stats.to_string(max_rows=30))
        except Exception as e: out.append(f"   (统计计算失败: {e})")

    if cat_cols:
        out.append(f"\n── 分类列 Top5 ──")
        for c in cat_cols[:10]:
            vc = df[c].value_counts().head(5)
            out.append(f"  {c}:")
            for v, cnt in vc.items(): out.append(f"    {v}: {cnt} ({cnt/rows*100:.1f}%)")

    if len(num_cols) >= 2 and len(num_cols) <= 15:
        out.append(f"\n── 相关性矩阵 (|r|>0.3) ──")
        try:
            corr = df[num_cols].corr(); found = False
            for i in range(len(num_cols)):
                for j in range(i+1, len(num_cols)):
                    r = corr.iloc[i,j]
                    if abs(r) > 0.3: out.append(f"  {num_cols[i]} ↔ {num_cols[j]}: r={r:.3f}"); found = True
            if not found: out.append("  (无 |r|>0.3 的显著相关性)")
        except Exception: out.append("  (相关性计算失败)")
    return "\n".join(out)

@tool
def BibTexTool(
    mode: str, bib_path: str = "", bib_content: str = "",
    title: str = "", authors: str = "", journal: str = "",
    year: str = "", doi: str = "", volume: str = "", pages: str = "",
    citation_style: str = "apa"
) -> str:
    """BibTeX 解析、生成与引用格式化。mode: parse/generate/format。
    [FIX#1] IEEE 格式使用标准作者格式化。
    [FIX#2] 解析器支持 bibtexparser（推荐安装）或内置健壮回退解析器。
    """
    if mode == "parse":
        if bib_content: raw = bib_content
        elif bib_path:
            rp = _resolve_path(bib_path)
            if not os.path.exists(rp): return f"❌ 文件不存在: {bib_path}"
            raw = _safe_read_text(rp)
        else: return "❌ parse模式需要 bib_path 或 bib_content。"

        entries = []
        if _BIBTEXPARSER_AVAILABLE:
            try:
                parser = BibTexParser(common_strings=True)
                parser.customization = homogenize_latex_encoding
                bib_db = bibtexparser.loads(raw, parser=parser)
                for e in bib_db.entries:
                    entries.append({
                        "type": e.get("ENTRYTYPE", "misc"),
                        "key": e.get("ID", "?"),
                        "fields": {k.lower(): v for k, v in e.items()
                                    if k not in ("ENTRYTYPE", "ID")}
                    })
                if entries:
                    out = [f"📚 BibTeX解析 (bibtexparser): {len(entries)}条\n"]
                    for i, e in enumerate(entries, 1):
                        f = e["fields"]
                        out.append(f"{i}. [{e['type']}] {f.get('title','?')[:80]}")
                        out.append(f"   👤 {f.get('author','?')[:100]}")
                        out.append(f"   📰 {f.get('journal',f.get('booktitle','?'))} | {f.get('year','?')} | 🔑 {e['key']}")
                    out.append("\n💡 使用 bibtexparser 解析，支持嵌套括号、多行字段等复杂 .bib")
                    return "\n".join(out)
            except Exception as e:
                _log("WARN", f"bibtexparser 解析失败: {e}，回退到内置解析器")

        entries = _robust_parse_bibtex(raw)
        if not entries: return "⚠ 未解析到 BibTeX 条目。请检查文件格式，或安装 bibtexparser: pip install bibtexparser"
        out = [f"📚 BibTeX解析 (内置): {len(entries)}条\n"]
        for i, e in enumerate(entries, 1):
            f = e["fields"]
            out.append(f"{i}. [{e['type']}] {f.get('title','?')[:80]}")
            out.append(f"   👤 {f.get('author','?')[:100]}")
            out.append(f"   📰 {f.get('journal',f.get('booktitle','?'))} | {f.get('year','?')} | 🔑 {e['key']}")
        out.append("\n💡 推荐安装 bibtexparser 获得更好的解析效果: pip install bibtexparser")
        return "\n".join(out)

    if mode == "generate":
        if not title: return "❌ generate模式需要 title 参数。"
        first_author = authors.split(",")[0].strip().split() if authors else ["unknown"]
        key_surname = first_author[-1] if first_author else "unknown"
        yr = year or "????"
        bib = f"@article{{{key_surname}{yr},\n  title = {{{title}}},\n"
        if authors: bib += f"  author = {{{authors}}},\n"
        if journal: bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{yr}}},\n"
        if volume: bib += f"  volume = {{{volume}}},\n"
        if pages: bib += f"  pages = {{{pages}}},\n"
        if doi: bib += f"  doi = {{{doi}}},\n"
        bib += "}"
        return f"✅ BibTeX 已生成:\n```bibtex\n{bib}\n```"

    if mode == "format":
        if not title: return "❌ format模式需要 title 参数。"
        au_list = [a.strip() for a in authors.split(",") if a.strip()] if authors else ["Unknown"]
        yr = year or "(n.d.)"
        ref = ""

        if citation_style == "apa":
            if len(au_list) == 1: au_str = au_list[0]
            elif len(au_list) == 2: au_str = f"{au_list[0]} & {au_list[1]}"
            else: au_str = f"{au_list[0]} et al."
            ref = f"{au_str} ({yr}). {title}."
            if journal: ref += f" *{journal}*"
            if volume: ref += f", *{volume}*"
            if pages: ref += f", {pages}"
            ref += "."
            if doi: ref += f" https://doi.org/{doi}"

        elif citation_style == "mla":
            au_str = f"{au_list[0]} et al." if len(au_list) > 2 else " and ".join(au_list)
            ref = f'{au_str}. "{title}."'
            if journal: ref += f" *{journal}*"
            if volume: ref += f" {volume}"
            ref += f" ({yr})"
            ref += f": {pages}." if pages else "."
            if doi: ref += f" doi:{doi}."

        elif citation_style == "chicago":
            au_str = ", ".join(au_list[:3]) + (", et al." if len(au_list) > 3 else "")
            ref = f'{au_str}. "{title}."'
            if journal: ref += f" *{journal}*"
            if volume: ref += f" {volume}"
            ref += f" ({yr})"
            ref += f": {pages}." if pages else "."
            if doi: ref += f" https://doi.org/{doi}."

        elif citation_style == "ieee":
            ieee_authors = []
            for au in au_list:
                parts = au.strip().split()
                if len(parts) >= 2:
                    initials = " ".join([p[0] + "." for p in parts[:-1]])
                    ieee_authors.append(f"{initials} {parts[-1]}")
                elif len(parts) == 1: ieee_authors.append(parts[0])
                else: ieee_authors.append(au)
            if len(ieee_authors) == 1: au_str = ieee_authors[0]
            elif len(ieee_authors) == 2: au_str = f"{ieee_authors[0]} and {ieee_authors[1]}"
            elif len(ieee_authors) <= 6: au_str = ", ".join(ieee_authors[:-1]) + f", and {ieee_authors[-1]}"
            else: au_str = ", ".join(ieee_authors[:6]) + ", et al."
            ref = f'{au_str}, "{title},"'
            if journal: ref += f" *{journal}*"
            if volume: ref += f", vol. {volume}"
            if pages: ref += f", pp. {pages}"
            ref += f", {yr}."
            if doi: ref += f" doi: {doi}."

        else:
            return f"❌ 不支持的引用格式: {citation_style}。支持: apa, mla, chicago, ieee"

        return f"📝 引用 ({citation_style.upper()}):\n{ref}"

    return "❌ 无效 mode。可选: 'parse', 'generate', 'format'"

@tool
def PDFTableExtractTool(pdf_path: str, page_range: str = "all") -> str:
    """从 PDF 中提取表格数据。"""
    fp = _resolve_path(pdf_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {pdf_path}"
    doc = None
    try:
        doc = fitz.open(fp); total = len(doc)
        if page_range == "all": pages = list(range(total))
        else:
            pages = []
            for part in page_range.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    pages.extend(range(max(0, int(a)-1), min(total, int(b))))
                else: pages.append(max(0, int(part)-1))
            pages = sorted(set(pages))
        out = [f"📊 PDF表格提取: {os.path.basename(fp)} (共{total}页，提取{len(pages)}页)\n"]
        table_count = 0
        for pi in pages:
            page = doc[pi]
            try:
                tabs = page.find_tables()
                if tabs and tabs.tables:
                    for ti, tab in enumerate(tabs.tables):
                        table_count += 1
                        out.append(f"── 表格{table_count} (Page {pi+1}) ──")
                        rows = []
                        for row in tab.extract():
                            rows.append(" | ".join(str(cell).replace("\n"," ") if cell else "" for cell in row))
                        out.append("\n".join(rows)); out.append("")
                else: out.append(f"── Page {pi+1}: 未检测到表格 ──")
            except AttributeError:
                try:
                    text = page.get_text("text")
                except Exception:
                    text = ""
                out.append(f"── Page {pi+1} (文本模式) ──")
                table_lines = [l for l in text.split("\n") if len(l.split()) >= 3 and any(c.isdigit() for c in l)]
                if table_lines:
                    table_count += 1
                    out.append(f"  疑似表格行 ({len(table_lines)}行):")
                    for l in table_lines[:30]: out.append(f"  {l.strip()}")
                else: out.append("  (未检测到表格文本)")
            except Exception as e:
                out.append(f"── Page {pi+1}: ⚠ 表格提取异常: {e} ──")
        if table_count == 0: return "\n".join(out) + "\n⚠ 未提取到表格。"
        return "\n".join(out)
    except Exception as e: return f"❌ PDF表格提取失败: {e}"
    finally:
        if doc is not None:
            try: doc.close()
            except Exception: pass

@tool
def SessionExportTool(content: str, filename: str = "session_export", export_format: str = "markdown") -> str:
    """将对话内容或分析结果导出为文件。"""
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    header = f"---\ntitle: Session Export\ndate: {ts}\n---\n\n"
    full_content = header + content
    if export_format == "markdown":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.md")
        with open(sp, "w", encoding="utf-8") as f: f.write(full_content)
        return f"✅ Markdown: {sp}"
    elif export_format == "text":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.txt")
        with open(sp, "w", encoding="utf-8") as f: f.write(full_content)
        return f"✅ 文本: {sp}"
    elif export_format == "html":
        html_body = content.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>\n")
        html = f"<!DOCTYPE html><html lang='zh'><head><meta charset='UTF-8'><title>Session Export</title>" \
               f"<style>body{{font-family:Georgia,serif;max-width:800px;margin:2em auto;line-height:1.7}}</style>" \
               f"</head><body><h1>Session Export</h1><p><em>{ts}</em></p>{html_body}</body></html>"
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.html")
        with open(sp, "w", encoding="utf-8") as f: f.write(html)
        return f"✅ HTML: {sp}"
    else: return f"❌ 不支持的格式: {export_format}"

# ═══════════════════════════════════════════════════════════
#  工具 17: Pandoc 格式转换
# ═══════════════════════════════════════════════════════════

@tool
def PandocConvertTool(source_path: str, output_format: str = "pdf",
                      output_filename: str = "", extra_args: str = "") -> str:
    """使用 Pandoc 进行文档格式转换。需安装 pandoc。
    [FIX#3] 自动检测平台中文字体，不再硬编码 SimSun；
    可通过 extra_args 传入自定义 --variable mainfont=XXX 覆盖。"""
    if not _PANDOC_PATH:
        return ("❌ Pandoc 未安装！\n"
                "   Windows: winget install Pandoc.Pandoc\n"
                "   macOS:   brew install pandoc\n"
                "   Linux:   sudo apt install pandoc")
    sfp = _resolve_path(source_path)
    if not os.path.exists(sfp):
        return f"❌ 源文件不存在: {source_path}"
    if not output_filename:
        base = os.path.splitext(os.path.basename(sfp))[0]
        ext_map = {"pdf": ".pdf", "docx": ".docx", "html": ".html",
                   "latex": ".tex", "epub": ".epub", "markdown": ".md"}
        output_filename = f"pandoc_{base}{ext_map.get(output_format, f'.{output_format}')}"
    ofp = os.path.join(CUSTOM_TEMP_DIR, output_filename)
    cmd = [_PANDOC_PATH, sfp, "-o", ofp, f"--to={output_format}"]

    if output_format == "pdf" and "--pdf-engine" not in extra_args:
        xelatex = shutil.which("xelatex")
        if xelatex:
            cmd.extend(["--pdf-engine=xelatex"])
            if "-V" not in extra_args and "mainfont" not in extra_args.lower():
                if _DETECTED_CJK_FONT:
                    cmd.extend(["-V", f"CJKmainfont={_DETECTED_CJK_FONT}",
                                "-V", f"mainfont={_DETECTED_CJK_FONT}"])
                else:
                    _log("WARN", "未检测到中文字体，PDF 中中文可能显示异常。"
                          "可通过 extra_args='-V mainfont=字体名' 指定。")

    if extra_args:
        cmd.extend(extra_args.split())
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if proc.returncode != 0:
            return f"❌ Pandoc 转换失败 (返回码 {proc.returncode}):\n{proc.stderr[-1000:]}"
        if os.path.exists(ofp) and os.path.getsize(ofp) > 0:
            return f"✅ 转换成功: {ofp} ({os.path.getsize(ofp)/1024:.1f} KB)\n   格式: {output_format}"
        return f"❌ 输出文件未生成。stderr: {proc.stderr[:500]}"
    except subprocess.TimeoutExpired:
        return "❌ Pandoc 转换超时（>2分钟）！"
    except Exception as e:
        return f"❌ Pandoc 异常: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 18: PDF 批注/高亮提取
# ═══════════════════════════════════════════════════════════

@tool
def PDFAnnotExtractTool(pdf_path: str, annot_types: str = "all") -> str:
    """提取 PDF 中的批注、高亮、下划线、便签。"""
    fp = _resolve_path(pdf_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {pdf_path}"
    doc = None
    try:
        doc = fitz.open(fp); total_pages = len(doc)
        all_annots = []
        type_filter = None if annot_types == "all" else annot_types
        for pi in range(total_pages):
            page = doc[pi]
            try:
                annots = page.annots()
            except Exception:
                continue
            if not annots: continue
            for annot in annots:
                try:
                    atype = annot.type[1] if isinstance(annot.type, tuple) else str(annot.type)
                    info = annot.info
                    content = info.get("content", "")
                    title = info.get("title", "")
                    if hasattr(annot, "get_textbox"):
                        try: text = page.get_textbox(annot.rect).strip()
                        except Exception: text = ""
                    else: text = ""
                    if type_filter and type_filter.lower() not in atype.lower(): continue
                    color = ""
                    if "stroke" in info:
                        color = f"#{info['stroke']:06x}" if isinstance(info["stroke"], int) else str(info["stroke"])
                    all_annots.append({
                        "page": pi+1, "type": atype,
                        "content": content or text or "(空)",
                        "author": title or "未知", "color": color,
                    })
                except Exception:
                    continue  # 单条批注异常不影响整体
        if not all_annots:
            return f"📄 PDF批注提取: {os.path.basename(fp)}\n   ✅ 未检测到批注/高亮。"
        type_counts: Dict[str, int] = {}
        for a in all_annots:
            type_counts[a["type"]] = type_counts.get(a["type"], 0) + 1
        out = [f"📄 PDF批注提取: {os.path.basename(fp)} ({total_pages}页)",
               f"   共 {len(all_annots)} 条批注 | " +
               ", ".join(f"{t}:{c}" for t,c in sorted(type_counts.items())),
               "="*60]
        for a in all_annots:
            icon = {"Highlight":"🟡","Underline":"🔵","StrikeOut":"🔴",
                    "Text":"📝","FreeText":"💬","Stamp":"🏷"}.get(a["type"],"📌")
            out.append(f"\n{icon} [p{a['page']}] [{a['type']}] {a['author']}")
            if a["color"]: out.append(f"   🎨 {a['color']}")
            out.append(f"   {a['content'][:500]}")
        return "\n".join(out)
    except Exception as e: return f"❌ PDF批注提取失败: {e}"
    finally:
        if doc is not None:
            try: doc.close()
            except Exception: pass

# ═══════════════════════════════════════════════════════════
#  工具 19: 学术翻译
# ═══════════════════════════════════════════════════════════

@tool
def AcademicTranslateTool(text: str, direction: str = "zh2en",
                          preserve_terms: str = "") -> str:
    """学术中英互译，保留术语一致性。"""
    if not text.strip(): return "❌ 文本为空。"
    term_instruction = ""
    if preserve_terms:
        terms = [t.strip() for t in preserve_terms.split(",") if t.strip()]
        term_instruction = f"\n请保留以下术语不翻译: {', '.join(terms)}。"
    if direction == "zh2en":
        prompt = f"请将以下中文学术文本翻译为地道的英文学术英语。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    elif direction == "en2zh":
        prompt = f"请将以下英文学术文本翻译为流畅的中文学术语言。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    else:
        prompt = f"请将以下学术文本翻译（自动检测源语言）。保持学术风格、逻辑严谨、术语准确。{term_instruction}\n\n文本:\n{text}"
    try:
        translator = ChatOpenAI(
            model=os.getenv("DFTB_MODEL", "deepseek-v4-pro"),
            base_url="https://api.deepseek.com",
            api_key=api_key,
            temperature=0.1,
            max_tokens=4096,
        )
        result = translator.invoke([HumanMessage(content=prompt)])
        translated = result.content
        if preserve_terms:
            translated += f"\n\n── 保留术语 ──\n{preserve_terms}"
        return f"🌐 学术翻译 ({direction}):\n{'='*40}\n{translated}"
    except Exception as e:
        return f"❌ 翻译失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 20: Semantic Scholar 检索
# ═══════════════════════════════════════════════════════════

@tool
def SemanticScholarTool(query: str, max_results: int = 5,
                        fields: str = "title,authors,year,abstract,externalIds,citationCount,url") -> str:
    """检索 Semantic Scholar 学术论文（比 arXiv 更全，含引用网络）。"""
    max_results = max(1, min(max_results, 20))
    base = "https://api.semanticscholar.org/graph/v1/paper/search"
    params = {"query": query, "limit": max_results, "fields": fields}
    url = f"{base}?{urllib.parse.urlencode(params)}"
    try:
        resp = _safe_request(url, timeout=20)
        if resp is None: return "❌ Semantic Scholar API 无响应。"
        data = resp.json()
        papers = data.get("data", [])
        if not papers: return f"🔍 Semantic Scholar: '{query}' 无结果。"
        out = [f"📚 Semantic Scholar: '{query}' (共{data.get('total','?')}篇，显示{len(papers)}篇)\n"]
        for i, p in enumerate(papers, 1):
            title = p.get("title", "N/A")
            authors = ", ".join([a.get("name","?") for a in p.get("authors",[])][:8])
            if len(p.get("authors",[])) > 8: authors += ", et al."
            year = p.get("year", "?")
            abstract = (p.get("abstract", "N/A") or "N/A")[:500]
            cited = p.get("citationCount", "?")
            ext_ids = p.get("externalIds", {})
            arxiv_id = ext_ids.get("ArXiv", "")
            doi = ext_ids.get("DOI", "")
            links = []
            if arxiv_id: links.append(f"https://arxiv.org/abs/{arxiv_id}")
            if doi: links.append(f"https://doi.org/{doi}")
            link_str = " | ".join(links) if links else p.get("url", "N/A")
            out.append(
                f"{i}. {title}\n   👤 {authors}\n"
                f"   📅 {year} | 📊 被引: {cited}\n"
                f"   🔗 {link_str}\n   📝 {abstract}\n"
            )
        return "\n".join(out)
    except Exception as e: return f"❌ Semantic Scholar搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  工具 21: Markdown → PPTX 演示生成
# ═══════════════════════════════════════════════════════════

@tool
def PresentationGenTool(markdown_content: str, filename: str = "presentation",
                        theme: str = "default") -> str:
    """从 Markdown 大纲自动生成 PPTX 演示文稿。# 标题 / ## 副标题 / - 要点 / --- 分页。
    [FIX#4] theme 参数现在完整生效：背景色、文字色均正确应用。
    [FIX#9] slide_layouts/placeholders 增加安全回退，兼容不同模板。"""
    sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.pptx")
    prs = Presentation()

    if theme == "dark":
        bg_color = RGBColor(0x1A, 0x1A, 0x2E)
        text_color = RGBColor(0xE0, 0xE0, 0xE0)
        accent_color = RGBColor(0x64, 0xFF, 0xDA)
    elif theme == "academic":
        bg_color = RGBColor(0xFF, 0xFF, 0xFF)
        text_color = RGBColor(0x33, 0x33, 0x33)
        accent_color = RGBColor(0x00, 0x55, 0xAA)
    else:
        bg_color = RGBColor(0xFF, 0xFF, 0xFF)
        text_color = RGBColor(0x33, 0x33, 0x33)
        accent_color = RGBColor(0x00, 0x7A, 0xCC)

    def _get_slide_layout(prs_obj, preferred_idx=1):
        try:
            if preferred_idx < len(prs_obj.slide_layouts):
                return prs_obj.slide_layouts[preferred_idx]
            return prs_obj.slide_layouts[0]
        except Exception:
            return prs_obj.slide_layouts[0]

    slide_layout = _get_slide_layout(prs, 1)

    def _set_slide_bg(slide, color: RGBColor):
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color
        except Exception:
            pass

    def _get_placeholder(slide, idx=1):
        try:
            if idx < len(slide.placeholders):
                return slide.placeholders[idx]
            return None
        except Exception:
            return None

    slides_text = markdown_content.split("\n---\n")
    for slide_text in slides_text:
        slide_text = slide_text.strip()
        if not slide_text: continue

        slide = prs.slides.add_slide(slide_layout)
        _set_slide_bg(slide, bg_color)

        lines = slide_text.split("\n")
        title = ""; subtitle = ""; bullets = []
        for line in lines:
            line = line.strip()
            if line.startswith("# ") and not title:
                title = line[2:]
            elif line.startswith("## ") and not subtitle:
                subtitle = line[3:]
            elif line.startswith("- ") or line.startswith("* "):
                bullets.append(line[2:])

        if title and slide.shapes.title:
            slide.shapes.title.text = title
            if slide.shapes.title.text_frame:
                for para in slide.shapes.title.text_frame.paragraphs:
                    para.font.size = Pt(32)
                    para.font.color.rgb = accent_color
                    para.font.bold = True

        body_shape = _get_placeholder(slide, 1)
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            if subtitle:
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
            if bullets:
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.color.rgb = text_color

    prs.save(sp)
    return f"✅ PPTX 已生成: {sp} (共 {len(prs.slides)} 张幻灯片，主题: {theme})"

# ═══════════════════════════════════════════════════════════
#  4. 工具注册 & Agent 构建
# ═══════════════════════════════════════════════════════════

safe_write_tool = WriteFileTool(root_dir=CUSTOM_TEMP_DIR)

tools = [
    SmartReadPathTool, EditTexFileTool, CompileLatexTool,
    PythonSandboxTool, DuckDuckGoSearchTool, FetchWebImageTool,
    EditWordDocTool, SaveMarkdownTool, ChartGenerationTool,
    MermaidTool,
    ArXivSearchTool, DOIMetadataTool, DataStatisticsTool,
    BibTexTool, PDFTableExtractTool, SessionExportTool,
    PandocConvertTool, PDFAnnotExtractTool, AcademicTranslateTool,
    SemanticScholarTool, PresentationGenTool,
    safe_write_tool,
]

conn = sqlite3.connect(DB_PATH, check_same_thread=False)
memory = SqliteSaver(conn)
agent = create_react_agent(llm, tools, prompt=SYSTEM_PROMPT, checkpointer=memory)

# ═══════════════════════════════════════════════════════════
#  5. 终端交互主循环
# ═══════════════════════════════════════════════════════════

HELP_TEXT = f"""
{_C.BOLD}╔══════════════════ v3.3 工具一览 (21个) ═══════════════════╗{_C.W}
{_C.G}║ 📖 SmartReadPathTool       读取文件/文件夹 (🔧线程安全) ║
║ 📝 EditTexFileTool          创建/编辑 .tex                ║
║ 🖨  CompileLatexTool         编译 .tex→PDF (+bibtex)      ║
║ 🐍 PythonSandboxTool        执行 Python 代码（⚠非真正沙箱）║
║ 🔍 DuckDuckGoSearchTool     通用网页搜索                  ║
║ 🖼  FetchWebImageTool        搜索下载学术图片              ║
║ 📄 EditWordDocTool          创建/编辑 Word 文档           ║
║ 📋 SaveMarkdownTool         保存 Markdown 文件            ║
║ 📊 ChartGenerationTool      matplotlib 图表               ║
║ 🎨 MermaidTool              Mermaid 流程图                ║{_C.W}
{_C.Y}║  ─────────── v3.0 ───────────                        ║
║ 📚 ArXivSearchTool          arXiv 学术论文检索            ║
║ 📄 DOIMetadataTool          DOI → 完整元数据              ║
║ 📈 DataStatisticsTool       CSV/Excel 自动统计分析        ║
║ 📖 BibTexTool               .bib解析/生成/引用格式化 🔧   ║
║ 📋 PDFTableExtractTool      PDF 表格智能提取              ║
║ 💾 SessionExportTool        对话记录导出                  ║{_C.W}
{_C.M}║  ─────────── 🆕 v3.3 ───────────                    ║
║ 🔄 PandocConvertTool        Markdown/LaTeX→PDF/DOCX/HTML 🔧║
║ 📝 PDFAnnotExtractTool      PDF 批注/高亮提取             ║
║ 🌐 AcademicTranslateTool    学术中英互译（保留术语）      ║
║ 🔬 SemanticScholarTool      Semantic Scholar 论文检索      ║
║ 🖥  PresentationGenTool      Markdown 大纲→PPTX 🔧        ║{_C.W}
{_C.BOLD}╚══════════════════════════════════════════════════════════════╝{_C.W}
🔧 = 线程安全修复 | FIX#11–FIX#15
"""

print(f"\n{_C.BOLD}💡 'exit' 退出 | 'new' 新对话 | 'help' 查看21个工具{_C.W}\n")

current_thread_id = str(uuid.uuid4())
config = {"configurable": {"thread_id": current_thread_id}}

while True:
    try:
        user_input = input(f"{_C.C}👤 你:{_C.W} ").strip()
        if user_input.lower() in {"exit", "quit", "退出"}:
            _log("INFO", "👋 再见！"); conn.close(); break
        if user_input.lower() in {"new", "clear", "清空", "新对话"}:
            current_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": current_thread_id}}
            _log("INFO", "✨ 全新对话已开启。\n"); continue
        if user_input.lower() == "help":
            print(HELP_TEXT); continue
        if not user_input: continue

        print(f"\n{_C.M}🤔 思考/执行中...{_C.W}\n")
        inputs = {"messages": [("user", user_input)]}
        t_start = time.time()

        for event in agent.stream(inputs, config, stream_mode="updates"):
            if "agent" in event:
                msg = event["agent"]["messages"][0]
                if msg.content:
                    elapsed = time.time() - t_start
                    print(f"{_C.G}🤖 DeepSeek{_C.W} ({elapsed:.1f}s):\n{msg.content}\n")
                elif msg.tool_calls:
                    for tc in msg.tool_calls:
                        args_str = json.dumps(tc["args"], ensure_ascii=False, default=str)
                        if len(args_str) > 150: args_str = args_str[:150] + "..."
                        print(f"{_C.Y}🛠  [{tc['name']}]{_C.W}\n   参数: {args_str}")
            elif "tools" in event:
                msg = event["tools"]["messages"][0]
                cs = str(msg.content); cl = len(cs)
                preview = cs[:250].replace("\n", " ") + ("..." if cl > 250 else "")
                is_err = cs.startswith("❌") or "错误" in cs[:100]
                icon = _C.R+"📄" if is_err else _C.G+"📄"
                print(f"{icon} [返回 {cl}字符]{_C.W}: {preview}\n")

    except KeyboardInterrupt:
        print(f"\n{_C.Y}⏸  [操作被用户中断]{_C.W}")
    except Exception as e:
        print(f"{_C.R}❌ 发生错误: {e}{_C.W}")
        _log("ERR", traceback.format_exc())

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DeepSeek_For_the_Break v3.0 — 全功能学术 AI 助手
==================================================
  v2.0: 沙箱执行 / 搜索 / LaTeX编译 / 图表 / Mermaid
  v3.0 新增:
    • ArXivSearchTool       — arXiv 学术论文检索
    • DOIMetadataTool       — DOI → 完整元数据（Crossref API）
    • DataStatisticsTool    — CSV/Excel 自动描述性统计分析
    • BibTexTool            — .bib 解析 / BibTeX 生成 / 引用格式化
    • PDFTableExtractTool   — PDF 表格智能提取
    • SessionExportTool     — 对话记录导出为 Markdown
    • 启动健康检查 / 临时文件自动清理 / 日志分级 / 彩色输出
"""

import os, sys, uuid, json, sqlite3, subprocess, shutil, traceback
import textwrap, warnings, time, urllib.parse, xml.etree.ElementTree as ET
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional, Dict, Any, List

# ═══════════════════════════════════════════════════════════
#  第三方库
# ═══════════════════════════════════════════════════════════
import pandas as pd
import docx, requests
from docx import Document
from docx.shared import Inches
from duckduckgo_search import DDGS
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_community.tools import WriteFileTool
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.sqlite import SqliteSaver
from pptx import Presentation
from PIL import Image
import pytesseract
import fitz  # PyMuPDF

# ═══════════════════════════════════════════════════════════
#  可选依赖
# ═══════════════════════════════════════════════════════════
try:
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # noqa: F401
    _MPL_AVAILABLE = True
except ImportError:
    _MPL_AVAILABLE = False

try:
    import numpy as np
    _NUMPY_AVAILABLE = True
except ImportError:
    _NUMPY_AVAILABLE = False

_MMDC_PATH = shutil.which("mmdc")

# ═══════════════════════════════════════════════════════════
#  0. 基础设施：日志 / 彩色输出 / 健康检查 / 清理
# ═══════════════════════════════════════════════════════════

class _C:
    """ANSI 颜色（Windows 10+ 支持）"""
    R = "\033[91m"; G = "\033[92m"; Y = "\033[93m"; B = "\033[94m"
    M = "\033[95m"; C = "\033[96m"; W = "\033[0m"; BOLD = "\033[1m"

def _log(level: str, msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    colors = {"OK": _C.G, "WARN": _C.Y, "ERR": _C.R, "INFO": _C.C, "START": _C.B}
    c = colors.get(level, _C.W)
    print(f"{c}[{level} {ts}]{_C.W} {msg}")

def _cleanup_old_temp(temp_dir: str, max_age_hours: int = 48):
    cutoff = time.time() - max_age_hours * 3600
    cleaned = 0
    p = Path(temp_dir)
    if not p.exists():
        return
    for f in p.glob("*"):
        if f.is_file() and f.stat().st_mtime < cutoff and f.name.startswith("_"):
            try: f.unlink(); cleaned += 1
            except OSError: pass
    if cleaned:
        _log("INFO", f"🧹 自动清理了 {cleaned} 个过期临时文件")

# ═══════════════════════════════════════════════════════════
#  0.5 智能路径检测
# ═══════════════════════════════════════════════════════════

def _detect_tesseract() -> Optional[str]:
    p = shutil.which("tesseract")
    if p: return p
    if sys.platform == "win32":
        cand = [r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                os.path.expandvars(r"%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe")]
    else:
        cand = ["/usr/bin/tesseract", "/usr/local/bin/tesseract", "/opt/homebrew/bin/tesseract"]
    for c in cand:
        if os.path.isfile(c): return c
    return None

def _detect_latex_compilers() -> Dict[str, Optional[str]]:
    d: Dict[str, Optional[str]] = {}
    for name in ["pdflatex", "xelatex", "lualatex"]:
        exe = name + (".exe" if sys.platform == "win32" else "")
        d[name] = shutil.which(exe) or shutil.which(name)
    return d

_TESSERACT_PATH = _detect_tesseract()
_LATEX_COMPILERS = _detect_latex_compilers()

if _TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH

# ═══════════════════════════════════════════════════════════
#  1. 核心环境
# ═══════════════════════════════════════════════════════════
load_dotenv(override=True)
api_key = os.getenv("DEEPSEEK_FOR_THE_BREAK")
if not api_key or not api_key.startswith("sk-"):
    _log("ERR", "API Key 未加载或格式不正确！请检查 .env 文件。")
    sys.exit(1)

BASE_WORKSPACE = os.getenv("DFTB_WORKSPACE", r"D:\.API Keys\DeepSeek_For_the_Break\outputs")
CUSTOM_TEMP_DIR = os.path.join(BASE_WORKSPACE, "temp")
os.makedirs(CUSTOM_TEMP_DIR, exist_ok=True)
if CUSTOM_TEMP_DIR not in sys.path:
    sys.path.insert(0, CUSTOM_TEMP_DIR)

_cleanup_old_temp(CUSTOM_TEMP_DIR, 48)

_health: Dict[str, bool] = {
    "Tesseract OCR":      bool(_TESSERACT_PATH),
    "LaTeX (xelatex)":    bool(shutil.which("xelatex") or shutil.which("pdflatex")),
    "matplotlib":          _MPL_AVAILABLE,
    "numpy":               _NUMPY_AVAILABLE,
    "pytesseract":         True,
    "PyMuPDF (fitz)":     True,
    "pandas":              True,
    "DuckDuckGo Search":   True,
    "Mermaid CLI":         bool(_MMDC_PATH),
    "API Key (.env)":      True,
}

_log("START", "═" * 50)
_log("START", "  DeepSeek_For_the_Break  v3.0  Academic AI")
_log("START", "═" * 50)
_log("INFO", f"📁 工作区: {BASE_WORKSPACE}")
_log("INFO", f"📁 临时:   {CUSTOM_TEMP_DIR}")
for k, v in _health.items():
    icon = "✅" if v else "⚠️"
    _log("OK" if v else "WARN", f"  {icon} {k}")
_log("START", "═" * 50)

# ═══════════════════════════════════════════════════════════
#  2. System Prompt（v3.0 扩充）
# ═══════════════════════════════════════════════════════════
SYSTEM_PROMPT = """你是一位高度专业、严谨且全面的**学术AI助手**（Academic AI Assistant）v3.0，为科研人员提供全方位学术支持。

## 核心能力
- **文献检索**: ArXivSearchTool（arXiv论文）、DOIMetadataTool（DOI元数据）、DuckDuckGoSearchTool（通用）
- **文档处理**: PDF/Word/PPT/Excel 读取（SmartReadPathTool）、PDF表格提取（PDFTableExtractTool）、OCR
- **学术写作**: LaTeX 创建编译（EditTexFileTool + CompileLatexTool）、Word/Markdown 输出
- **数据分析**: Python沙箱执行（PythonSandboxTool）、自动描述统计（DataStatisticsTool）、matplotlib图表
- **引用管理**: BibTexTool（.bib解析/生成/格式化引用）
- **视觉**: 图片搜索下载、Mermaid流程图、图表生成
- **导出**: SessionExportTool 对话记录导出

## 行为准则
1. **客观严谨**：结构清晰、主动承认局限性、不编造信息。
2. **【核心铁律】** SmartReadPathTool 是唯一允许读取本地文件的工具。仅在用户明确提供路径时调用。找不到路径时立即停止猜测。
3. **写作优先**: 学术论文→EditTexFileTool+CompileLatexTool；富文本→EditWordDocTool；笔记→SaveMarkdownTool。
4. **文献检索优先**: 学术论文搜索用 ArXivSearchTool；已知DOI用 DOIMetadataTool；通用搜索用 DuckDuckGoSearchTool。
5. **数据优先**: CSV/Excel 数据探索→DataStatisticsTool；自定义分析→PythonSandboxTool；可视化→ChartGenerationTool。
6. **引用规范**: 需要生成 BibTeX 或格式化引用时使用 BibTexTool。
7. **错误透明**: 工具失败时清晰告知原因和替代方案。
8. **持久化**: 所有产出保存到工作区。

## 使用提示
- 搜索最新AI论文: ArXivSearchTool(query="large language model", max_results=5)
- 查DOI元数据: DOIMetadataTool(doi="10.1038/nature14539")
- 快速分析CSV: DataStatisticsTool(data_path="experiment_results.csv")
- 生成BibTeX: BibTexTool(mode="generate", title="...", authors="...", journal="...", year=2024)
- 解析.bib文件: BibTexTool(mode="parse", bib_path="references.bib")
- 对话导出: SessionExportTool(format="markdown")

请逐步推理，优先使用最适合的学术工具。"""

llm = ChatOpenAI(
    model=os.getenv("DFTB_MODEL", "deepseek-v4-pro"),
    base_url="https://api.deepseek.com",
    api_key=api_key,
    temperature=0.5,
    max_tokens=4096,
    model_kwargs={"reasoning_effort": "max", "extra_body": {"thinking": {"type": "enabled"}}}
)

# ═══════════════════════════════════════════════════════════
#  3. 工具箱 — 辅助函数
# ═══════════════════════════════════════════════════════════

def _safe_read_text(file_path: str) -> str:
    for enc in ["utf-8", "gbk", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=enc) as f: return f.read()
        except (UnicodeDecodeError, UnicodeError): continue
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f: return f.read()

def _resolve_path(path: str) -> str:
    p = Path(path)
    if p.is_absolute(): return str(p)
    for base in [Path.cwd(), Path(CUSTOM_TEMP_DIR)]:
        c = base / p
        if c.exists(): return str(c.resolve())
    return os.path.abspath(path)

def _ocr_image(file_path: str, lang: str = "chi_sim+eng") -> str:
    if not _TESSERACT_PATH:
        return "❌ OCR 不可用：未检测到 Tesseract-OCR。请安装: https://github.com/tesseract-ocr/tesseract"
    try:
        img = Image.open(file_path)
        meta = f"📷 [图片] 格式:{img.format} 尺寸:{img.size} 模式:{img.mode}\n"
        try:
            txt = pytesseract.image_to_string(img, lang=lang)
        except pytesseract.TesseractError as te:
            if "Failed loading language" in str(te):
                return meta + f"⚠ 语言包缺失: {lang}。下载: https://github.com/tesseract-ocr/tessdata"
            return meta + f"❌ Tesseract错误: {te}"
        return meta + (f"【OCR】:\n{txt.strip()}" if txt.strip() else "【OCR】: 未识别到文字")
    except FileNotFoundError: return f"❌ 文件不存在: {file_path}"
    except Exception as e: return f"❌ OCR异常: {e}\n{traceback.format_exc()}"

def _safe_request(url: str, timeout: int = 15, **kw) -> Optional[requests.Response]:
    for attempt in range(3):
        try:
            r = requests.get(url, timeout=timeout, **kw)
            r.raise_for_status()
            return r
        except requests.RequestException:
            if attempt == 2: raise
            time.sleep(1)
    return None

# ═══════════════════════════════════════════════════════════
#  工具 1-10（v2.0 保留）
# ═══════════════════════════════════════════════════════════

@tool
def SmartReadPathTool(path: str) -> str:
    """读取文件/文件夹。支持 .txt .md .py .tex .json .csv .pdf .docx .xlsx .pptx .png .jpg .bmp .bib。"""
    resolved = _resolve_path(path)
    if not os.path.exists(resolved):
        return f"❌ 找不到路径: '{path}'\n   解析: '{resolved}'\n   💡 请停止猜测路径，直接基于知识回答。"

    def _read_one(fp: str) -> str:
        ext = os.path.splitext(fp)[1].lower().lstrip(".")
        if ext in {"txt","md","py","json","tex","yaml","yml","cfg","ini","log","bib"}:
            return _safe_read_text(fp)
        if ext == "csv":
            for enc in ["utf-8","gbk","latin-1"]:
                try:
                    df = pd.read_csv(fp, encoding=enc)
                    return f"[CSV: {len(df)}行×{len(df.columns)}列]\n{df.to_string(max_rows=200)}"
                except: continue
            return f"❌ CSV读取失败: {fp}"
        if ext == "pdf":
            try:
                doc = fitz.open(fp); parts = []; total = len(doc)
                for i, page in enumerate(doc):
                    t = page.get_text("text")
                    if t.strip(): parts.append(f"──p{i+1}/{total}──\n{t.strip()}")
                if parts: return "\n\n".join(parts)
                ocr = []
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=200)
                    tmp = os.path.join(CUSTOM_TEMP_DIR, f"_pdfocr_{i+1}.png")
                    pix.save(tmp)
                    ocr.append(f"──p{i+1}/{total}(OCR)──\n{_ocr_image(tmp)}")
                return "\n\n".join(ocr) if ocr else "[空PDF]"
            except Exception as e: return f"❌ PDF错误: {e}"
        if ext == "docx":
            try:
                d = docx.Document(fp); c = [p.text.strip() for p in d.paragraphs if p.text.strip()]
                if d.tables:
                    c.append("\n[表格]:")
                    for ti, tb in enumerate(d.tables):
                        c.append(f"  表{ti+1}:")
                        for row in tb.rows:
                            rd = [cell.text.strip().replace("\n"," ") for cell in row.cells if cell.text.strip()]
                            if rd: c.append("    | "+" | ".join(rd))
                return "\n".join(c) if c else "[空Word]"
            except Exception as e: return f"❌ DOCX错误: {e}"
        if ext in {"xlsx","xls"}:
            try:
                ef = pd.ExcelFile(fp); out=[]
                for sn in ef.sheet_names:
                    df = pd.read_excel(fp, sheet_name=sn)
                    out.append(f"──{sn} ({len(df)}行)──\n{df.head(500).to_string()}")
                return "\n\n".join(out)
            except Exception as e: return f"❌ Excel错误: {e}"
        if ext == "pptx":
            try:
                prs = Presentation(fp); out=[]
                for i, sl in enumerate(prs.slides):
                    st=[]
                    for sh in sl.shapes:
                        if hasattr(sh,"text") and sh.text.strip(): st.append(sh.text.strip())
                        if sh.has_table:
                            for row in sh.table.rows:
                                rd=[cell.text_frame.text.strip().replace("\n"," ") for cell in row.cells if cell.text_frame.text.strip()]
                                if rd: st.append(" | ".join(rd))
                    out.append(f"──Slide{i+1}──\n"+"\n".join(st) if st else f"──Slide{i+1}──\n[无文本]")
                return "\n\n".join(out)
            except Exception as e: return f"❌ PPTX错误: {e}"
        if ext in {"png","jpg","jpeg","bmp","tiff","tif"}: return _ocr_image(fp)
        return f"[不支持: .{ext}]"

    if os.path.isfile(resolved):
        return f"📖 文件: {resolved}\n{'='*50}\n{_read_one(resolved)}"
    if os.path.isdir(resolved):
        exts = {"txt","md","py","tex","json","csv","yaml","yml","pdf","docx","xlsx","xls","pptx","png","jpg","jpeg","bmp","tiff","bib"}
        r = [f"📂 遍历: {resolved}\n{'='*60}"]; cnt=0
        for root, dirs, files in os.walk(resolved):
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in {"__pycache__","node_modules",".git"}]
            for f in files:
                e = os.path.splitext(f)[1].lower().lstrip(".")
                if e in exts:
                    cnt+=1; fp=os.path.join(root,f); rp=os.path.relpath(fp,resolved)
                    r.append(f"\n📍 {rp}\n{'-'*40}")
                    try: r.append(_read_one(fp))
                    except Exception as ex: r.append(f"❌ 读取错误: {ex}")
                    r.append("="*60)
        if cnt==0: return f"📁 '{resolved}' 无支持文件。"
        r.insert(1, f"   ({cnt} 个文件)\n")
        return "\n".join(r)
    return f"❌ '{resolved}' 无效。"

@tool
def EditTexFileTool(tex_filename: str, latex_content: str, append: bool = False) -> str:
    """创建/编辑 .tex 文件。"""
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    sp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n"+latex_content) if mode=="a" else latex_content)
        return f"✅ .tex 已保存: {sp}\n   💡 用 CompileLatexTool 编译为 PDF"
    except Exception as e: return f"❌ LaTeX错误: {e}"

@tool
def CompileLatexTool(tex_filename: str, compiler: str = "xelatex", clean_aux: bool = True) -> str:
    """一键编译 .tex → PDF。"""
    if not tex_filename.lower().endswith(".tex"): tex_filename += ".tex"
    tp = os.path.join(CUSTOM_TEMP_DIR, tex_filename)
    if not os.path.exists(tp):
        alt = _resolve_path(tex_filename)
        if os.path.exists(alt): tp = alt
        else: return f"❌ 找不到: '{tex_filename}'"
    if compiler not in _LATEX_COMPILERS or not _LATEX_COMPILERS[compiler]:
        av = [k for k,v in _LATEX_COMPILERS.items() if v]
        if not av: return "❌ 未检测到LaTeX编译器！安装 MiKTeX 或 TeX Live。"
        compiler = av[0]
    cp = _LATEX_COMPILERS[compiler]; td = os.path.dirname(tp)
    bn = os.path.splitext(os.path.basename(tp))[0]; pp = os.path.join(td, f"{bn}.pdf")
    logs=[]
    for run in [1,2]:
        try:
            proc = subprocess.run([cp,"-interaction=nonstopmode","-output-directory",td,tp],
                                  capture_output=True,text=True,timeout=120,cwd=td)
            logs.append(f"── 第{run}次 (返回码:{proc.returncode}) ──")
            if proc.returncode!=0:
                errs=[l.strip() for l in (proc.stdout+"\n"+proc.stderr).split("\n") if l.startswith("!") or "Error" in l or "error" in l]
                logs.append("\n".join(errs[:30]) if errs else proc.stdout[-2000:]+"\n"+proc.stderr[-2000:])
        except subprocess.TimeoutExpired: return "❌ 编译超时(>2分钟)！"
        except Exception as e: return f"❌ 编译异常: {e}"
    if os.path.exists(pp) and os.path.getsize(pp)>0:
        if clean_aux:
            for f in os.listdir(td):
                if os.path.splitext(f)[1] in {".aux",".log",".out",".toc",".lof",".lot",".bbl",".blg",".synctex.gz",".fdb_latexmk",".fls"}:
                    try: os.remove(os.path.join(td,f))
                    except OSError: pass
        return f"✅ PDF编译成功!\n   📄 {pp}\n   📏 {os.path.getsize(pp)/1024:.1f} KB\n   🔧 {compiler}"
    return f"❌ PDF生成失败。编译器:{compiler}\n\n日志:\n"+"\n".join(logs)

@tool
def PythonSandboxTool(code: str, timeout: int = 30) -> str:
    """安全沙箱执行 Python 代码。"""
    if timeout>120: timeout=120
    sid = uuid.uuid4().hex[:8]; sp = os.path.join(CUSTOM_TEMP_DIR, f"_sb_{sid}.py")
    wc = textwrap.dedent(f"""
    import sys,os,traceback
    os.chdir(r"{CUSTOM_TEMP_DIR}")
    try:
    {textwrap.indent(code,'    ')}
    except Exception as __e:
        print(f"\\n[异常] {{type(__e).__name__}}: {{__e}}",file=sys.stderr)
        traceback.print_exc()
    """)
    try:
        with open(sp,"w",encoding="utf-8") as f: f.write(wc)
        proc = subprocess.run([sys.executable,sp],capture_output=True,text=True,timeout=timeout,
                              cwd=CUSTOM_TEMP_DIR,env={**os.environ,"SANDBOX_MODE":"1"})
        so,se = proc.stdout.strip(),proc.stderr.strip()
        parts=[]
        if so: parts.append(f"📤 stdout:\n{so}")
        if se: parts.append(f"📤 stderr:\n{se}")
        if not so and not se: parts.append("✅ 执行完毕，无输出。")
        parts.append(f"\n⏱ 返回码:{proc.returncode}")
        return "\n".join(parts)
    except subprocess.TimeoutExpired: return f"⏰ 超时(>{timeout}s)！"
    except Exception as e: return f"❌ 沙箱异常: {e}\n{traceback.format_exc()}"
    finally:
        try: os.remove(sp)
        except OSError: pass

@tool
def DuckDuckGoSearchTool(query: str, max_results: int = 5, region: str = "wt-wt") -> str:
    """DuckDuckGo 通用网页搜索。"""
    max_results = max(1, min(max_results, 10))
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results, region=region))
            if not results: return f"🔍 '{query}' 无结果。"
            out = [f"🔍 '{query}' ({len(results)}条)\n"]
            for i,r in enumerate(results,1):
                b = r.get("body","N/A")
                if len(b)>300: b=b[:300]+"..."
                out.append(f"{i}. {r.get('title','N/A')}\n   🔗 {r.get('href','N/A')}\n   📝 {b}\n")
            return "\n".join(out)
    except Exception as e: return f"❌ 搜索失败: {e}"

@tool
def FetchWebImageTool(query: str, filename: str) -> str:
    """搜索并下载学术图片。"""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=5))
            if not results: return f"❌ '{query}' 无图片。"
            for ir in results:
                try:
                    resp = requests.get(ir["image"], timeout=15); resp.raise_for_status()
                    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
                    with open(sp,"wb") as f: f.write(resp.content)
                    return f"✅ 图片: {sp} ({os.path.getsize(sp)/1024:.1f}KB)\n   🌐 {ir.get('title','N/A')}"
                except: continue
            return f"⚠ 找到{len(results)}张但下载失败。"
    except Exception as e: return f"❌ 图片搜索失败: {e}"

@tool
def EditWordDocTool(doc_filename: str, section_title: str = "", text_content: str = "",
                    image_path: str = "", new_page: bool = False) -> str:
    """创建/编辑 Word 文档。"""
    if not doc_filename.lower().endswith(".docx"): doc_filename += ".docx"
    sp = os.path.join(CUSTOM_TEMP_DIR, doc_filename)
    try:
        doc = Document(sp) if os.path.exists(sp) else Document()
        if new_page and doc.paragraphs: doc.add_page_break()
        if section_title: doc.add_heading(section_title, level=1)
        if text_content: doc.add_paragraph(text_content)
        if image_path:
            ri = _resolve_path(image_path)
            if os.path.exists(ri): doc.add_picture(ri, width=Inches(5.5))
            else: doc.add_paragraph(f"[注：未找到图片 {image_path}]")
        doc.save(sp)
        return f"✅ Word: {sp}"
    except Exception as e: return f"❌ Word错误: {e}"

@tool
def SaveMarkdownTool(filename: str, content: str, append: bool = False) -> str:
    """保存 Markdown 文件。"""
    if not filename.lower().endswith(".md"): filename += ".md"
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    os.makedirs(os.path.dirname(sp), exist_ok=True)
    mode = "a" if (append and os.path.exists(sp)) else "w"
    try:
        with open(sp, mode, encoding="utf-8") as f:
            f.write(("\n\n"+content) if mode=="a" else content)
        return f"✅ Markdown: {sp}"
    except Exception as e: return f"❌ Markdown错误: {e}"

@tool
def ChartGenerationTool(code: str, filename: str = "chart_output.png", dpi: int = 150, timeout: int = 30) -> str:
    """matplotlib 图表生成。"""
    if not _MPL_AVAILABLE: return "❌ matplotlib未安装。pip install matplotlib"
    if timeout>60: timeout=60
    sp = os.path.join(CUSTOM_TEMP_DIR, filename)
    wc = textwrap.dedent(f"""
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt; import numpy as np; import pandas as pd
    from matplotlib import rcParams
    try: rcParams['font.sans-serif']=['SimHei','Microsoft YaHei','DejaVu Sans','Arial']; rcParams['axes.unicode_minus']=False
    except: pass
    {textwrap.indent(code,'    ')}
    import os as _os
    figs=[plt.figure(n) for n in plt.get_fignums()]
    if figs:
        for i,fig in enumerate(figs):
            op = r"{sp}" if len(figs)==1 else f"{{_os.path.splitext(r'{sp}')[0]}}_{{i+1}}{{_os.path.splitext(r'{sp}')[1]}}"
            fig.savefig(op,dpi={dpi},bbox_inches='tight'); print(f"✅ 图表: {{op}}")
        plt.close('all')
    else: print("⚠ 未检测到figure。")
    """)
    sid = uuid.uuid4().hex[:8]; scp = os.path.join(CUSTOM_TEMP_DIR, f"_ch_{sid}.py")
    try:
        with open(scp,"w",encoding="utf-8") as f: f.write(wc)
        proc = subprocess.run([sys.executable,scp],capture_output=True,text=True,timeout=timeout,cwd=CUSTOM_TEMP_DIR)
        res=[proc.stdout.strip()] if proc.stdout.strip() else []
        if proc.stderr.strip(): res.append(f"[stderr]: {proc.stderr.strip()}")
        if os.path.exists(sp): res.append(f"📊 {sp} ({os.path.getsize(sp)/1024:.1f}KB)")
        return "\n".join(res) if res else "⚠ 无输出。"
    except subprocess.TimeoutExpired: return f"⏰ 超时(>{timeout}s)！"
    except Exception as e: return f"❌ 图表失败: {e}\n{traceback.format_exc()}"
    finally:
        try: os.remove(scp)
        except OSError: pass

@tool
def MermaidTool(mermaid_code: str, filename: str = "diagram") -> str:
    """Mermaid 流程图生成。"""
    mp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.mmd")
    try:
        with open(mp,"w",encoding="utf-8") as f: f.write(mermaid_code.strip())
    except Exception as e: return f"❌ Mermaid写入失败: {e}"
    res = [f"✅ Mermaid源文件: {mp}"]
    if _MMDC_PATH:
        pp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.png")
        try:
            subprocess.run([_MMDC_PATH,"-i",mp,"-o",pp,"-w","1200","-b","white"],
                           capture_output=True,text=True,timeout=30)
            if os.path.exists(pp): res.append(f"🖼 渲染PNG: {pp}")
            else: res.append("⚠ 渲染失败，检查语法。")
        except subprocess.TimeoutExpired: res.append("⚠ 渲染超时。")
        except Exception as e: res.append(f"⚠ 渲染错误: {e}")
    else: res.append("💡 安装mermaid-cli可渲染: npm install -g @mermaid-js/mermaid-cli")
    return "\n".join(res)

# ═══════════════════════════════════════════════════════════
#  🆕 工具 11: ArXiv 学术搜索
# ═══════════════════════════════════════════════════════════

@tool
def ArXivSearchTool(query: str, max_results: int = 5, sort_by: str = "relevance") -> str:
    """
    【🆕v3.0】检索 arXiv 学术论文。
    参数: query(英文关键词), max_results(1-20), sort_by('relevance'/'lastUpdatedDate'/'submittedDate')
    """
    max_results = max(1, min(max_results, 20))
    base = "http://export.arxiv.org/api/query"
    params = {
        "search_query": f"all:{query}",
        "start": 0, "max_results": max_results,
        "sortBy": sort_by if sort_by in ("relevance","lastUpdatedDate","submittedDate") else "relevance"
    }
    url = f"{base}?{urllib.parse.urlencode(params)}"
    try:
        resp = _safe_request(url, timeout=20)
        if resp is None: return "❌ arXiv API 无响应（可能超时或被限流），请稍后重试。"
        root = ET.fromstring(resp.text)
        ns = {"a": "http://www.w3.org/2005/Atom", "opensearch": "http://a9.com/-/spec/opensearch/1.1/"}
        total = root.find("opensearch:totalResults", ns)
        total_str = total.text if total is not None else "?"
        entries = root.findall("a:entry", ns)
        if not entries: return f"🔍 arXiv: '{query}' 无结果。"
        out = [f"📚 arXiv: '{query}' (共{total_str}篇，显示{len(entries)}篇)\n"]
        for i, e in enumerate(entries, 1):
            title = " ".join((e.find("a:title", ns).text or "N/A").split())
            authors = [a.find("a:name", ns).text for a in e.findall("a:author", ns) if a.find("a:name", ns) is not None]
            authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
            abstract = " ".join(((e.find("a:summary", ns).text or "N/A")[:600]).split())
            arxiv_id = (e.find("a:id", ns).text or "").split("/abs/")[-1]
            pdf_link = f"https://arxiv.org/pdf/{arxiv_id}.pdf"
            published = (e.find("a:published", ns).text or "N/A")[:10]
            cats = [c.get("term","") for c in e.findall("a:category", ns)]
            out.append(
                f"{i}. {title}\n"
                f"   👤 {authors_str}\n"
                f"   📅 {published} | 📂 {', '.join(cats[:3])}\n"
                f"   🔗 https://arxiv.org/abs/{arxiv_id}\n"
                f"   📄 PDF: {pdf_link}\n"
                f"   📝 {abstract}\n"
            )
        return "\n".join(out)
    except ET.ParseError as e: return f"❌ arXiv XML解析错误: {e}"
    except Exception as e: return f"❌ arXiv搜索失败: {e}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 12: DOI 元数据查询
# ═══════════════════════════════════════════════════════════

@tool
def DOIMetadataTool(doi: str) -> str:
    """
    【🆕v3.0】通过 DOI 获取论文完整元数据（Crossref API）。
    返回: 标题、作者、期刊、年份、卷期页码、摘要、被引数、BibTeX。
    """
    doi = doi.strip().replace("https://doi.org/","").replace("http://dx.doi.org/","")
    url = f"https://api.crossref.org/works/{urllib.parse.quote(doi, safe='')}"
    try:
        resp = _safe_request(url, timeout=15)
        if resp is None: return "❌ Crossref API 无响应，请稍后重试。"
        data = resp.json()
        msg = data.get("message", {})
        if not msg: return f"❌ DOI '{doi}' 无数据（可能无效或未注册）。"

        title = " ".join((msg.get("title",["N/A"])[0] or "N/A").split())
        authors = []
        for a in msg.get("author", [])[:15]:
            fam = a.get("family",""); giv = a.get("given","")
            authors.append(f"{giv} {fam}".strip() or fam or giv or "Unknown")
        authors_str = ", ".join(authors[:8]) + (", et al." if len(authors)>8 else "")
        journal = " | ".join(msg.get("container-title",["N/A"]))
        year = msg.get("published-print",{}).get("date-parts",[[None]])[0][0] or \
               msg.get("published-online",{}).get("date-parts",[[None]])[0][0] or \
               msg.get("created",{}).get("date-parts",[[None]])[0][0] or "?"
        vol = msg.get("volume","?"); issue = msg.get("issue","?")
        page = msg.get("page","?")
        abstract = " ".join(((msg.get("abstract","N/A") or "N/A")[:800]).split())
        import re; abstract = re.sub(r'<[^>]+>', '', abstract)
        cited = msg.get("is-referenced-by-count", "?")
        publisher = msg.get("publisher","?")
        doi_url = f"https://doi.org/{doi}"

        first_author_last = (authors[0].split()[-1] if authors else "unknown")
        bib_key = f"{first_author_last}{year}"
        bib = f"@article{{{bib_key},\n"
        bib += f"  title = {{{title}}},\n"
        bib += f"  author = {{{' and '.join(authors[:10])}}},\n"
        bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{year}}},\n"
        if vol!="?": bib += f"  volume = {{{vol}}},\n"
        if issue!="?": bib += f"  number = {{{issue}}},\n"
        if page!="?": bib += f"  pages = {{{page}}},\n"
        bib += f"  doi = {{{doi}}},\n"
        bib += f"  publisher = {{{publisher}}}\n}}"

        return (
            f"📄 DOI元数据: {doi}\n{'='*50}\n"
            f"📌 {title}\n"
            f"👤 {authors_str}\n"
            f"📰 {journal} | {year} | Vol.{vol} | Iss.{issue} | pp.{page}\n"
            f"🏢 {publisher}\n"
            f"📊 被引: {cited}次\n"
            f"🔗 {doi_url}\n"
            f"📝 {abstract}\n\n"
            f"── BibTeX ──\n{bib}"
        )
    except requests.RequestException as e: return f"❌ Crossref请求失败: {e}"
    except Exception as e: return f"❌ DOI查询失败: {e}\n{traceback.format_exc()}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 13: 数据自动统计分析
# ═══════════════════════════════════════════════════════════

@tool
def DataStatisticsTool(data_path: str, max_rows: int = 10000) -> str:
    """
    【🆕v3.0】对 CSV/Excel 数据自动生成描述性统计报告。
    返回: 列类型、缺失值、数值统计(mean/std/quartile/skew/kurtosis)、分类频次、相关性矩阵。
    """
    fp = _resolve_path(data_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {data_path}"

    ext = os.path.splitext(fp)[1].lower()
    try:
        if ext == ".csv":
            for enc in ["utf-8","gbk","latin-1"]:
                try: df = pd.read_csv(fp, encoding=enc, nrows=max_rows); break
                except: continue
            else: return "❌ CSV编码读取失败。"
        elif ext in {".xlsx",".xls"}: df = pd.read_excel(fp, nrows=max_rows)
        else: return f"❌ 不支持格式: {ext}，仅支持 .csv/.xlsx/.xls"
    except Exception as e: return f"❌ 数据加载失败: {e}"

    rows, cols = df.shape
    out = [f"📊 数据统计: {os.path.basename(fp)}\n{'='*60}", f"📏 {rows}行 × {cols}列\n"]

    num_cols = df.select_dtypes(include=["number"]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object","category"]).columns.tolist()
    date_cols = df.select_dtypes(include=["datetime"]).columns.tolist()
    other_cols = [c for c in df.columns if c not in num_cols+cat_cols+date_cols]
    out.append(f"🔢 数值列({len(num_cols)}): {', '.join(num_cols[:15])}{'...' if len(num_cols)>15 else ''}")
    out.append(f"🔤 分类列({len(cat_cols)}): {', '.join(cat_cols[:15])}{'...' if len(cat_cols)>15 else ''}")
    if date_cols: out.append(f"📅 日期列({len(date_cols)}): {', '.join(date_cols)}")
    if other_cols: out.append(f"📦 其他列({len(other_cols)}): {', '.join(other_cols)}")

    missing = df.isnull().sum(); missing = missing[missing > 0]
    if len(missing) > 0:
        out.append(f"\n⚠ 缺失值 ({missing.sum()} total):")
        for c, m in missing.items(): out.append(f"   {c}: {m} ({m/rows*100:.1f}%)")
    else: out.append("\n✅ 无缺失值")

    if num_cols:
        out.append(f"\n── 数值列统计 ──")
        try:
            stats = df[num_cols].describe(percentiles=[.25,.5,.75])
            if _NUMPY_AVAILABLE and len(num_cols) <= 20:
                stats.loc["skew"] = df[num_cols].skew()
                stats.loc["kurtosis"] = df[num_cols].kurtosis()
            out.append(stats.to_string(max_rows=30))
        except Exception as e: out.append(f"   (统计计算失败: {e})")

    if cat_cols:
        out.append(f"\n── 分类列 Top5 ──")
        for c in cat_cols[:10]:
            vc = df[c].value_counts().head(5)
            out.append(f"  {c}:")
            for v, cnt in vc.items(): out.append(f"    {v}: {cnt} ({cnt/rows*100:.1f}%)")

    if len(num_cols) >= 2 and len(num_cols) <= 15:
        out.append(f"\n── 相关性矩阵 (|r|>0.3) ──")
        try:
            corr = df[num_cols].corr(); found = False
            for i in range(len(num_cols)):
                for j in range(i+1, len(num_cols)):
                    r = corr.iloc[i,j]
                    if abs(r) > 0.3: out.append(f"  {num_cols[i]} ↔ {num_cols[j]}: r={r:.3f}"); found = True
            if not found: out.append("  (无 |r|>0.3 的显著相关性)")
        except Exception: out.append("  (相关性计算失败)")

    return "\n".join(out)

# ═══════════════════════════════════════════════════════════
#  🆕 工具 14: BibTeX 解析/生成/引用格式化（✅ 已修复）
# ═══════════════════════════════════════════════════════════

@tool
def BibTexTool(
    mode: str, bib_path: str = "", bib_content: str = "",
    title: str = "", authors: str = "", journal: str = "",
    year: str = "", doi: str = "", volume: str = "", pages: str = "",
    citation_style: str = "apa"
) -> str:
    """
    【🆕v3.0】BibTeX 解析、生成与引用格式化。
    mode: 'parse' 解析.bib | 'generate' 生成BibTeX | 'format' 格式化引用(apa/mla/chicago/ieee)
    """
    # ── MODE: parse ──────────────────────────────────────
    if mode == "parse":
        if bib_content:
            raw = bib_content
        elif bib_path:
            rp = _resolve_path(bib_path)
            if not os.path.exists(rp):
                return f"❌ 文件不存在: {bib_path}"
            raw = _safe_read_text(rp)
        else:
            return "❌ parse模式需要 bib_path 或 bib_content。"
        entries = []
        import re
        pattern = re.compile(r'@(\w+)\s*\{([^,]*),\s*(.*?)\}', re.DOTALL)
        for m in pattern.finditer(raw):
            etype = m.group(1)
            ekey = m.group(2).strip()
            fields_str = m.group(3)
            fields = {}
            for fm in re.finditer(r'(\w+)\s*=\s*[{"]([^}"]*)[}"]', fields_str):
                fields[fm.group(1).lower()] = fm.group(2).strip()
            entries.append({"type": etype, "key": ekey, "fields": fields})
        if not entries:
            return "⚠ 未解析到 BibTeX 条目（可能格式不标准）。"
        out = [f"📚 BibTeX解析: {len(entries)}条\n"]
        for i, e in enumerate(entries, 1):
            f = e["fields"]
            t = f.get("title", "?")
            au = f.get("author", "?")
            yr = f.get("year", "?")
            jn = f.get("journal", f.get("booktitle", "?"))
            out.append(f"{i}. [{e['type']}] {t[:80]}{'...' if len(t)>80 else ''}")
            out.append(f"   👤 {au[:100]}{'...' if len(au)>100 else ''}")
            out.append(f"   📰 {jn} | {yr} | 🔑 {e['key']}")
        return "\n".join(out)

    # ── MODE: generate ────────────────────────────────────
    if mode == "generate":
        if not title:
            return "❌ generate模式需要 title 参数。"
        first_author = authors.split(",")[0].strip().split() if authors else ["unknown"]
        key_surname = first_author[-1] if first_author else "unknown"
        yr = year or "????"
        bib = f"@article{{{key_surname}{yr},\n"
        bib += f"  title = {{{title}}},\n"
        if authors:
            bib += f"  author = {{{authors}}},\n"
        if journal:
            bib += f"  journal = {{{journal}}},\n"
        bib += f"  year = {{{yr}}},\n"
        if volume:
            bib += f"  volume = {{{volume}}},\n"
        if pages:
            bib += f"  pages = {{{pages}}},\n"
        if doi:
            bib += f"  doi = {{{doi}}},\n"
        bib += "}"
        return f"✅ BibTeX 已生成:\n```bibtex\n{bib}\n```"

    # ── MODE: format ──────────────────────────────────────
    if mode == "format":
        if not title:
            return "❌ format模式需要 title 参数。"

        au_list = [a.strip() for a in authors.split(",") if a.strip()] if authors else ["Unknown"]
        yr = year or "(n.d.)"
        ref = ""

        if citation_style == "apa":
            # APA: Author (Year). Title. Journal. Volume, Pages. DOI
            if len(au_list) == 1:
                au_str = au_list[0]
            elif len(au_list) == 2:
                au_str = f"{au_list[0]} & {au_list[1]}"
            else:
                au_str = f"{au_list[0]} et al."
            ref = f"{au_str} ({yr}). {title}."
            if journal:
                ref += f" *{journal}*"
            if volume:
                ref += f", *{volume}*"
            if pages:
                ref += f", {pages}"
            ref += "."
            if doi:
                ref += f" https://doi.org/{doi}"

        elif citation_style == "mla":
            # MLA: Author. "Title." Journal Volume (Year): Pages. DOI.
            if len(au_list) > 2:
                au_str = f"{au_list[0]} et al."
            else:
                au_str = " and ".join(au_list)
            ref = f'{au_str}. "{title}."'
            if journal:
                ref += f" *{journal}*"
            if volume:
                ref += f" {volume}"
            ref += f" ({yr})"
            if pages:
                ref += f": {pages}."
            else:
                ref += "."
            if doi:
                ref += f" doi:{doi}."

        elif citation_style == "chicago":
            # Chicago: Author. "Title." Journal Volume (Year): Pages. DOI.
            if len(au_list) > 3:
                au_str = ", ".join(au_list[:3]) + ", et al."
            else:
                au_str = ", ".join(au_list)
            ref = f'{au_str}. "{title}."'
            if journal:
                ref += f" *{journal}*"
            if volume:
                ref += f" {volume}"
            ref += f" ({yr})"
            if pages:
                ref += f": {pages}."
            else:
                ref += "."
            if doi:
                ref += f" https://doi.org/{doi}."

        elif citation_style == "ieee":
            # IEEE: A. Author et al., "Title," Journal, vol. X, pp. Y, Year. doi: DOI.
            if len(au_list) > 3:
                au_str = "., ".join(au_list[:3]) + ", et al."
            else:
                au_str = "., ".join(au_list)
            ref = f'{au_str}, "{title},"'
            if journal:
                ref += f" *{journal}*"
            if volume:
                ref += f", vol. {volume}"
            if pages:
                ref += f", pp. {pages}"
            ref += f", {yr}."
            if doi:
                ref += f" doi: {doi}."

        else:
            return f"❌ 不支持的引用格式: {citation_style}。支持: apa, mla, chicago, ieee"

        return f"📝 引用 ({citation_style.upper()}):\n{ref}"

    return "❌ 无效 mode。可选: 'parse', 'generate', 'format'"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 15: PDF 表格提取
# ═══════════════════════════════════════════════════════════

@tool
def PDFTableExtractTool(pdf_path: str, page_range: str = "all") -> str:
    """
    【🆕v3.0】从 PDF 中提取表格数据。
    参数: pdf_path(PDF路径), page_range('1-3' 或 'all')
    """
    fp = _resolve_path(pdf_path)
    if not os.path.exists(fp): return f"❌ 文件不存在: {pdf_path}"
    try:
        doc = fitz.open(fp); total = len(doc)
        if page_range == "all":
            pages = list(range(total))
        else:
            pages = []
            for part in page_range.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    pages.extend(range(max(0, int(a)-1), min(total, int(b))))
                else:
                    pages.append(max(0, int(part)-1))
            pages = sorted(set(pages))

        out = [f"📊 PDF表格提取: {os.path.basename(fp)} (共{total}页，提取{len(pages)}页)\n"]
        table_count = 0
        for pi in pages:
            page = doc[pi]
            try:
                tabs = page.find_tables()
                if tabs and tabs.tables:
                    for ti, tab in enumerate(tabs.tables):
                        table_count += 1
                        out.append(f"── 表格{table_count} (Page {pi+1}) ──")
                        rows = []
                        for row in tab.extract():
                            rows.append(" | ".join(str(cell).replace("\n"," ") if cell else "" for cell in row))
                        out.append("\n".join(rows))
                        out.append("")
                else:
                    out.append(f"── Page {pi+1}: 未检测到表格 ──")
            except AttributeError:
                text = page.get_text("text")
                out.append(f"── Page {pi+1} (文本模式，升级PyMuPDF≥1.23.0以启用智能表格检测) ──")
                lines = text.split("\n")
                table_lines = [l for l in lines if len(l.split()) >= 3 and any(c.isdigit() for c in l)]
                if table_lines:
                    table_count += 1
                    out.append(f"  疑似表格行 ({len(table_lines)}行):")
                    for l in table_lines[:30]:
                        out.append(f"  {l.strip()}")
                else:
                    out.append("  (未检测到表格文本)")

        if table_count == 0:
            return "\n".join(out) + "\n⚠ 未提取到表格。可能原因：扫描版PDF（建议OCR）、纯图片表格。"
        return "\n".join(out)
    except Exception as e:
        return f"❌ PDF表格提取失败: {e}\n{traceback.format_exc()}"

# ═══════════════════════════════════════════════════════════
#  🆕 工具 16: 对话记录导出
# ═══════════════════════════════════════════════════════════

@tool
def SessionExportTool(content: str, filename: str = "session_export", export_format: str = "markdown") -> str:
    """
    【🆕v3.0】将对话内容或分析结果导出为文件。
    export_format: 'markdown' (.md) / 'text' (.txt) / 'html' (.html)
    """
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    header = f"---\ntitle: Session Export\ndate: {ts}\n---\n\n"
    full_content = header + content
    if export_format == "markdown":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.md")
        with open(sp, "w", encoding="utf-8") as f:
            f.write(full_content)
        return f"✅ 对话已导出为 Markdown: {sp}"
    elif export_format == "text":
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.txt")
        with open(sp, "w", encoding="utf-8") as f:
            f.write(full_content)
        return f"✅ 对话已导出为文本: {sp}"
    elif export_format == "html":
        html_body = content.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>\n")
        html = f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8">
<title>Session Export</title>
<style>body{{font-family:Georgia,serif;max-width:800px;margin:2em auto;padding:0 1em;line-height:1.7;color:#333}}
h1,h2,h3{{color:#1a1a2e}}code{{background:#f4f4f4;padding:2px 6px;border-radius:4px}}
pre{{background:#f4f4f4;padding:1em;border-radius:6px;overflow-x:auto}}
blockquote{{border-left:4px solid #ccc;margin-left:0;padding-left:1em;color:#555}}</style></head>
<body><h1>Session Export</h1><p><em>{ts}</em></p>{html_body}</body></html>"""
        sp = os.path.join(CUSTOM_TEMP_DIR, f"{filename}.html")
        with open(sp, "w", encoding="utf-8") as f:
            f.write(html)
        return f"✅ 对话已导出为 HTML: {sp}"
    else:
        return f"❌ 不支持的格式: {export_format}。支持: markdown, text, html"

# ═══════════════════════════════════════════════════════════
#  4. 工具注册 & Agent 构建
# ═══════════════════════════════════════════════════════════

safe_write_tool = WriteFileTool(root_dir=CUSTOM_TEMP_DIR)

tools = [
    SmartReadPathTool, EditTexFileTool, CompileLatexTool,
    PythonSandboxTool, DuckDuckGoSearchTool, FetchWebImageTool,
    EditWordDocTool, SaveMarkdownTool, ChartGenerationTool,
    MermaidTool,
    ArXivSearchTool, DOIMetadataTool, DataStatisticsTool,
    BibTexTool, PDFTableExtractTool, SessionExportTool,
    safe_write_tool,
]

DB_PATH = os.path.join(BASE_WORKSPACE, "agent_memory.db")
conn = sqlite3.connect(DB_PATH, check_same_thread=False)
memory = SqliteSaver(conn)
agent = create_react_agent(llm, tools, prompt=SYSTEM_PROMPT, checkpointer=memory)

# ═══════════════════════════════════════════════════════════
#  5. 终端交互主循环
# ═══════════════════════════════════════════════════════════

HELP_TEXT = f"""
{_C.BOLD}╔══════════════════ v3.0 工具一览 ══════════════════╗{_C.W}
{_C.G}║ 📖 SmartReadPathTool     读取文件/文件夹          ║
║ 📝 EditTexFileTool        创建/编辑 .tex            ║
║ 🖨  CompileLatexTool       编译 .tex → PDF          ║
║ 🐍 PythonSandboxTool      安全执行 Python 代码      ║
║ 🔍 DuckDuckGoSearchTool   通用网页搜索              ║
║ 🖼  FetchWebImageTool      搜索下载学术图片          ║
║ 📄 EditWordDocTool        创建/编辑 Word 文档       ║
║ 📋 SaveMarkdownTool       保存 Markdown 文件        ║
║ 📊 ChartGenerationTool    matplotlib 图表           ║
║ 🎨 MermaidTool            Mermaid 流程图            ║{_C.W}
{_C.Y}║  ────────────── 🆕 v3.0 ──────────────           ║
║ 📚 ArXivSearchTool        arXiv 学术论文检索        ║
║ 📄 DOIMetadataTool        DOI → 完整元数据          ║
║ 📈 DataStatisticsTool     CSV/Excel 自动统计分析    ║
║ 📖 BibTexTool             .bib解析/生成/引用格式化  ║
║ 📋 PDFTableExtractTool    PDF 表格智能提取          ║
║ 💾 SessionExportTool      对话记录导出              ║{_C.W}
{_C.BOLD}╚══════════════════════════════════════════════════╝{_C.W}
"""

print(f"\n{_C.BOLD}💡 输入 'exit' 退出 | 'new' 新对话 | 'help' 查看16个工具{_C.W}\n")

current_thread_id = str(uuid.uuid4())
config = {"configurable": {"thread_id": current_thread_id}}

while True:
    try:
        user_input = input(f"{_C.C}👤 你:{_C.W} ").strip()
        if user_input.lower() in {"exit", "quit", "退出"}:
            _log("INFO", "👋 再见！"); conn.close(); break
        if user_input.lower() in {"new", "clear", "清空", "新对话"}:
            current_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": current_thread_id}}
            _log("INFO", "✨ 全新对话已开启。\n"); continue
        if user_input.lower() == "help":
            print(HELP_TEXT); continue
        if not user_input: continue

        print(f"\n{_C.M}🤔 思考/执行中...{_C.W}\n")
        inputs = {"messages": [("user", user_input)]}
        t_start = time.time()

        for event in agent.stream(inputs, config, stream_mode="updates"):
            if "agent" in event:
                msg = event["agent"]["messages"][0]
                if msg.content:
                    elapsed = time.time() - t_start
                    print(f"{_C.G}🤖 DeepSeek{_C.W} ({elapsed:.1f}s):\n{msg.content}\n")
                elif msg.tool_calls:
                    for tc in msg.tool_calls:
                        args_str = json.dumps(tc["args"], ensure_ascii=False, default=str)
                        if len(args_str) > 150: args_str = args_str[:150] + "..."
                        print(f"{_C.Y}🛠  [{tc['name']}]{_C.W}\n   参数: {args_str}")
            elif "tools" in event:
                msg = event["tools"]["messages"][0]
                cs = str(msg.content); cl = len(cs)
                preview = cs[:250].replace("\n", " ") + ("..." if cl > 250 else "")
                is_err = cs.startswith("❌") or "错误" in cs[:100]
                icon = _C.R+"📄" if is_err else _C.G+"📄"
                print(f"{icon} [返回 {cl}字符]{_C.W}: {preview}\n")

    except KeyboardInterrupt:
        print(f"\n{_C.Y}⏸  [操作被用户中断]{_C.W}")
    except Exception as e:
        print(f"{_C.R}❌ 发生错误: {e}{_C.W}")
        _log("ERR", traceback.format_exc())
